"""Coverage for the RAG upload / document / citations-export endpoints in
src/nyxgpt/app.py (issue #3234, scoped split of #3218):

- POST /rag/ingest (failure path -- success is covered by test_metrics.py)
- GET /rag/documents/{doc_id} (rag_document_info)
- GET /rag/documents (rag_documents_list)
- GET /sessions/{name}/documents (list_session_documents non-list metadata edge case)
- POST /rag/upload (rag_upload_file) -- multi text-format parsing (txt, json, md,
  pptx, docx, epub, html, pdf incl. OCR fallback), size-limit boundary, and each
  format's validation/error branches.

Citations export (GET /sessions/{name}/citations/export) is already fully
covered by test_message_format_endpoints.py and is intentionally not
duplicated here.
"""

from __future__ import annotations

import io
import sys
from typing import Any
from unittest.mock import Mock, patch

import pytest
from fastapi.testclient import TestClient

from nyxgpt.app import app

pytestmark = pytest.mark.unit

UPLOAD_URL = "/api/v1/rag/upload"


def _client() -> TestClient:
    return TestClient(app)


def _fake_ingest_result(**overrides: Any) -> dict[str, Any]:
    result = {
        "chunks_ingested": 1,
        "status": "ingested",
        "doc_hash": "hash123",
        "previous_hash": None,
    }
    result.update(overrides)
    return result


# ============================================================================
# POST /rag/ingest
# ============================================================================


def test_rag_ingest_failure_returns_400_and_increments_failure_metric() -> None:
    client = _client()
    with patch("nyxgpt.app.ingest_document", side_effect=ValueError("embedding backend down")):
        response = client.post(
            "/api/v1/rag/ingest",
            json={"doc_id": "doc-fail", "text": "hello world"},
        )

    assert response.status_code == 400
    assert response.json()["error"]["message"] == "embedding backend down"

    metrics_text = client.get("/metrics").text
    assert 'nyxgpt_rag_ingests_total{result="failure",source="document"}' in metrics_text.replace(
        " ", ""
    )


def test_rag_ingest_missing_doc_id_returns_422() -> None:
    client = _client()
    response = client.post("/api/v1/rag/ingest", json={"text": "hello world"})
    assert response.status_code == 422


# ============================================================================
# GET /rag/documents/{doc_id}
# ============================================================================


def test_rag_document_info_found() -> None:
    mock_store = Mock()
    mock_store.get_document_info.return_value = {
        "doc_id": "doc-1",
        "doc_hash": "abc",
        "ingested_at": "2026-01-01T00:00:00",
        "updated_at": "2026-01-02T00:00:00",
        "chunks": 3,
        "embedding_model": "nomic-embed-text",
    }
    client = _client()
    with patch(
        "nyxgpt.rag.vectorstore_cassandra.CassandraVectorStore",
        lambda **kwargs: mock_store,
    ):
        response = client.get("/api/v1/rag/documents/doc-1")

    assert response.status_code == 200
    body = response.json()
    assert body["doc_id"] == "doc-1"
    assert body["chunks"] == 3
    mock_store.close.assert_called_once()


def test_rag_document_info_not_found_returns_404() -> None:
    mock_store = Mock()
    mock_store.get_document_info.return_value = None
    client = _client()
    with patch(
        "nyxgpt.rag.vectorstore_cassandra.CassandraVectorStore",
        lambda **kwargs: mock_store,
    ):
        response = client.get("/api/v1/rag/documents/missing-doc")

    assert response.status_code == 404
    mock_store.close.assert_called_once()


def test_rag_document_info_store_error_returns_400() -> None:
    mock_store = Mock()
    mock_store.get_document_info.side_effect = RuntimeError("connection reset")
    client = _client()
    with patch(
        "nyxgpt.rag.vectorstore_cassandra.CassandraVectorStore",
        lambda **kwargs: mock_store,
    ):
        response = client.get("/api/v1/rag/documents/doc-err")

    assert response.status_code == 400
    assert "connection reset" in response.json()["error"]["message"]
    mock_store.close.assert_called_once()


# ============================================================================
# GET /rag/documents
# ============================================================================


def test_rag_documents_list_enriches_with_metadata() -> None:
    """`metadata` is a Cassandra `text` column holding a JSON-encoded object --
    the driver hands back a raw string, never a pre-parsed dict."""
    mock_store = Mock()
    mock_store.list_docs.return_value = [
        {"doc_id": "doc-1", "chunks": 2, "embedding_model": "nomic-embed-text"},
    ]
    mock_store.table_name = "documents"
    row = Mock()
    row.metadata = '{"filename": "report.pdf", "tags": ["finance"]}'
    row.ingested_at = None
    mock_store.session.execute.return_value = [row]

    client = _client()
    with patch(
        "nyxgpt.rag.vectorstore_cassandra.CassandraVectorStore",
        lambda **kwargs: mock_store,
    ):
        response = client.get("/api/v1/rag/documents")

    assert response.status_code == 200
    docs = response.json()["documents"]
    assert docs == [
        {
            "doc_id": "doc-1",
            "chunks": 2,
            "embedding_model": "nomic-embed-text",
            "filename": "report.pdf",
            "tags": ["finance"],
            "ingested_at": None,
        }
    ]
    mock_store.close.assert_called_once()


def test_rag_documents_list_enriches_legacy_string_metadata_without_warning_path() -> None:
    """Regression fixture for #3541: pre-citation-enrichment rows (like
    `smoke-test-doc`) store metadata as a bare, non-JSON string rather than a
    JSON object. Enrichment must degrade filename/tags to None instead of
    raising `'str' object has no attribute 'get'`."""
    mock_store = Mock()
    mock_store.list_docs.return_value = [
        {"doc_id": "smoke-test-doc", "chunks": 1, "embedding_model": "nomic-embed-text"},
    ]
    mock_store.table_name = "documents"
    row = Mock()
    row.metadata = "smoke-test-doc"
    row.ingested_at = None
    mock_store.session.execute.return_value = [row]

    client = _client()
    with patch(
        "nyxgpt.rag.vectorstore_cassandra.CassandraVectorStore",
        lambda **kwargs: mock_store,
    ):
        response = client.get("/api/v1/rag/documents")

    assert response.status_code == 200
    docs = response.json()["documents"]
    assert docs == [
        {
            "doc_id": "smoke-test-doc",
            "chunks": 1,
            "embedding_model": "nomic-embed-text",
            "filename": None,
            "tags": None,
            "ingested_at": None,
        }
    ]
    mock_store.close.assert_called_once()


def test_rag_documents_list_enriches_malformed_metadata_string() -> None:
    """Metadata that isn't even valid JSON also degrades gracefully."""
    mock_store = Mock()
    mock_store.list_docs.return_value = [
        {"doc_id": "doc-broken", "chunks": 1, "embedding_model": "nomic-embed-text"},
    ]
    mock_store.table_name = "documents"
    row = Mock()
    row.metadata = "{not valid json"
    row.ingested_at = None
    mock_store.session.execute.return_value = [row]

    client = _client()
    with patch(
        "nyxgpt.rag.vectorstore_cassandra.CassandraVectorStore",
        lambda **kwargs: mock_store,
    ):
        response = client.get("/api/v1/rag/documents")

    assert response.status_code == 200
    docs = response.json()["documents"]
    assert docs == [
        {
            "doc_id": "doc-broken",
            "chunks": 1,
            "embedding_model": "nomic-embed-text",
            "filename": None,
            "tags": None,
            "ingested_at": None,
        }
    ]


def test_rag_documents_list_falls_back_when_metadata_enrichment_fails() -> None:
    mock_store = Mock()
    mock_store.list_docs.return_value = [
        {"doc_id": "doc-2", "chunks": 5, "embedding_model": "nomic-embed-text"},
    ]
    mock_store.table_name = "documents"
    mock_store.session.execute.side_effect = RuntimeError("query failed")

    client = _client()
    with patch(
        "nyxgpt.rag.vectorstore_cassandra.CassandraVectorStore",
        lambda **kwargs: mock_store,
    ):
        response = client.get("/api/v1/rag/documents")

    assert response.status_code == 200
    docs = response.json()["documents"]
    assert docs == [
        {
            "doc_id": "doc-2",
            "chunks": 5,
            "embedding_model": "nomic-embed-text",
            "filename": None,
            "tags": None,
            "ingested_at": None,
        }
    ]


def test_rag_documents_list_store_failure_returns_500() -> None:
    mock_store = Mock()
    mock_store.list_docs.side_effect = RuntimeError("keyspace unavailable")

    client = _client()
    with patch(
        "nyxgpt.rag.vectorstore_cassandra.CassandraVectorStore",
        lambda **kwargs: mock_store,
    ):
        response = client.get("/api/v1/rag/documents")

    assert response.status_code == 500
    mock_store.close.assert_called_once()


# ============================================================================
# GET /sessions/{name}/documents
# ============================================================================


def test_list_session_documents_resets_non_list_attached_doc_ids(tmp_path, monkeypatch) -> None:
    from nyxgpt import sessions

    sessions_dir = tmp_path / "sessions"
    sessions_dir.mkdir(parents=True, exist_ok=True)
    sf = sessions.session_file_for("weird-meta", sessions_dir)
    mf = sessions.meta_file_for(sf)
    sessions.save_session_messages(sf, [])
    sessions.save_session_meta(mf, {"attached_doc_ids": "not-a-list"})

    monkeypatch.setattr(
        "nyxgpt.app.get_sessions_dir",
        lambda cfg: sessions_dir,
    )

    client = _client()
    response = client.get("/api/v1/sessions/weird-meta/documents")

    assert response.status_code == 200
    assert response.json()["attached_doc_ids"] == []


# ============================================================================
# POST /rag/upload -- validation and size limits
# ============================================================================


def test_rag_upload_rejects_unsupported_file_type() -> None:
    client = _client()
    response = client.post(
        UPLOAD_URL,
        files={"file": ("archive.zip", b"binary-content", "application/zip")},
    )
    assert response.status_code == 400
    assert "not supported" in response.json()["error"]["message"]


def test_rag_upload_exactly_at_size_limit_succeeds(monkeypatch: pytest.MonkeyPatch) -> None:
    # The global request-size middleware (MAX_BODY_BYTES, 1 MiB default) would
    # reject this before it ever reaches rag_upload_file's own 10 MB check, so
    # raise it out of the way to isolate the endpoint's own boundary.
    monkeypatch.setattr("nyxgpt.app.MAX_BODY_BYTES", 20 * 1024 * 1024)
    content = b"a" * (10 * 1024 * 1024)
    client = _client()
    with patch("nyxgpt.app.ingest_document", return_value=_fake_ingest_result()):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("big.txt", content, "text/plain")},
        )
    assert response.status_code == 200


def test_rag_upload_over_size_limit_returns_413(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr("nyxgpt.app.MAX_BODY_BYTES", 20 * 1024 * 1024)
    content = b"a" * (10 * 1024 * 1024 + 1)
    client = _client()
    response = client.post(
        UPLOAD_URL,
        files={"file": ("toobig.txt", content, "text/plain")},
    )
    assert response.status_code == 413
    assert "too large" in response.json()["error"]["message"].lower()
    assert "10485761 bytes" in response.json()["error"]["message"]


# ============================================================================
# POST /rag/upload -- .txt
# ============================================================================


def test_rag_upload_txt_success_uses_filename_as_doc_id() -> None:
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("notes.txt", b"plain text content", "text/plain")},
        )

    assert response.status_code == 200
    assert response.json()["doc_id"] == "notes.txt"
    assert captured["doc_id"] == "notes.txt"
    assert captured["text"] == "plain text content"


def test_rag_upload_txt_explicit_doc_id_overrides_filename() -> None:
    client = _client()
    with patch("nyxgpt.app.ingest_document", return_value=_fake_ingest_result()):
        response = client.post(
            f"{UPLOAD_URL}?doc_id=custom-id",
            files={"file": ("notes.txt", b"plain text content", "text/plain")},
        )
    assert response.status_code == 200
    assert response.json()["doc_id"] == "custom-id"


def test_rag_upload_txt_decode_error_returns_400() -> None:
    client = _client()
    response = client.post(
        UPLOAD_URL,
        files={"file": ("bad.txt", b"\xff\xfe\x00\x01", "text/plain")},
    )
    assert response.status_code == 400
    assert "encoding" in response.json()["error"]["message"].lower()


def test_rag_upload_ingestion_failure_returns_500() -> None:
    client = _client()
    with patch("nyxgpt.app.ingest_document", side_effect=RuntimeError("vector store down")):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("notes.txt", b"plain text content", "text/plain")},
        )
    assert response.status_code == 500
    assert "Ingestion failed" in response.json()["error"]["message"]


# ============================================================================
# POST /rag/upload -- collection targeting (#3463)
# ============================================================================


def test_rag_upload_with_collection_param_ingests_into_that_collection() -> None:
    """Chat-uploaded documents must be able to target a non-default
    collection -- previously `/rag/upload` had no `collection` param at
    all and always ingested into "default"."""
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    mock_store = Mock()
    mock_store.list_collections.return_value = ["default", "research-notes"]

    with (
        patch("nyxgpt.rag.vectorstore_cassandra.CassandraVectorStore", return_value=mock_store),
        patch("nyxgpt.app.ingest_document", side_effect=fake_ingest),
    ):
        response = client.post(
            f"{UPLOAD_URL}?collection=research-notes",
            files={"file": ("notes.txt", b"plain text content", "text/plain")},
        )

    assert response.status_code == 200
    assert response.json()["collection"] == "research-notes"
    assert captured["collection"] == "research-notes"


def test_rag_upload_without_collection_defaults_to_default() -> None:
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("notes.txt", b"plain text content", "text/plain")},
        )

    assert response.status_code == 200
    assert response.json()["collection"] == "default"
    assert captured["collection"] == "default"


def test_rag_upload_with_unknown_collection_returns_400() -> None:
    """A typo'd collection name must be rejected before it silently opens a
    brand-new, unrelated empty table."""
    client = _client()
    mock_store = Mock()
    mock_store.list_collections.return_value = ["default"]

    with patch("nyxgpt.rag.vectorstore_cassandra.CassandraVectorStore", return_value=mock_store):
        response = client.post(
            f"{UPLOAD_URL}?collection=does-not-exist",
            files={"file": ("notes.txt", b"plain text content", "text/plain")},
        )

    assert response.status_code == 400
    assert "Unknown collection" in response.json()["error"]["message"]


# ============================================================================
# POST /rag/upload -- .json
# ============================================================================


def test_rag_upload_json_success_reformats_content() -> None:
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("data.json", b'{"a":1,"b":2}', "application/json")},
        )

    assert response.status_code == 200
    assert '"a": 1' in captured["text"]


def test_rag_upload_json_invalid_returns_400() -> None:
    client = _client()
    response = client.post(
        UPLOAD_URL,
        files={"file": ("bad.json", b"{not valid json", "application/json")},
    )
    assert response.status_code == 400
    assert "Invalid JSON" in response.json()["error"]["message"]


# ============================================================================
# POST /rag/upload -- .md
# ============================================================================


def test_rag_upload_md_success_with_frontmatter() -> None:
    md_content = b"""---
title: My Doc
author: Jane
---

# Heading

Some **body** text.
"""
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("doc.md", md_content, "text/markdown")},
        )

    assert response.status_code == 200
    assert "title: My Doc" in captured["text"]
    assert "Heading" in captured["text"]
    assert "body" in captured["text"]


def test_rag_upload_md_falls_back_to_plain_text_on_import_error() -> None:
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with (
        patch.dict(sys.modules, {"frontmatter": None}),
        patch("nyxgpt.app.ingest_document", side_effect=fake_ingest),
    ):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("doc.md", b"# Just markdown", "text/markdown")},
        )

    assert response.status_code == 200
    assert captured["text"] == "# Just markdown"


def test_rag_upload_md_parsing_failure_returns_400() -> None:
    client = _client()
    with patch("frontmatter.loads", side_effect=RuntimeError("boom")):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("doc.md", b"# Heading", "text/markdown")},
        )
    assert response.status_code == 400
    assert "Markdown parsing failed" in response.json()["error"]["message"]


# ============================================================================
# POST /rag/upload -- .pptx
# ============================================================================


def _build_pptx(*, with_text: bool, with_notes: bool = False) -> bytes:
    from pptx import Presentation

    prs = Presentation()
    if with_text:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Slide Title"
        if with_notes:
            slide.notes_slide.notes_text_frame.text = "Speaker notes here"
    else:
        blank_layout = prs.slide_layouts[6]
        prs.slides.add_slide(blank_layout)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_rag_upload_pptx_success_with_notes() -> None:
    content = _build_pptx(with_text=True, with_notes=True)
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={
                "file": (
                    "deck.pptx",
                    content,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )

    assert response.status_code == 200
    assert "[Slide 1]" in captured["text"]
    assert "Slide Title" in captured["text"]
    assert "[Speaker Notes]" in captured["text"]
    assert "Speaker notes here" in captured["text"]


def test_rag_upload_pptx_empty_returns_400() -> None:
    content = _build_pptx(with_text=False)
    client = _client()
    response = client.post(
        UPLOAD_URL,
        files={
            "file": (
                "empty.pptx",
                content,
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert response.status_code == 400
    assert "no extractable text" in response.json()["error"]["message"]


def test_rag_upload_pptx_import_error_returns_400() -> None:
    content = _build_pptx(with_text=True)
    client = _client()
    with patch.dict(sys.modules, {"pptx": None}):
        response = client.post(
            UPLOAD_URL,
            files={
                "file": (
                    "deck.pptx",
                    content,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
    assert response.status_code == 400
    assert "PPTX support not available" in response.json()["error"]["message"]


def test_rag_upload_pptx_parsing_failure_returns_400() -> None:
    content = _build_pptx(with_text=True)
    client = _client()
    with patch("pptx.Presentation", side_effect=RuntimeError("corrupt")):
        response = client.post(
            UPLOAD_URL,
            files={
                "file": (
                    "deck.pptx",
                    content,
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )
    assert response.status_code == 400
    assert "PPTX parsing failed" in response.json()["error"]["message"]


# ============================================================================
# POST /rag/upload -- .docx
# ============================================================================


def _build_docx(*, with_content: bool = True, with_image: bool = False) -> bytes:
    from docx import Document

    doc = Document()
    if with_content:
        doc.add_heading("Doc Title", level=1)
        doc.add_paragraph("A normal paragraph.")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "H1"
        table.cell(0, 1).text = "H2"
        table.cell(1, 0).text = "v1"
        table.cell(1, 1).text = "v2"
        if with_image:
            from PIL import Image

            img_buf = io.BytesIO()
            Image.new("RGB", (10, 10), color="red").save(img_buf, format="PNG")
            img_buf.seek(0)

            # Image 1: python-docx only ever sets docPr's `name` attribute
            # by default -- exercises rag_upload_file's "elif ...name"
            # branch.
            para1 = doc.add_paragraph()
            para1.add_run().add_picture(img_buf, width=100000)

            # Image 2: set `descr` (alt text) directly via the XML (python-
            # docx has no API for it) to exercise the "if ...descr" branch,
            # which is checked first.
            img_buf.seek(0)
            para2 = doc.add_paragraph()
            run2 = para2.add_run()
            run2.add_picture(img_buf, width=100000)
            docpr_ns = "{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}"
            for docpr in run2._element.findall(f".//{docpr_ns}docPr"):
                docpr.set("descr", "A red square")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def test_rag_upload_docx_success_with_headings_table_and_image() -> None:
    content = _build_docx(with_content=True, with_image=True)
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={
                "file": (
                    "report.docx",
                    content,
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
        )

    assert response.status_code == 200
    assert "## Doc Title" in captured["text"]
    assert "A normal paragraph." in captured["text"]
    assert "[Table]" in captured["text"]
    assert "H1 | H2" in captured["text"]
    assert "[Image 1: Picture" in captured["text"]
    assert "[Image 2: A red square]" in captured["text"]


def test_rag_upload_docx_empty_returns_400() -> None:
    content = _build_docx(with_content=False)
    client = _client()
    response = client.post(
        UPLOAD_URL,
        files={
            "file": (
                "empty.docx",
                content,
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
    )
    assert response.status_code == 400
    assert "empty" in response.json()["error"]["message"].lower()


def test_rag_upload_docx_package_not_found_returns_400() -> None:
    client = _client()
    response = client.post(
        UPLOAD_URL,
        files={
            "file": (
                "corrupt.docx",
                b"not a real docx file at all",
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            )
        },
    )
    assert response.status_code == 400
    assert "corrupted" in response.json()["error"]["message"].lower()


def test_rag_upload_docx_bad_zip_returns_400() -> None:
    from docx.opc.exceptions import PackageNotFoundError

    content = _build_docx(with_content=True)
    client = _client()
    with patch("docx.Document", side_effect=PackageNotFoundError("bad zip")):
        response = client.post(
            UPLOAD_URL,
            files={
                "file": (
                    "report.docx",
                    content,
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
        )
    assert response.status_code == 400
    assert "corrupted" in response.json()["error"]["message"].lower()


def test_rag_upload_docx_import_error_returns_400() -> None:
    content = _build_docx(with_content=True)
    client = _client()
    with patch.dict(sys.modules, {"docx": None}):
        response = client.post(
            UPLOAD_URL,
            files={
                "file": (
                    "report.docx",
                    content,
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
        )
    assert response.status_code == 400
    assert "DOCX support not available" in response.json()["error"]["message"]


def test_rag_upload_docx_generic_parsing_failure_returns_400() -> None:
    from unittest.mock import PropertyMock

    content = _build_docx(with_content=True)
    client = _client()
    mock_doc = Mock()
    type(mock_doc).paragraphs = PropertyMock(side_effect=RuntimeError("xml corrupt"))
    with patch("docx.Document", return_value=mock_doc):
        response = client.post(
            UPLOAD_URL,
            files={
                "file": (
                    "report.docx",
                    content,
                    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                )
            },
        )
    assert response.status_code == 400
    assert "DOCX parsing failed" in response.json()["error"]["message"]


# ============================================================================
# POST /rag/upload -- .epub
# ============================================================================


def _build_epub(*, with_content: bool = True) -> bytes:
    import tempfile
    from pathlib import Path as _Path

    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier("id-1")
    book.set_title("My Book")
    book.set_language("en")
    book.add_author("Author One")
    book.add_metadata("DC", "description", "A test description")
    book.add_metadata("DC", "publisher", "Pub House")
    book.add_metadata("DC", "date", "2024-01-01")

    # No EpubNav here deliberately: ebooklib classifies nav.xhtml as
    # ITEM_DOCUMENT too, which would make the single-chapter case below
    # always look like a >1-document book and always take the "[Chapter N]"
    # prefixed branch, never the unprefixed one.
    items: list[Any] = [epub.EpubNcx()]
    spine: list[Any] = []

    if with_content:
        chapter = epub.EpubHtml(title="Chapter 1", file_name="chap1.xhtml", lang="en")
        chapter.content = (
            "<script>alert('xss')</script><h1>Chapter 1</h1><p>" + ("word " * 30) + "</p>"
        )
        items.append(chapter)
        spine.append(chapter)
        book.toc = (epub.Link("chap1.xhtml", "Chapter 1", "chap1"),)

    for item in items:
        book.add_item(item)
    book.spine = spine

    with tempfile.TemporaryDirectory() as tmp:
        path = _Path(tmp) / "book.epub"
        epub.write_epub(str(path), book)
        return path.read_bytes()


def test_rag_upload_epub_success_with_metadata_and_chapter() -> None:
    content = _build_epub(with_content=True)
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("book.epub", content, "application/epub+zip")},
        )

    assert response.status_code == 200
    text = captured["text"]
    assert "[Metadata]" in text
    # NOTE: extract_metadata_values() in rag_upload_file assumes ebooklib
    # returns (namespace, value) tuples, but it actually returns
    # (value, attrs_dict) -- so the metadata section renders the (usually
    # empty) attrs dict rather than the real value. This is pre-existing
    # behavior in src/nyxgpt/app.py, not something introduced by this test.
    assert "Title: {}" in text
    assert "Author: {'id': 'creator'}" in text
    # NOTE: the heading-hierarchy extraction loop above builds
    # `chapter_texts` ("# Chapter 1") but that variable is never joined into
    # `chapter_text` -- only the raw get_text() lines are used -- so the
    # "#"-prefixed heading never actually appears in the ingested text.
    assert "Chapter 1" in text
    assert "alert" not in text  # <script> content is stripped before extraction
    assert "word word word" in text


def test_rag_upload_epub_success_with_multiple_chapters_prefixes_each() -> None:
    """With >1 content item, each chapter is prefixed with "[Chapter N]" --
    the single-chapter test above exercises the un-prefixed branch instead."""
    import tempfile
    from pathlib import Path as _Path

    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier("id-2")
    book.set_title("Two Chapter Book")
    book.set_language("en")

    chapters = []
    for i in (1, 2):
        chapter = epub.EpubHtml(title=f"Chapter {i}", file_name=f"chap{i}.xhtml", lang="en")
        chapter.content = f"<h1>Chapter {i}</h1><p>" + (f"content{i} " * 30) + "</p>"
        book.add_item(chapter)
        chapters.append(chapter)

    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", *chapters]

    with tempfile.TemporaryDirectory() as tmp:
        path = _Path(tmp) / "book.epub"
        epub.write_epub(str(path), book)
        content = path.read_bytes()

    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("two.epub", content, "application/epub+zip")},
        )

    assert response.status_code == 200
    text = captured["text"]
    assert "[Chapter 1]" in text
    assert "[Chapter 2]" in text
    assert "content1" in text
    assert "content2" in text


def test_rag_upload_epub_invalid_file_returns_400() -> None:
    client = _client()
    response = client.post(
        UPLOAD_URL,
        files={"file": ("bad.epub", b"not a real epub", "application/epub+zip")},
    )
    assert response.status_code == 400
    assert "Invalid ePUB" in response.json()["error"]["message"]


def test_rag_upload_epub_no_content_returns_400() -> None:
    content = _build_epub(with_content=False)
    client = _client()
    response = client.post(
        UPLOAD_URL,
        files={"file": ("empty.epub", content, "application/epub+zip")},
    )
    assert response.status_code == 400
    assert "no text" in response.json()["error"]["message"].lower()


def test_rag_upload_epub_generic_processing_failure_returns_400() -> None:
    """A failure *after* epub.read_epub() succeeds (e.g. while parsing a
    chapter's HTML body) hits the outer generic except, distinct from the
    "Invalid ePUB file" handler wrapping read_epub() itself."""
    content = _build_epub(with_content=True)
    client = _client()
    with patch("bs4.BeautifulSoup", side_effect=RuntimeError("soup exploded")):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("book.epub", content, "application/epub+zip")},
        )
    assert response.status_code == 400
    assert "ePUB parsing failed" in response.json()["error"]["message"]


def test_rag_upload_epub_import_error_returns_400() -> None:
    content = _build_epub(with_content=True)
    client = _client()
    with patch.dict(sys.modules, {"ebooklib": None}):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("book.epub", content, "application/epub+zip")},
        )
    assert response.status_code == 400
    assert "ePUB support not available" in response.json()["error"]["message"]


# ============================================================================
# POST /rag/upload -- .html / .htm
# ============================================================================


def test_rag_upload_html_success_with_structured_content() -> None:
    html = b"""<html>
<head>
<title>Page Title</title>
<meta name="description" content="A test page">
<meta name="author" content="Jane Doe">
</head>
<body>
<nav>Skip this nav</nav>
<main>
<h1>Main Heading</h1>
<p>A paragraph of body text.</p>
<p></p>
<div class="advertisement">Buy now!</div>
<div><p>Nested inside container</p></div>
<blockquote>A quoted block.</blockquote>
<pre>code block</pre>
<ul><li>item one</li><li>item two</li></ul>
<ol><li>first</li><li>second</li></ol>
<table><tr><th>Col1</th><th>Col2</th></tr><tr><td>a</td><td>b</td></tr></table>
</main>
<footer>Skip this footer</footer>
</body>
</html>"""
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("page.html", html, "text/html")},
        )

    assert response.status_code == 200
    text = captured["text"]
    assert "Title: Page Title" in text
    assert "# Main Heading" in text
    assert "A paragraph of body text." in text
    assert "> A quoted block." in text
    assert "```\ncode block\n```" in text
    assert "• item one" in text
    assert "1. first" in text
    assert "[Table]" in text
    assert "Skip this nav" not in text
    assert "Skip this footer" not in text
    assert "Buy now!" not in text
    assert "Nested inside container" in text


def test_rag_upload_html_fallback_content_extraction() -> None:
    html = b"<html><body><span>just a span with no block tags</span></body></html>"
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("page.htm", html, "text/html")},
        )

    assert response.status_code == 200
    assert "just a span with no block tags" in captured["text"]


def test_rag_upload_html_fallback_encoding_decodes_latin1() -> None:
    html = "<html><body><p>Café au lait</p></body></html>".encode("iso-8859-1")
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("page.html", html, "text/html")},
        )

    assert response.status_code == 200
    assert "Caf" in captured["text"]


def test_rag_upload_html_empty_returns_400() -> None:
    html = b"<html><body></body></html>"
    client = _client()
    response = client.post(
        UPLOAD_URL,
        files={"file": ("empty.html", html, "text/html")},
    )
    assert response.status_code == 400
    assert "no text" in response.json()["error"]["message"].lower()


def test_rag_upload_html_import_error_returns_400() -> None:
    html = b"<html><body><p>content</p></body></html>"
    client = _client()
    with patch.dict(sys.modules, {"bs4": None}):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("page.html", html, "text/html")},
        )
    assert response.status_code == 400
    assert "HTML support not available" in response.json()["error"]["message"]


def test_rag_upload_html_parsing_failure_returns_400() -> None:
    html = b"<html><body><p>content</p></body></html>"
    client = _client()
    with patch("bs4.BeautifulSoup", side_effect=RuntimeError("parser exploded")):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("page.html", html, "text/html")},
        )
    assert response.status_code == 400
    assert "HTML parsing failed" in response.json()["error"]["message"]


# ============================================================================
# POST /rag/upload -- .pdf
# ============================================================================


def _build_text_pdf(text: str = "Hello PDF World") -> bytes:
    from reportlab.pdfgen import canvas

    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.setTitle("Test PDF")
    c.setAuthor("Test Author")
    c.drawString(100, 750, text)
    c.showPage()
    c.save()
    return buf.getvalue()


def _build_image_only_pdf() -> bytes:
    from reportlab.pdfgen import canvas

    buf = io.BytesIO()
    c = canvas.Canvas(buf)
    c.rect(100, 100, 50, 50, fill=1)
    c.showPage()
    c.save()
    return buf.getvalue()


def _build_table_pdf() -> bytes:
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=(300, 300))
    table = Table([["H1", "H2"], ["v1", "v2"]])
    table.setStyle(TableStyle([("GRID", (0, 0), (-1, -1), 1, colors.black)]))
    doc.build([table])
    return buf.getvalue()


def _no_pdf_metadata():
    """Patch pypdf.PdfReader so `reader.metadata` is empty.

    reportlab always stamps default /Title, /Producer, etc. onto generated
    PDFs, so without this, rag_upload_file's `[Metadata]\\n...` preamble
    alone is enough non-empty "text" to prevent needs_ocr from ever being
    True and to prevent the "no text extracted" 400 from ever firing --
    masking the OCR and no-text branches entirely.
    """
    mock_reader = Mock()
    mock_reader.metadata = None
    return patch("pypdf.PdfReader", return_value=mock_reader)


def test_rag_upload_pdf_success_via_pdfplumber() -> None:
    content = _build_text_pdf("Hello PDF World")
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("doc.pdf", content, "application/pdf")},
        )

    assert response.status_code == 200
    assert "Hello PDF World" in captured["text"]
    assert "Title: Test PDF" in captured["text"]


def test_rag_upload_pdf_success_extracts_table() -> None:
    content = _build_table_pdf()
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with patch("nyxgpt.app.ingest_document", side_effect=fake_ingest):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("table.pdf", content, "application/pdf")},
        )

    assert response.status_code == 200
    assert "[Table]" in captured["text"]
    assert "H1 | H2" in captured["text"]


def test_rag_upload_pdf_pdfplumber_failure_falls_back_to_pypdf() -> None:
    content = _build_text_pdf("Fallback text")
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    with (
        patch("pdfplumber.open", side_effect=RuntimeError("plumber broke")),
        patch("nyxgpt.app.ingest_document", side_effect=fake_ingest),
    ):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("doc.pdf", content, "application/pdf")},
        )

    assert response.status_code == 200
    assert "Fallback text" in captured["text"]


def test_rag_upload_pdf_ocr_improves_extraction() -> None:
    content = _build_image_only_pdf()
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    fake_image = Mock()
    with (
        _no_pdf_metadata(),
        patch("pdf2image.convert_from_bytes", return_value=[fake_image]),
        patch("pytesseract.image_to_string", return_value="OCR extracted text " * 5),
        patch("nyxgpt.app.ingest_document", side_effect=fake_ingest),
    ):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("scan.pdf", content, "application/pdf")},
        )

    assert response.status_code == 200
    assert "OCR extracted text" in captured["text"]
    assert "[Page 1 (OCR)]" in captured["text"]


def test_rag_upload_pdf_ocr_prepends_metadata_and_honors_tesseract_cmd() -> None:
    """Covers the OCR-block metadata prefix and custom tesseract_cmd config,
    which the other OCR tests deliberately avoid by nulling out metadata."""
    import pytesseract as _pytesseract

    content = _build_image_only_pdf()
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    # Real reportlab-stamped metadata is short but non-empty, so a very high
    # threshold is needed to still trigger needs_ocr despite it being present.
    fake_cfg = _pdf_cfg(ocr_min_text_threshold=100000, tesseract_cmd="/usr/local/bin/tesseract")
    fake_image = Mock()
    original_cmd = _pytesseract.pytesseract.tesseract_cmd
    try:
        with (
            patch("nyxgpt.app.load_config", return_value=fake_cfg),
            patch("pdf2image.convert_from_bytes", return_value=[fake_image]),
            patch("pytesseract.image_to_string", return_value="OCR text with metadata " * 5),
            patch("nyxgpt.app.ingest_document", side_effect=fake_ingest),
        ):
            response = client.post(
                UPLOAD_URL,
                files={"file": ("scan.pdf", content, "application/pdf")},
            )

        assert response.status_code == 200
        assert "OCR text with metadata" in captured["text"]
        assert "[Metadata]" in captured["text"]
        assert _pytesseract.pytesseract.tesseract_cmd == "/usr/local/bin/tesseract"
    finally:
        _pytesseract.pytesseract.tesseract_cmd = original_cmd


def test_rag_upload_pdf_ocr_import_error_other_missing_lib() -> None:
    """The ImportError handler's else branch: a missing OCR dependency whose
    message names neither "pytesseract" nor "pdf2image" by name."""
    content = _build_image_only_pdf()
    client = _client()
    with (
        _no_pdf_metadata(),
        patch("pdf2image.convert_from_bytes", side_effect=ImportError("libpoppler.so not found")),
    ):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("scan.pdf", content, "application/pdf")},
        )

    assert response.status_code == 400
    assert "no text" in response.json()["error"]["message"].lower()


def test_rag_upload_pdf_ocr_does_not_improve_extraction_keeps_standard_text() -> None:
    # 40 real chars clears the OCR trigger (< the default 50-char threshold,
    # so OCR still runs) but the formatted single-page OCR block below
    # ("[Page 1 (OCR)]\nx", 17 chars) is shorter, so it's discarded in favor
    # of the original.
    original_text = "A" * 40
    content = _build_text_pdf(original_text)
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    fake_image = Mock()
    with (
        _no_pdf_metadata(),
        patch("pdf2image.convert_from_bytes", return_value=[fake_image]),
        patch("pytesseract.image_to_string", return_value="x"),
        patch("nyxgpt.app.ingest_document", side_effect=fake_ingest),
    ):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("scan.pdf", content, "application/pdf")},
        )

    assert response.status_code == 200
    assert captured["text"] == original_text
    assert "(OCR)" not in captured["text"]


def test_rag_upload_pdf_ocr_page_failure_is_skipped() -> None:
    content = _build_image_only_pdf()
    client = _client()
    captured: dict[str, Any] = {}

    def fake_ingest(**kwargs: Any) -> dict[str, Any]:
        captured.update(kwargs)
        return _fake_ingest_result()

    fake_image = Mock()
    with (
        _no_pdf_metadata(),
        patch("pdf2image.convert_from_bytes", return_value=[fake_image, fake_image]),
        patch(
            "pytesseract.image_to_string",
            side_effect=[RuntimeError("ocr engine crash"), "second page ok " * 5],
        ),
        patch("nyxgpt.app.ingest_document", side_effect=fake_ingest),
    ):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("scan.pdf", content, "application/pdf")},
        )

    assert response.status_code == 200
    assert "second page ok" in captured["text"]
    assert "[Page 2 (OCR)]" in captured["text"]


def test_rag_upload_pdf_ocr_produces_no_text_keeps_going() -> None:
    content = _build_image_only_pdf()
    client = _client()
    with (
        _no_pdf_metadata(),
        patch("pdf2image.convert_from_bytes", return_value=[Mock()]),
        patch("pytesseract.image_to_string", return_value="   "),
    ):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("scan.pdf", content, "application/pdf")},
        )

    assert response.status_code == 400
    assert "no text" in response.json()["error"]["message"].lower()


def test_rag_upload_pdf_ocr_import_error_pytesseract_missing_skips_ocr() -> None:
    content = _build_image_only_pdf()
    client = _client()
    with _no_pdf_metadata(), patch.dict(sys.modules, {"pytesseract": None}):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("scan.pdf", content, "application/pdf")},
        )

    assert response.status_code == 400
    assert "no text" in response.json()["error"]["message"].lower()


def test_rag_upload_pdf_ocr_import_error_pdf2image_missing_skips_ocr() -> None:
    content = _build_image_only_pdf()
    client = _client()
    with _no_pdf_metadata(), patch.dict(sys.modules, {"pdf2image": None}):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("scan.pdf", content, "application/pdf")},
        )

    assert response.status_code == 400
    assert "no text" in response.json()["error"]["message"].lower()


def test_rag_upload_pdf_ocr_generic_exception_skips_ocr() -> None:
    content = _build_image_only_pdf()
    client = _client()
    with (
        _no_pdf_metadata(),
        patch("pdf2image.convert_from_bytes", side_effect=RuntimeError("poppler missing")),
    ):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("scan.pdf", content, "application/pdf")},
        )

    assert response.status_code == 400
    assert "no text" in response.json()["error"]["message"].lower()


def test_rag_upload_pdf_no_text_without_ocr_returns_400() -> None:
    content = _build_image_only_pdf()
    client = _client()
    fake_cfg = _pdf_cfg(ocr_enabled=False)
    with (
        _no_pdf_metadata(),
        patch("nyxgpt.app.load_config", return_value=fake_cfg),
    ):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("scan.pdf", content, "application/pdf")},
        )

    assert response.status_code == 400
    assert "no text" in response.json()["error"]["message"].lower()


def test_rag_upload_pdf_import_error_returns_400() -> None:
    content = _build_text_pdf()
    client = _client()
    with patch.dict(sys.modules, {"pypdf": None}):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("doc.pdf", content, "application/pdf")},
        )
    assert response.status_code == 400
    assert "PDF support not available" in response.json()["error"]["message"]


def test_rag_upload_pdf_parsing_failure_returns_400() -> None:
    content = _build_text_pdf()
    client = _client()
    with patch("pypdf.PdfReader", side_effect=RuntimeError("malformed pdf")):
        response = client.post(
            UPLOAD_URL,
            files={"file": ("doc.pdf", content, "application/pdf")},
        )
    assert response.status_code == 400
    assert "PDF parsing failed" in response.json()["error"]["message"]


def _pdf_cfg(**pdf_overrides: Any):
    from configparser import ConfigParser

    cfg = ConfigParser()
    defaults = {
        "ocr_enabled": "true",
        "ocr_min_text_threshold": "50",
        "ocr_dpi": "300",
        "ocr_lang": "eng",
        "ocr_psm": "3",
    }
    cfg["pdf"] = {k: str(v).lower() if isinstance(v, bool) else str(v) for k, v in defaults.items()}
    for key, value in pdf_overrides.items():
        cfg["pdf"][key] = str(value).lower() if isinstance(value, bool) else str(value)
    return cfg
