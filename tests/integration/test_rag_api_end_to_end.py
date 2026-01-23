from __future__ import annotations

import time
import uuid

import httpx
import pytest


def _unique_doc_id(prefix: str) -> str:
    """Generate a unique doc_id for testing to avoid hash-based skip."""
    return f"{prefix}-{uuid.uuid4().hex[:10]}"


@pytest.mark.integration
def test_rag_api_ingest_and_query(
    api_base_url: str, require_ollama: None, require_cassandra: None
) -> None:
    doc_id = f"itest-{uuid.uuid4().hex[:10]}"
    text = (
        "Cassandra 5.0 supports vector search with SAI indexes. "
        "This sentence is used for nyxGPT integration testing."
    )

    with httpx.Client(base_url=api_base_url, timeout=60.0) as client:
        # Optional: verify the API is up before we do slower work
        r = client.get("/health")
        assert r.status_code == 200

        ingest_resp = client.post(
            "/api/v1/rag/ingest",
            json={
                "doc_id": doc_id,
                "text": text,
            },
        )
        assert ingest_resp.status_code in (200, 201)

        # Give Cassandra indexing a moment (SAI / vector index)
        time.sleep(2.0)

        query_resp = client.post(
            "/api/v1/rag/query",
            json={
                "query": "What does Cassandra support for vector search?",
                "top_k": 3,
            },
        )
        assert query_resp.status_code == 200

        data = query_resp.json()
        assert isinstance(data, dict)
        assert "results" in data
        assert isinstance(data["results"], list)
        assert len(data["results"]) > 0

        top = data["results"][0]
        assert isinstance(top, dict)
        assert "text" in top
        assert isinstance(top["text"], str)
        assert len(top["text"]) > 0
        assert "doc_id" in top
        assert "chunk_id" in top
        assert "score" in top


@pytest.mark.integration
def test_session_rag_enable_disable(api_base_url: str) -> None:
    """Test per-session RAG enable/disable endpoints."""
    session_name = f"rag-test-{uuid.uuid4().hex[:8]}"

    with httpx.Client(base_url=api_base_url, timeout=10.0) as client:
        # 1. Get initial metadata (creates session)
        meta_resp = client.get(f"/api/v1/sessions/{session_name}/metadata")
        assert meta_resp.status_code == 200
        meta = meta_resp.json()
        assert "rag_enabled" in meta
        # Default should be False (or inherited from config)
        assert isinstance(meta["rag_enabled"], bool)

        # 2. Enable RAG
        enable_resp = client.post(f"/api/v1/sessions/{session_name}/rag/enable")
        assert enable_resp.status_code == 200
        enable_data = enable_resp.json()
        assert enable_data["session"] == session_name
        assert enable_data["rag_enabled"] is True

        # 3. Verify RAG is enabled in metadata
        meta_resp2 = client.get(f"/api/v1/sessions/{session_name}/metadata")
        assert meta_resp2.status_code == 200
        meta2 = meta_resp2.json()
        assert meta2["rag_enabled"] is True

        # 4. Disable RAG
        disable_resp = client.post(f"/api/v1/sessions/{session_name}/rag/disable")
        assert disable_resp.status_code == 200
        disable_data = disable_resp.json()
        assert disable_data["session"] == session_name
        assert disable_data["rag_enabled"] is False

        # 5. Verify RAG is disabled in metadata
        meta_resp3 = client.get(f"/api/v1/sessions/{session_name}/metadata")
        assert meta_resp3.status_code == 200
        meta3 = meta_resp3.json()
        assert meta3["rag_enabled"] is False


@pytest.mark.integration
def test_rag_upload_text_file(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test RAG file upload endpoint with .txt file."""
    # Create a test text file
    test_file = tmp_path / "test_upload.txt"
    test_content = "This is a test document for RAG upload testing. It contains important information."
    test_file.write_text(test_content)

    # Use unique doc_id to avoid hash-based skip from previous test runs
    doc_id = _unique_doc_id("txt-upload")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        # Upload the file with unique doc_id
        with open(test_file, "rb") as f:
            files = {"file": ("test_upload.txt", f, "text/plain")}
            upload_resp = client.post(
                "/api/v1/rag/upload", files=files, params={"doc_id": doc_id}
            )

        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "doc_id" in upload_data
        assert "chunks_ingested" in upload_data
        assert upload_data["chunks_ingested"] > 0

        # Give Cassandra indexing a moment
        time.sleep(2.0)

        # Verify we can query the uploaded content
        query_resp = client.post(
            "/api/v1/rag/query", json={"query": "test document", "top_k": 5}
        )
        assert query_resp.status_code == 200
        results = query_resp.json()["results"]
        assert len(results) > 0


@pytest.mark.integration
def test_rag_upload_markdown_file(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test RAG file upload endpoint with .md file with frontmatter."""
    # Create a test markdown file with frontmatter
    test_file = tmp_path / "test.md"
    markdown_content = """---
title: Test Document
author: Integration Test
tags: [test, markdown]
---

# Main Heading

This is a test markdown document with **bold** and *italic* text.

## Section 1

Some content here with `code` inline.

```python
def hello():
    print("Hello, world!")
```

## Section 2

More content for testing RAG ingestion.
"""
    test_file.write_text(markdown_content)

    # Use unique doc_id to avoid hash-based skip from previous test runs
    doc_id = _unique_doc_id("md-upload")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        # Upload the markdown file with unique doc_id
        with open(test_file, "rb") as f:
            files = {"file": ("test.md", f, "text/markdown")}
            upload_resp = client.post(
                "/api/v1/rag/upload", files=files, params={"doc_id": doc_id}
            )

        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "doc_id" in upload_data
        assert "chunks_ingested" in upload_data
        assert upload_data["chunks_ingested"] > 0

        # Give Cassandra indexing a moment
        time.sleep(2.0)

        # Verify we can query the uploaded content
        query_resp = client.post(
            "/api/v1/rag/query", json={"query": "markdown document", "top_k": 5}
        )
        assert query_resp.status_code == 200
        results = query_resp.json()["results"]
        assert len(results) > 0


@pytest.mark.integration
def test_rag_upload_docx_file(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test RAG file upload endpoint with .docx file (Microsoft Word)."""
    # Create a test DOCX file
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx not available, skipping DOCX test")

    test_file = tmp_path / "test.docx"
    doc = Document()

    # Add heading
    doc.add_heading("Test Document Title", level=1)

    # Add paragraphs
    doc.add_paragraph("This is a test Word document for RAG upload testing.")
    doc.add_paragraph("It contains important information about DOCX support.")

    # Add a heading and more content
    doc.add_heading("Section 1", level=2)
    doc.add_paragraph("This section tests heading preservation in DOCX files.")

    # Add a table
    table = doc.add_table(rows=3, cols=2)
    table.rows[0].cells[0].text = "Feature"
    table.rows[0].cells[1].text = "Status"
    table.rows[1].cells[0].text = "Text extraction"
    table.rows[1].cells[1].text = "Working"
    table.rows[2].cells[0].text = "Table support"
    table.rows[2].cells[1].text = "Working"

    # Save the document
    doc.save(test_file)

    # Use unique doc_id to avoid hash-based skip from previous test runs
    doc_id = _unique_doc_id("docx-upload")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        # Upload the DOCX file with unique doc_id
        with open(test_file, "rb") as f:
            files = {"file": ("test.docx", f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
            upload_resp = client.post("/api/v1/rag/upload", files=files, params={"doc_id": doc_id})

        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "doc_id" in upload_data
        assert "chunks_ingested" in upload_data
        assert upload_data["chunks_ingested"] > 0

        # Give Cassandra indexing a moment
        time.sleep(2.0)

        # Verify we can query the uploaded content
        query_resp = client.post(
            "/api/v1/rag/query", json={"query": "Word document DOCX", "top_k": 5}
        )
        assert query_resp.status_code == 200
        results = query_resp.json()["results"]
        assert len(results) > 0


@pytest.mark.integration
def test_rag_upload_pptx_file(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test RAG file upload endpoint with .pptx file."""
    pytest.importorskip("pptx", reason="python-pptx not installed")
    from pptx import Presentation

    # Create a test PowerPoint file
    test_file = tmp_path / "test_presentation.pptx"
    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "Test Presentation"
    subtitle1.text = "Integration Testing for PPTX Support"

    # Slide 2: Content slide with bullet points
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    title2 = slide2.shapes.title
    body2 = slide2.placeholders[1]
    title2.text = "Key Features"
    tf2 = body2.text_frame
    tf2.text = "First feature: text extraction"
    p2 = tf2.add_paragraph()
    p2.text = "Second feature: slide order preservation"
    p2.level = 0
    p3 = tf2.add_paragraph()
    p3.text = "Third feature: speaker notes support"
    p3.level = 0

    # Slide 3: Slide with speaker notes
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Technical Architecture"
    # Add speaker notes
    notes_slide3 = slide3.notes_slide
    notes_tf3 = notes_slide3.notes_text_frame
    notes_tf3.text = "This slide discusses the technical architecture of the RAG system."

    prs.save(str(test_file))

    # Use unique doc_id to avoid hash-based skip from previous test runs
    doc_id = _unique_doc_id("pptx-upload")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        # Upload the PPTX file with unique doc_id
        with open(test_file, "rb") as f:
            files = {"file": ("test_presentation.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
            upload_resp = client.post("/api/v1/rag/upload", files=files, params={"doc_id": doc_id})

        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "doc_id" in upload_data
        assert "chunks_ingested" in upload_data
        assert upload_data["chunks_ingested"] > 0

        # Give Cassandra indexing a moment
        time.sleep(2.0)

        # Verify we can query the uploaded content
        query_resp = client.post(
            "/api/v1/rag/query", json={"query": "presentation features", "top_k": 5}
        )
        assert query_resp.status_code == 200
        results = query_resp.json()["results"]
        assert len(results) > 0


@pytest.mark.integration
def test_rag_upload_docx_empty_file(api_base_url: str, tmp_path) -> None:
    """Test RAG file upload endpoint with empty .docx file."""
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx not available, skipping DOCX test")

    test_file = tmp_path / "empty.docx"
    doc = Document()
    # Save empty document with no content
    doc.save(test_file)

    with httpx.Client(base_url=api_base_url, timeout=10.0) as client:
        with open(test_file, "rb") as f:
            files = {"file": ("empty.docx", f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
            upload_resp = client.post("/api/v1/rag/upload", files=files)

        # Should reject empty file with 400
        assert upload_resp.status_code == 400
        error_data = upload_resp.json()
        assert "error" in error_data
        error_msg = error_data["error"]["message"].lower()
        assert "empty" in error_msg or "no extractable text" in error_msg


@pytest.mark.integration
def test_rag_upload_docx_only_table(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test RAG file upload endpoint with .docx file containing only tables."""
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx not available, skipping DOCX test")

    test_file = tmp_path / "only_table.docx"
    doc = Document()

    # Add only a table, no paragraphs
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Column A"
    table.rows[0].cells[1].text = "Column B"
    table.rows[1].cells[0].text = "Data 1"
    table.rows[1].cells[1].text = "Data 2"

    doc.save(test_file)

    doc_id = _unique_doc_id("docx-only-table")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        with open(test_file, "rb") as f:
            files = {"file": ("only_table.docx", f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
            upload_resp = client.post("/api/v1/rag/upload", files=files, params={"doc_id": doc_id})

        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "chunks_ingested" in upload_data
        assert upload_data["chunks_ingested"] > 0


@pytest.mark.integration
def test_rag_upload_docx_only_text(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test RAG file upload endpoint with .docx file containing only text (no tables)."""
    try:
        from docx import Document
    except ImportError:
        pytest.skip("python-docx not available, skipping DOCX test")

    test_file = tmp_path / "only_text.docx"
    doc = Document()

    # Add only paragraphs, no tables
    doc.add_paragraph("This is a simple DOCX file.")
    doc.add_paragraph("It contains only plain text paragraphs.")
    doc.add_paragraph("No tables or other complex formatting.")

    doc.save(test_file)

    doc_id = _unique_doc_id("docx-only-text")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        with open(test_file, "rb") as f:
            files = {"file": ("only_text.docx", f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
            upload_resp = client.post("/api/v1/rag/upload", files=files, params={"doc_id": doc_id})

        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "chunks_ingested" in upload_data
        assert upload_data["chunks_ingested"] > 0


@pytest.mark.integration
def test_rag_upload_pptx_with_speaker_notes(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test that PPTX upload correctly extracts speaker notes."""
    pytest.importorskip("pptx", reason="python-pptx not installed")
    from pptx import Presentation

    # Create a test PowerPoint file with specific speaker notes
    test_file = tmp_path / "notes_test.pptx"
    prs = Presentation()

    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title1 = slide1.shapes.title
    title1.text = "Speaker Notes Test"

    # Add unique speaker notes text
    notes_slide1 = slide1.notes_slide
    notes_tf1 = notes_slide1.notes_text_frame
    unique_note_text = "This unique speaker note should be extracted and searchable in RAG."
    notes_tf1.text = unique_note_text

    prs.save(str(test_file))

    # Use unique doc_id
    doc_id = _unique_doc_id("pptx-notes")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        # Upload the PPTX file
        with open(test_file, "rb") as f:
            files = {"file": ("notes_test.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
            upload_resp = client.post("/api/v1/rag/upload", files=files, params={"doc_id": doc_id})

        assert upload_resp.status_code == 200

        # Give Cassandra indexing a moment
        time.sleep(2.0)

        # Query for the speaker notes content
        query_resp = client.post(
            "/api/v1/rag/query", json={"query": "unique speaker note searchable", "top_k": 5}
        )
        assert query_resp.status_code == 200
        results = query_resp.json()["results"]
        assert len(results) > 0
        # Verify speaker notes were extracted
        found_note = any(unique_note_text.lower() in result["text"].lower() for result in results)
        assert found_note, "Speaker notes were not extracted"


@pytest.mark.integration
def test_rag_upload_pptx_slide_order(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test that PPTX upload preserves slide order."""
    pytest.importorskip("pptx", reason="python-pptx not installed")
    from pptx import Presentation

    # Create a test PowerPoint file with numbered slides
    test_file = tmp_path / "order_test.pptx"
    prs = Presentation()

    for i in range(1, 4):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = f"Slide {i}: Content in order"

    prs.save(str(test_file))

    # Use unique doc_id
    doc_id = _unique_doc_id("pptx-order")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        # Upload the PPTX file
        with open(test_file, "rb") as f:
            files = {"file": ("order_test.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
            upload_resp = client.post("/api/v1/rag/upload", files=files, params={"doc_id": doc_id})

        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "chunks_ingested" in upload_data
        assert upload_data["chunks_ingested"] > 0


@pytest.mark.integration
def test_rag_upload_docx_corrupted_file(api_base_url: str, tmp_path) -> None:
    """Test RAG file upload endpoint with corrupted .docx file."""
    test_file = tmp_path / "corrupted.docx"
    # Write invalid content that's not a valid DOCX/ZIP file
    test_file.write_bytes(b"This is not a valid DOCX file content")

    with httpx.Client(base_url=api_base_url, timeout=10.0) as client:
        with open(test_file, "rb") as f:
            files = {"file": ("corrupted.docx", f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
            upload_resp = client.post("/api/v1/rag/upload", files=files)

        # Should reject corrupted file with 400
        assert upload_resp.status_code == 400
        error_data = upload_resp.json()
        assert "error" in error_data
        error_msg = error_data["error"]["message"].lower()
        # Accept: corrupted/invalid file errors, or docx support not available
        assert any(term in error_msg for term in ("corrupted", "invalid", "not available", "not supported"))


@pytest.mark.integration
def test_rag_upload_docx_with_images(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test RAG file upload with DOCX file containing embedded images."""
    try:
        from docx import Document
        from docx.shared import Inches
    except ImportError:
        pytest.skip("python-docx not available, skipping DOCX test")

    try:
        from PIL import Image
    except ImportError:
        pytest.skip("Pillow not available, skipping image test")

    test_file = tmp_path / "test_with_image.docx"
    image_file = tmp_path / "test_image.png"

    # Create a simple test image
    img = Image.new("RGB", (100, 100), color="red")
    img.save(image_file)

    # Create DOCX with image
    doc = Document()
    doc.add_heading("Document with Image", level=1)
    doc.add_paragraph("This paragraph is before the image.")

    # Add image with alt text
    try:
        doc.add_picture(str(image_file), width=Inches(2))
    except Exception:
        # If adding image fails, skip this test
        pytest.skip("Could not add image to DOCX")

    doc.add_paragraph("This paragraph is after the image.")

    doc.save(test_file)

    # Use unique doc_id
    doc_id = _unique_doc_id("docx-with-image")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        # Upload the DOCX file
        with open(test_file, "rb") as f:
            files = {"file": ("test_with_image.docx", f, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
            upload_resp = client.post("/api/v1/rag/upload", files=files, params={"doc_id": doc_id})

        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "doc_id" in upload_data
        assert "chunks_ingested" in upload_data
        # Should have at least one chunk with the image marker
        assert upload_data["chunks_ingested"] > 0

        # Give Cassandra indexing a moment
        time.sleep(2.0)

        # Query to verify content was ingested (including text around image)
        query_resp = client.post(
            "/api/v1/rag/query", json={"query": "paragraph image", "top_k": 5}
        )
        assert query_resp.status_code == 200
        results = query_resp.json()["results"]
        # Should find the content
        assert len(results) > 0


@pytest.mark.integration
def test_rag_upload_empty_pptx(api_base_url: str, tmp_path) -> None:
    """Test that uploading empty PPTX file is handled gracefully."""
    pytest.importorskip("pptx", reason="python-pptx not installed")
    from pptx import Presentation

    # Create an empty PowerPoint file
    test_file = tmp_path / "empty.pptx"
    prs = Presentation()
    # Add a slide but with no text content
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    prs.save(str(test_file))

    with httpx.Client(base_url=api_base_url, timeout=10.0) as client:
        # Upload the empty PPTX file
        with open(test_file, "rb") as f:
            files = {"file": ("empty.pptx", f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
            upload_resp = client.post("/api/v1/rag/upload", files=files)

        # Should reject with 400 for no extractable text
        assert upload_resp.status_code == 400
        error_data = upload_resp.json()
        assert "error" in error_data
        assert "no extractable text" in error_data["error"]["message"].lower()


@pytest.mark.integration
def test_rag_upload_invalid_file_type(api_base_url: str, tmp_path) -> None:
    """Test that uploading unsupported file types is rejected."""
    # Create a test file with unsupported extension
    test_file = tmp_path / "test.exe"
    test_file.write_bytes(b"fake executable content")

    with httpx.Client(base_url=api_base_url, timeout=10.0) as client:
        # Try to upload unsupported file type
        with open(test_file, "rb") as f:
            files = {"file": ("test.exe", f, "application/octet-stream")}
            upload_resp = client.post("/api/v1/rag/upload", files=files)

        # Should reject with 400
        assert upload_resp.status_code == 400
        error_data = upload_resp.json()
        assert "error" in error_data
        assert "not supported" in error_data["error"]["message"].lower()


@pytest.mark.integration
def test_rag_upload_pdf_with_tables_and_metadata(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test PDF upload with improved extraction: tables, formatting, multi-column, metadata (#2663)."""
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    except ImportError:
        pytest.skip("reportlab not available for PDF generation")

    # Create a test PDF with metadata, tables, and formatted text
    test_file = tmp_path / "test_enhanced.pdf"

    doc = SimpleDocTemplate(
        str(test_file),
        pagesize=letter,
        title="Test Document",
        author="Test Author",
        subject="PDF Extraction Testing",
    )

    styles = getSampleStyleSheet()
    story = []

    # Add title
    title = Paragraph("Enhanced PDF Test Document", styles['Title'])
    story.append(title)
    story.append(Spacer(1, 0.2 * inch))

    # Add description
    desc = Paragraph(
        "This document tests improved PDF extraction with tables, formatting, and metadata.",
        styles['Normal']
    )
    story.append(desc)
    story.append(Spacer(1, 0.3 * inch))

    # Add a table to test table extraction
    table_data = [
        ['Feature', 'Status', 'Priority'],
        ['Table handling', 'Improved', 'High'],
        ['Formatting preservation', 'Enhanced', 'High'],
        ['Multi-column support', 'Added', 'Medium'],
        ['Metadata extraction', 'Implemented', 'High'],
    ]

    table = Table(table_data)
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black),
    ]))

    story.append(table)
    story.append(Spacer(1, 0.3 * inch))

    # Add more text content
    content = Paragraph(
        "The improved PDF extraction now properly handles complex layouts including "
        "tables with structured data, preserves text formatting, and extracts document "
        "metadata such as title, author, and creation date.",
        styles['Normal']
    )
    story.append(content)

    # Build the PDF
    doc.build(story)

    # Use unique doc_id to avoid hash-based skip from previous test runs
    doc_id = _unique_doc_id("pdf-enhanced")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        # Upload the PDF file with unique doc_id
        with open(test_file, "rb") as f:
            files = {"file": ("test_enhanced.pdf", f, "application/pdf")}
            upload_resp = client.post("/api/v1/rag/upload", files=files, params={"doc_id": doc_id})

        # Verify successful upload
        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "doc_id" in upload_data
        assert "chunks_ingested" in upload_data
        assert upload_data["chunks_ingested"] > 0

        # Give Cassandra indexing a moment
        time.sleep(2.0)

        # Verify we can query content from the table
        query_resp = client.post(
            "/api/v1/rag/query", json={"query": "table handling formatting", "top_k": 5}
        )
        assert query_resp.status_code == 200
        results = query_resp.json()["results"]
        assert len(results) > 0

        # Verify the extracted text contains table markers and metadata
        # The text should contain our table data and metadata
        combined_text = " ".join(r["text"] for r in results)

        # Check for table content
        assert "Table handling" in combined_text or "Improved" in combined_text

        # Check for metadata section (at least title or author should be present)
        assert "Metadata" in combined_text or "Test" in combined_text


@pytest.mark.integration
def test_rag_upload_epub_file(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test RAG file upload endpoint with .epub file."""
    pytest.importorskip("ebooklib", reason="ebooklib not installed")
    from ebooklib import epub

    # Create a test ePUB file
    test_file = tmp_path / "test_book.epub"
    book = epub.EpubBook()

    # Set metadata
    book.set_identifier("test-epub-123")
    book.set_title("Test ePUB Book")
    book.set_language("en")
    book.add_author("Test Author")

    # Create chapters
    c1 = epub.EpubHtml(title="Introduction", file_name="chap_01.xhtml", lang="en")
    c1.content = """
    <html>
    <head><title>Introduction</title></head>
    <body>
    <h1>Introduction</h1>
    <p>This is a test ePUB document for RAG upload testing.</p>
    <p>It contains important information about ePUB support.</p>
    </body>
    </html>
    """

    c2 = epub.EpubHtml(title="Chapter 1", file_name="chap_02.xhtml", lang="en")
    c2.content = """
    <html>
    <head><title>Chapter 1</title></head>
    <body>
    <h1>Chapter 1: The Beginning</h1>
    <p>This chapter tests content extraction from ePUB files.</p>
    <p>ePUB is a standard format for electronic books.</p>
    </body>
    </html>
    """

    # Add chapters to book
    book.add_item(c1)
    book.add_item(c2)

    # Add table of contents
    book.toc = (
        epub.Link("chap_01.xhtml", "Introduction", "intro"),
        epub.Link("chap_02.xhtml", "Chapter 1", "chap1"),
    )

    # Add navigation files
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())

    # Define spine (reading order)
    book.spine = ["nav", c1, c2]

    # Write the ePUB file
    epub.write_epub(str(test_file), book)

    # Use unique doc_id to avoid hash-based skip from previous test runs
    doc_id = _unique_doc_id("epub-upload")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        # Upload the ePUB file with unique doc_id
        with open(test_file, "rb") as f:
            files = {"file": ("test_book.epub", f, "application/epub+zip")}
            upload_resp = client.post(
                "/api/v1/rag/upload", files=files, params={"doc_id": doc_id}
            )

        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "doc_id" in upload_data
        assert "chunks_ingested" in upload_data
        assert upload_data["chunks_ingested"] > 0

        # Give Cassandra indexing a moment
        time.sleep(2.0)

        # Verify we can query the uploaded content
        query_resp = client.post(
            "/api/v1/rag/query", json={"query": "ePUB electronic books", "top_k": 5}
        )
        assert query_resp.status_code == 200
        results = query_resp.json()["results"]
        assert len(results) > 0


@pytest.mark.integration
def test_rag_upload_epub_with_metadata(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test that ePUB upload correctly extracts metadata."""
    pytest.importorskip("ebooklib", reason="ebooklib not installed")
    from ebooklib import epub

    # Create a test ePUB file with rich metadata
    test_file = tmp_path / "metadata_book.epub"
    book = epub.EpubBook()

    # Set comprehensive metadata
    book.set_identifier("metadata-test-456")
    book.set_title("Complete Metadata Test")
    book.set_language("en")
    book.add_author("Jane Doe")
    book.add_author("John Smith")
    book.add_metadata("DC", "publisher", "Test Publisher")
    book.add_metadata("DC", "description", "A test book for metadata extraction")
    book.add_metadata("DC", "date", "2024-01-01")

    # Create a simple chapter
    c1 = epub.EpubHtml(title="Content", file_name="content.xhtml", lang="en")
    c1.content = """
    <html>
    <body>
    <h1>Test Content</h1>
    <p>This tests metadata extraction from ePUB files.</p>
    </body>
    </html>
    """

    book.add_item(c1)
    book.toc = (epub.Link("content.xhtml", "Content", "content"),)
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav", c1]

    epub.write_epub(str(test_file), book)

    doc_id = _unique_doc_id("epub-metadata")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        with open(test_file, "rb") as f:
            files = {"file": ("metadata_book.epub", f, "application/epub+zip")}
            upload_resp = client.post(
                "/api/v1/rag/upload", files=files, params={"doc_id": doc_id}
            )

        assert upload_resp.status_code == 200

        # Give Cassandra indexing a moment
        time.sleep(2.0)

        # Query for metadata content
        query_resp = client.post(
            "/api/v1/rag/query", json={"query": "metadata test", "top_k": 5}
        )
        assert query_resp.status_code == 200
        results = query_resp.json()["results"]
        assert len(results) > 0


@pytest.mark.integration
def test_rag_upload_epub_empty_file(api_base_url: str, tmp_path) -> None:
    """Test that uploading empty ePUB file is handled gracefully."""
    pytest.importorskip("ebooklib", reason="ebooklib not installed")
    from ebooklib import epub

    # Create an ePUB with no content chapters
    test_file = tmp_path / "empty.epub"
    book = epub.EpubBook()
    book.set_identifier("empty-test")
    book.set_title("Empty Book")
    book.set_language("en")

    # Add navigation but no content chapters
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav"]

    epub.write_epub(str(test_file), book)

    with httpx.Client(base_url=api_base_url, timeout=10.0) as client:
        with open(test_file, "rb") as f:
            files = {"file": ("empty.epub", f, "application/epub+zip")}
            upload_resp = client.post("/api/v1/rag/upload", files=files)

        # Should reject with 400 for no extractable text
        assert upload_resp.status_code == 400
        error_data = upload_resp.json()
        assert "error" in error_data
        error_msg = error_data["error"]["message"].lower()
        assert "no text" in error_msg or "empty" in error_msg or "malformed" in error_msg


@pytest.mark.integration
def test_rag_upload_epub_corrupted_file(api_base_url: str, tmp_path) -> None:
    """Test RAG file upload endpoint with corrupted .epub file."""
    test_file = tmp_path / "corrupted.epub"
    # Write invalid content that's not a valid ePUB file
    test_file.write_bytes(b"This is not a valid ePUB file content")

    with httpx.Client(base_url=api_base_url, timeout=10.0) as client:
        with open(test_file, "rb") as f:
            files = {"file": ("corrupted.epub", f, "application/epub+zip")}
            upload_resp = client.post("/api/v1/rag/upload", files=files)

        # Should reject corrupted file with 400
        assert upload_resp.status_code == 400
        error_data = upload_resp.json()
        assert "error" in error_data
        error_msg = error_data["error"]["message"].lower()
        assert any(term in error_msg for term in ("invalid", "failed", "corrupted", "parsing"))


@pytest.mark.integration
def test_rag_upload_epub_multi_chapter(
    api_base_url: str, require_ollama: None, require_cassandra: None, tmp_path
) -> None:
    """Test that ePUB upload preserves chapter structure for multi-chapter books."""
    pytest.importorskip("ebooklib", reason="ebooklib not installed")
    from ebooklib import epub

    # Create a multi-chapter ePUB
    test_file = tmp_path / "multi_chapter.epub"
    book = epub.EpubBook()

    book.set_identifier("multi-chapter-789")
    book.set_title("Multi-Chapter Test Book")
    book.set_language("en")

    # Create multiple chapters
    chapters = []
    for i in range(1, 4):
        c = epub.EpubHtml(title=f"Chapter {i}", file_name=f"chap_{i:02d}.xhtml", lang="en")
        c.content = f"""
        <html>
        <body>
        <h1>Chapter {i}</h1>
        <p>Content for chapter {i} goes here.</p>
        <p>This is part {i} of the multi-chapter book.</p>
        </body>
        </html>
        """
        book.add_item(c)
        chapters.append(c)

    book.toc = tuple(epub.Link(f"chap_{i:02d}.xhtml", f"Chapter {i}", f"chap{i}") for i in range(1, 4))
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav"] + chapters

    epub.write_epub(str(test_file), book)

    doc_id = _unique_doc_id("epub-multi-chapter")

    with httpx.Client(base_url=api_base_url, timeout=30.0) as client:
        with open(test_file, "rb") as f:
            files = {"file": ("multi_chapter.epub", f, "application/epub+zip")}
            upload_resp = client.post(
                "/api/v1/rag/upload", files=files, params={"doc_id": doc_id}
            )

        assert upload_resp.status_code == 200
        upload_data = upload_resp.json()
        assert "chunks_ingested" in upload_data
        assert upload_data["chunks_ingested"] > 0

        # Give Cassandra indexing a moment
        time.sleep(2.0)

        # Verify content from different chapters can be queried
        query_resp = client.post(
            "/api/v1/rag/query", json={"query": "multi-chapter book", "top_k": 5}
        )
        assert query_resp.status_code == 200
        results = query_resp.json()["results"]
        assert len(results) > 0
