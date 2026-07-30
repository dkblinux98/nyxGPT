"""Core chat orchestration: building context, calling the model, and persisting turns.

Ties together session state, RAG retrieval, prompt-mode selection, message
budget truncation, and the response cache to produce a single chat reply
(`chat`) or a streamed one (`chat_stream`).
"""

from __future__ import annotations

import json
import logging
import time
import uuid
from collections.abc import Callable, Iterator
from configparser import ConfigParser
from dataclasses import dataclass
from datetime import UTC, datetime
from typing import Any, cast

from nyxgpt.cache import CacheBackend, DiskCache, MemoryCache, NoOpCache, hash_text
from nyxgpt.config import (
    get_context_warning_threshold,
    get_context_window_size,
    get_prompt_mode_enabled,
    get_prompt_mode_long_threshold,
    get_prompt_mode_short_threshold,
    get_rag_context_format,
    get_rag_enabled,
    get_rag_instruction_template,
    get_system_prompt_minimize,
    load_config,
)
from nyxgpt.ollama_client import ollama_chat, ollama_chat_stream_tokens
from nyxgpt.rag.rag import compose_context, retrieve_context
from nyxgpt.rag.vectorstore_cassandra import MetadataFilter
from nyxgpt.sessions import load_session, save_session
from nyxgpt.token_counter import count_message_tokens, count_tokens

logger = logging.getLogger(__name__)


def _safe_token_count(text: str) -> int | None:
    """Return `count_tokens(text)`, or None if the tokenizer isn't available.

    Token counts are a "nice to have" on the chat lifecycle log record, not
    load-bearing -- an unavailable tokenizer must never break the log line.
    """
    try:
        return count_tokens(text)
    except Exception:
        return None


def _safe_message_token_count(messages: list[dict[str, Any]]) -> int | None:
    """Return `count_message_tokens(messages)`, or None if unavailable (see `_safe_token_count`)."""
    try:
        return count_message_tokens(messages)
    except Exception:
        return None


# Global response cache instance (initialized lazily)
_response_cache: CacheBackend[str] | None = None


@dataclass
class ChatResult:
    """Result of a completed (non-streaming) chat exchange."""

    session: str
    model: str
    reply: str
    rag_used: bool
    rag_chunks: int
    rag_context: list[dict] | None = (
        None  # List of {text, score, doc_id, chunk_id, similarity_score}
    )


@dataclass
class ChatContext:
    """Prepared context for a chat interaction.

    This consolidates all the setup work shared between
    streaming and non-streaming chat.
    """

    messages: list[dict[str, Any]]
    state: Any  # SessionState
    chosen_model: str
    base_url: str
    chat_timeout_s: int
    rag_used: bool
    rag_chunks: int
    rag_context: list[dict] | None = None  # RAG retrieval results
    output_format: dict[str, Any] | None = None  # JSON schema for structured output


def _cfg(config_path: str | None) -> Any:
    """Load the config from `config_path`, or the default config if None."""
    return load_config(config_path)


def _get_bool(cfg: Any, section: str, key: str, default: bool) -> bool:
    """Read a boolean config value, tolerating missing/non-standard representations."""
    try:
        result: bool = cfg.getboolean(section, key, fallback=default)
        return result
    except Exception:
        # Robust against missing/invalid types
        try:
            v = cfg.get(section, key, fallback=str(default)).strip().lower()
            return v in {"1", "true", "yes", "on"}
        except Exception:
            return default


def _get_int(cfg: Any, section: str, key: str, default: int) -> int:
    """Read an integer config value, falling back to `default` if missing/invalid."""
    try:
        result: int = cfg.getint(section, key, fallback=default)
        return result
    except Exception:
        try:
            return int(cfg.get(section, key, fallback=str(default)))
        except Exception:
            return default


def _get_str(cfg: Any, section: str, key: str, default: str) -> str:
    """Read a string config value, falling back to `default` if missing."""
    try:
        result: str = cfg.get(section, key, fallback=default)
        return result
    except Exception:
        return default


def _get_response_cache() -> CacheBackend[str]:
    """Get or initialize the global response cache.

    Returns:
        Initialized cache backend based on config settings
    """
    global _response_cache

    if _response_cache is not None:
        return _response_cache

    # Load config to determine cache settings
    cfg = load_config(None)

    # Check if response cache is enabled
    cache_enabled = _get_bool(cfg, "cache", "response_cache_enabled", False)

    if not cache_enabled:
        logger.debug("Response cache disabled")
        _response_cache = NoOpCache()
        return _response_cache

    # Get cache backend type
    cache_backend = _get_str(cfg, "cache", "response_cache_backend", "memory").lower()

    if cache_backend == "memory":
        max_size = _get_int(cfg, "cache", "response_cache_max_size", 100)
        ttl = _get_int(cfg, "cache", "response_cache_ttl_seconds", 1800)
        _response_cache = MemoryCache(max_size=max_size, default_ttl=ttl, name="chat_response")
        logger.debug(f"Response cache initialized: memory (max_size={max_size}, ttl={ttl}s)")

    elif cache_backend == "disk":
        cache_dir = _get_str(cfg, "cache", "response_cache_dir", "~/.nyxGPT/cache/responses")
        ttl = _get_int(cfg, "cache", "response_cache_ttl_seconds", 3600)
        _response_cache = DiskCache(cache_dir=cache_dir, default_ttl=ttl, name="chat_response")
        logger.debug(f"Response cache initialized: disk (dir={cache_dir}, ttl={ttl}s)")

    else:
        logger.warning(f"Unknown cache backend '{cache_backend}', disabling cache")
        _response_cache = NoOpCache()

    return _response_cache


def clear_response_cache() -> None:
    """Clear the global response cache.

    This is useful for testing or when you want to force fresh responses.
    """
    global _response_cache
    if _response_cache is not None:
        _response_cache.clear()
        logger.info("Response cache cleared")


def _minimize_system_prompt(prompt: str) -> str:
    """Minimize system prompt to reduce token usage.

    Applies conservative optimizations that preserve semantic meaning:
    - Normalizes whitespace (removes extra spaces, newlines)
    - Removes redundant filler words and phrases
    - Condenses common verbose patterns

    Args:
        prompt: Original system prompt text

    Returns:
        Minimized version of the prompt
    """
    import re

    if not prompt or not prompt.strip():
        return prompt

    text = prompt.strip()

    # Normalize whitespace: collapse multiple spaces/newlines
    text = re.sub(r"\s+", " ", text)

    # Remove redundant courtesies and filler words (case-insensitive)
    # These patterns preserve meaning while reducing tokens
    patterns = [
        (r"\bPlease\s+", ""),  # "Please respond" -> "Respond"
        (r"\bYou are\s+", ""),  # "You are a helpful assistant" -> "helpful assistant"
        (r"\bYou should\s+", ""),  # "You should answer" -> "Answer"
        (r"\bI want you to\s+", ""),  # "I want you to act" -> "Act"
        (r"\bYour task is to\s+", ""),  # "Your task is to help" -> "Help"
        (r"\bMake sure to\s+", ""),  # "Make sure to respond" -> "Respond"
        (r"\bBe sure to\s+", ""),  # "Be sure to answer" -> "Answer"
        (r"\s+in order to\s+", " to "),  # "do X in order to Y" -> "do X to Y"
        (r"\s+as well as\s+", " and "),  # "X as well as Y" -> "X and Y"
    ]

    for pattern, replacement in patterns:
        text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)

    # Clean up any double spaces created by replacements
    text = re.sub(r"\s{2,}", " ", text).strip()

    logger.debug(
        "Minimized system prompt: %d -> %d chars (%.1f%% reduction)",
        len(prompt),
        len(text),
        100 * (1 - len(text) / len(prompt)) if len(prompt) > 0 else 0,
    )

    return text


def _detect_prompt_mode(message_count: int, cfg: Any) -> str:
    """Detect appropriate prompt mode based on conversation length.

    Args:
        message_count: Number of messages in conversation (excluding system)
        cfg: Config instance

    Returns:
        One of: "short", "medium", "long"
    """
    short_threshold = get_prompt_mode_short_threshold(cfg)
    long_threshold = get_prompt_mode_long_threshold(cfg)

    if message_count < short_threshold:
        return "short"
    elif message_count >= long_threshold:
        return "long"
    else:
        return "medium"


def _get_prompt_template(mode: str) -> str:
    """Get system prompt template for the given mode.

    Args:
        mode: One of "short", "medium", "long"

    Returns:
        Prompt template string appropriate for the mode
    """
    templates = {
        "short": ("You are a helpful AI assistant. Provide clear, concise responses."),
        "medium": (
            "You are a helpful AI assistant engaged in a conversation. "
            "Provide informative responses while maintaining context from previous messages. "
            "Be concise but thorough when needed."
        ),
        "long": (
            "You are a helpful AI assistant engaged in an extended conversation. "
            "Maintain awareness of the full conversation history and refer back to earlier points when relevant. "
            "Provide comprehensive responses that build on previous exchanges. "
            "Be thoughtful about context and continuity throughout the discussion."
        ),
    }
    return templates.get(mode, templates["medium"])


def _truncate_messages_to_budget(
    messages: list[dict[str, Any]],
    max_tokens: int,
    cfg: Any,
    model: str,
) -> list[dict[str, Any]]:
    """Truncate message history to fit within token budget.

    Preserves system messages and recent history. Removes oldest user/assistant
    turns until the context fits within the budget.

    Args:
        messages: List of message dicts to truncate
        max_tokens: Maximum token budget
        cfg: Config instance for threshold lookup
        model: Model name for logging

    Returns:
        Truncated message list that fits within budget
    """
    # Try to count tokens, fall back gracefully if tiktoken not available
    try:
        token_count = count_message_tokens(messages)
    except ImportError:
        logger.warning(
            "tiktoken not installed, cannot enforce context budget. "
            "Install with: pip install tiktoken"
        )
        return messages

    # If we're within budget, no truncation needed
    if token_count <= max_tokens:
        # Check if we're approaching the warning threshold
        warning_threshold = get_context_warning_threshold(cfg)
        warning_tokens = int(max_tokens * warning_threshold)
        if token_count >= warning_tokens:
            pct = (token_count / max_tokens) * 100
            logger.warning(
                f"Context approaching limit: {token_count}/{max_tokens} tokens ({pct:.1f}%) "
                f"for model {model}"
            )
        return messages

    logger.warning(
        f"Context exceeds budget: {token_count}/{max_tokens} tokens "
        f"for model {model}. Truncating conversation history."
    )

    # Separate system messages from conversation
    system_messages = [m for m in messages if m.get("role") == "system"]
    conversation_messages = [m for m in messages if m.get("role") != "system"]

    # Always preserve system messages and the current user prompt (last message)
    if not conversation_messages:
        # Edge case: only system messages, shouldn't exceed budget but handle gracefully
        return messages

    # Keep the last message (current prompt) and work backwards
    preserved = conversation_messages[-1:]
    remaining = conversation_messages[:-1]

    # Try progressively smaller history until we fit
    while remaining:
        # Try keeping system + remaining + current prompt
        candidate = system_messages + remaining + preserved

        try:
            candidate_tokens = count_message_tokens(candidate)
        except ImportError:
            # If tiktoken fails mid-truncation, return what we have
            break

        if candidate_tokens <= max_tokens:
            logger.info(
                f"Truncated to {len(candidate)} messages "
                f"({candidate_tokens} tokens, removed {len(messages) - len(candidate)} messages)"
            )
            return candidate

        # Remove the oldest conversation message
        remaining = remaining[1:]

    # If we still don't fit, try just system + current prompt
    minimal = system_messages + preserved
    try:
        minimal_tokens = count_message_tokens(minimal)
        logger.warning(
            f"Truncated to minimal context: {len(minimal)} messages "
            f"({minimal_tokens} tokens). All history removed."
        )
    except ImportError:
        pass

    return minimal


def _extract_document_text(raw_bytes: bytes, media_type: str, filename: str) -> str:
    """Extract readable text from a document's raw bytes.

    Supports: PDF, DOCX, PPTX, ePUB, HTML, Markdown, plain text.
    Falls back to UTF-8 decode for unrecognised types.
    """
    import io

    fname = filename.lower()

    # PDF
    if media_type == "application/pdf" or fname.endswith(".pdf"):
        import pdfplumber

        with pdfplumber.open(io.BytesIO(raw_bytes)) as pdf:
            text = "\n".join(page.extract_text() or "" for page in pdf.pages).strip()
        return text or f"[PDF: {filename} — no extractable text found]"

    # DOCX
    if media_type == (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ) or fname.endswith(".docx"):
        from docx import Document

        doc = Document(io.BytesIO(raw_bytes))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text:
                    parts.append(row_text)
        return "\n\n".join(parts) or f"[DOCX: {filename} — no extractable text found]"

    # PPTX
    if media_type == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ) or fname.endswith(".pptx"):
        from pptx import Presentation

        prs = Presentation(io.BytesIO(raw_bytes))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if slide_texts:
                parts.append(f"[Slide {i}]\n" + "\n\n".join(slide_texts))
        return "\n\n".join(parts) or f"[PPTX: {filename} — no extractable text found]"

    # ePUB
    if media_type == "application/epub+zip" or fname.endswith(".epub"):
        import ebooklib
        from bs4 import BeautifulSoup
        from ebooklib import epub

        book = epub.read_epub(io.BytesIO(raw_bytes))
        parts = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            soup = BeautifulSoup(item.get_content(), "html.parser")
            page_text = soup.get_text(separator="\n").strip()
            if page_text:
                parts.append(page_text)
        return "\n\n".join(parts) or f"[ePUB: {filename} — no extractable text found]"

    # HTML
    if media_type == "text/html" or fname.endswith((".html", ".htm")):
        from bs4 import BeautifulSoup

        html_content = raw_bytes.decode("utf-8", errors="replace")
        soup = BeautifulSoup(html_content, "html.parser")
        for tag in soup(["script", "style", "nav", "header", "footer"]):
            tag.decompose()
        return soup.get_text(separator="\n").strip()

    # Plain text / Markdown / fallback
    return raw_bytes.decode("utf-8", errors="replace")


def _build_user_message(
    prompt: str,
    attachments: list[dict[str, Any]] | None,
) -> dict[str, Any]:
    """Build a user message dict, including multimodal content blocks when attachments are present.

    For image attachments, the base64 data is placed in the Ollama ``images`` field.
    For document attachments (PDF, plain-text), the decoded text is prepended to the prompt.

    Args:
        prompt: The user's text prompt.
        attachments: Optional list of attachment dicts with keys ``type``, ``media_type``,
            ``data`` (base64), and optionally ``filename``.

    Returns:
        A message dict suitable for the Ollama /api/chat payload.
    """
    if not attachments:
        return {"role": "user", "content": prompt}

    import base64

    image_data: list[str] = []
    doc_text_parts: list[str] = []

    for att in attachments:
        att_type = att.get("type", "")
        raw_data = att.get("data", "")
        filename = att.get("filename") or "attachment"

        if att_type == "image":
            # Ollama accepts base64 image strings in the `images` array
            image_data.append(raw_data)
        elif att_type == "document":
            # Extract document text and prepend to prompt
            try:
                raw_bytes = base64.b64decode(raw_data)
                media_type = att.get("media_type", "")
                text = _extract_document_text(raw_bytes, media_type, filename)
                doc_text_parts.append(f"[Attached document: {filename}]\n{text}")
            except Exception:
                logger.warning("Failed to decode document attachment: %s", filename)

    full_prompt = prompt
    if doc_text_parts:
        full_prompt = "\n\n".join(doc_text_parts) + "\n\n" + prompt

    msg: dict[str, Any] = {"role": "user", "content": full_prompt}
    if image_data:
        msg["images"] = image_data
    return msg


def _prepare_chat_context(
    prompt: str,
    *,
    session: str = "default",
    new: bool = False,
    model: str | None = None,
    system: str | None = None,
    config_path: str | None = None,
    sessions_dir: str | None = None,
    rag_enabled: bool | None = None,
    rag_filters: dict[str, Any] | None = None,
    attachments: list[dict[str, Any]] | None = None,
    output_format: dict[str, Any] | None = None,
) -> ChatContext:
    """Prepare messages and context for a chat interaction.

    This function consolidates all the setup logic shared between
    chat() and chat_stream(), eliminating code duplication and
    making the codebase easier to maintain.

    Args:
        prompt: User's input message
        session: Session name to use
        new: Whether to start a new session (clearing history)
        model: Model to use (overrides config default)
        system: System prompt (overrides config default)
        config_path: Path to config file
        sessions_dir: Path to sessions directory
        rag_enabled: Enable/disable RAG for this request
        rag_filters: Metadata filters for RAG document selection
        attachments: Optional inline file attachments (images/documents)

    Returns:
        ChatContext with all prepared data needed for the LLM call
    """
    cfg = _cfg(config_path)

    base_url = _get_str(cfg, "ollama", "base_url", "http://127.0.0.1:11434")
    default_model = _get_str(cfg, "nyxgpt", "default_model", "llama3.1:8b")
    chosen_model = model or default_model
    chat_timeout_s = _get_int(cfg, "nyxgpt", "chat_timeout_seconds", 180)

    # Load session messages
    state = load_session(session, cfg, sessions_dir_override=sessions_dir)
    if new:
        state.messages = []

    messages: list[dict[str, Any]] = []

    # System prompt with adaptive mode
    sys_msg = system or _get_str(cfg, "nyxgpt", "system_prompt", "")

    # Apply adaptive prompt mode if enabled and no custom system prompt
    if not sys_msg.strip() and get_prompt_mode_enabled(cfg):
        # Count conversation messages to determine mode
        conversation_msg_count = len(state.messages)
        mode = _detect_prompt_mode(conversation_msg_count, cfg)
        sys_msg = _get_prompt_template(mode)
        logger.debug(
            f"Adaptive prompt mode: {mode} " f"(conversation has {conversation_msg_count} messages)"
        )

    if sys_msg.strip():
        # Apply minimization if enabled
        if get_system_prompt_minimize(cfg):
            sys_msg = _minimize_system_prompt(sys_msg)
        messages.append({"role": "system", "content": sys_msg.strip()})

    # Optional RAG context injection
    # Priority: 1) explicit rag_enabled param, 2) session metadata, 3) global config
    if rag_enabled is not None:
        should_use_rag = rag_enabled
    else:
        session_rag = state.meta.get("rag_enabled")
        should_use_rag = session_rag if isinstance(session_rag, bool) else get_rag_enabled(cfg)

    rag_context = ""
    rag_chunks = 0
    rag_rows = None  # Store raw RAG results

    if should_use_rag:
        # Build metadata filter from rag_filters dict if provided
        metadata_filter = None
        if rag_filters:
            from datetime import datetime

            # Parse dates if provided
            date_from_dt = None
            date_to_dt = None
            if rag_filters.get("date_from"):
                try:
                    date_from_dt = datetime.fromisoformat(rag_filters["date_from"])
                except (ValueError, TypeError):
                    logger.warning(f"Invalid date_from format: {rag_filters['date_from']}")
            if rag_filters.get("date_to"):
                try:
                    date_to_dt = datetime.fromisoformat(rag_filters["date_to"])
                except (ValueError, TypeError):
                    logger.warning(f"Invalid date_to format: {rag_filters['date_to']}")

            metadata_filter = MetadataFilter(
                doc_ids=rag_filters.get("doc_ids"),
                filename=rag_filters.get("filename"),
                tags=rag_filters.get("tags"),
                date_from=date_from_dt,
                date_to=date_to_dt,
            )

        # Disable debug mode for chat path - we don't need debug info here
        rows_result = retrieve_context(prompt, debug_mode=False, metadata_filter=metadata_filter)
        # Type narrowing: debug_mode=False means result is list[dict], not tuple
        rows = cast(list[dict], rows_result)

        # Force-include attached documents from session metadata.
        # These are always retrieved regardless of rag_filters, merged with normal results.
        attached_doc_ids = state.meta.get("attached_doc_ids", [])
        if attached_doc_ids and isinstance(attached_doc_ids, list):
            force_filter = MetadataFilter(doc_ids=list(attached_doc_ids))
            force_result = retrieve_context(prompt, debug_mode=False, metadata_filter=force_filter)
            force_rows = cast(list[dict], force_result)

            # Merge: deduplicate by (doc_id, chunk_id), force-included rows take precedence
            seen_keys: set[tuple] = set()
            merged: list[dict] = []
            for r in force_rows:
                key = (r.get("doc_id"), r.get("chunk_id"))
                if key not in seen_keys:
                    seen_keys.add(key)
                    merged.append(r)
            for r in rows:
                key = (r.get("doc_id"), r.get("chunk_id"))
                if key not in seen_keys:
                    seen_keys.add(key)
                    merged.append(r)
            rows = merged

        rag_chunks = len(rows)
        rag_rows = rows  # Save raw results
        rag_context = compose_context(rows)

        if rag_context:
            # Load configurable templates
            instruction_template = get_rag_instruction_template(cfg)
            context_format = get_rag_context_format(cfg)

            # Format the context using the template
            formatted_context = context_format.format(context=rag_context)

            # Build the full instruction with formatted context
            full_instruction = instruction_template.format(context=formatted_context)

            messages.append(
                {
                    "role": "system",
                    "content": full_instruction,
                }
            )

    # Add prior conversation
    for m in state.messages:
        role = str(m.get("role", "user"))
        content = str(m.get("content", ""))
        if content:
            messages.append({"role": role, "content": content})

    # Add this turn (with optional multimodal attachments)
    messages.append(_build_user_message(prompt, attachments))

    # Enforce context window budget
    max_tokens = get_context_window_size(cfg, chosen_model)
    messages = _truncate_messages_to_budget(messages, max_tokens, cfg, chosen_model)

    # Update session metadata with chosen model
    state.meta["model"] = chosen_model

    return ChatContext(
        messages=messages,
        state=state,
        chosen_model=chosen_model,
        base_url=base_url,
        chat_timeout_s=chat_timeout_s,
        rag_used=should_use_rag,
        rag_chunks=rag_chunks,
        rag_context=rag_rows,
        output_format=output_format,
    )


def _persist_chat_turn(
    context: ChatContext,
    prompt: str,
    reply: str,
    cfg: ConfigParser,
    sessions_dir: str | None = None,
) -> None:
    """Persist a completed chat turn to session storage."""
    timestamp = datetime.now(UTC).isoformat()
    context.state.messages.append(
        {
            "role": "user",
            "content": prompt,
            "id": str(uuid.uuid4()),
            "timestamp": timestamp,
        }
    )

    # Build assistant message with optional RAG chunks
    assistant_msg: dict[str, Any] = {
        "role": "assistant",
        "content": reply,
        "id": str(uuid.uuid4()),
        "timestamp": timestamp,
    }

    # Include RAG chunks if RAG was used
    if context.rag_used:
        assistant_msg["rag_chunks"] = [
            {
                "text": chunk.get("text", ""),
                "score": chunk.get("score", 0.0),
                "doc_id": chunk.get("doc_id"),
                "chunk_id": chunk.get("chunk_id"),
                "similarity_score": chunk.get("similarity_score"),
            }
            for chunk in (context.rag_context or [])
        ]

    context.state.messages.append(assistant_msg)
    save_session(context.state, cfg, sessions_dir_override=sessions_dir)


def chat(
    prompt: str,
    *,
    session: str = "default",
    new: bool = False,
    model: str | None = None,
    system: str | None = None,
    config_path: str | None = None,
    sessions_dir: str | None = None,
    rag_enabled: bool | None = None,
    rag_filters: dict[str, Any] | None = None,
    attachments: list[dict[str, Any]] | None = None,
    output_format: dict[str, Any] | None = None,
) -> ChatResult:
    """Run a chat turn with optional response caching, persisting session history.

    Response caching is based on the full conversation context (messages + model).
    Cache is keyed by hash of messages to ensure identical prompts get cached responses.

    Config:
      - `[cache] response_cache_enabled` (enable/disable caching)
      - `[cache] response_cache_backend` (memory or disk)
      - `[cache] response_cache_ttl_seconds` (cache expiration time)

    Args:
        prompt: User's input message
        session: Session name to use
        new: Whether to start a new session (clearing history)
        model: Model to use (overrides config default)
        system: System prompt (overrides config default)
        config_path: Path to config file
        sessions_dir: Path to sessions directory
        rag_enabled: Enable/disable RAG for this request
        rag_filters: Metadata filters for RAG document selection
        attachments: Optional inline file attachments (images/documents)

    Returns:
        ChatResult with reply and metadata
    """

    request_start = time.monotonic()
    context = _prepare_chat_context(
        prompt=prompt,
        session=session,
        new=new,
        model=model,
        system=system,
        config_path=config_path,
        sessions_dir=sessions_dir,
        rag_enabled=rag_enabled,
        rag_filters=rag_filters,
        attachments=attachments,
        output_format=output_format,
    )
    logger.info(
        "Chat request started",
        extra={"session": session, "model": context.chosen_model, "streaming": False},
    )

    # Try to retrieve from cache first
    cache = _get_response_cache()
    cache_key_data = {
        "messages": context.messages,
        "model": context.chosen_model,
    }
    cache_key = hash_text(json.dumps(cache_key_data, sort_keys=True))

    cached_reply = cache.get(cache_key)
    if cached_reply is not None:
        logger.debug(f"Response cache hit for session={session}")
        # Still persist the cached reply to session
        cfg = _cfg(config_path)
        _persist_chat_turn(context, prompt, cached_reply, cfg, sessions_dir)

        logger.info(
            "Chat request completed",
            extra={
                "session": session,
                "model": context.chosen_model,
                "streaming": False,
                "outcome": "cache_hit",
                "duration_ms": round((time.monotonic() - request_start) * 1000, 1),
                "output_tokens": _safe_token_count(cached_reply),
            },
        )
        return ChatResult(
            session=context.state.name,
            model=context.chosen_model,
            reply=cached_reply,
            rag_used=context.rag_used,
            rag_chunks=context.rag_chunks,
            rag_context=context.rag_context,
        )

    # Cache miss - call LLM
    logger.debug(f"Response cache miss for session={session}, calling LLM...")
    ollama_start = time.monotonic()
    try:
        reply = ollama_chat(
            base_url=context.base_url,
            model=context.chosen_model,
            messages=context.messages,
            timeout_s=context.chat_timeout_s,
            output_format=context.output_format,
        )
    except Exception as e:
        logger.error(
            "Chat request failed",
            extra={
                "session": session,
                "model": context.chosen_model,
                "streaming": False,
                "outcome": "error",
                "duration_ms": round((time.monotonic() - request_start) * 1000, 1),
                "ollama_duration_ms": round((time.monotonic() - ollama_start) * 1000, 1),
                "error": str(e),
                "error_type": type(e).__name__,
            },
        )
        raise
    ollama_duration_ms = round((time.monotonic() - ollama_start) * 1000, 1)

    # Store in cache
    cache.set(cache_key, reply)

    cfg = _cfg(config_path)
    _persist_chat_turn(context, prompt, reply, cfg, sessions_dir)

    logger.info(
        "Chat request completed",
        extra={
            "session": session,
            "model": context.chosen_model,
            "streaming": False,
            "outcome": "success",
            "duration_ms": round((time.monotonic() - request_start) * 1000, 1),
            "ollama_duration_ms": ollama_duration_ms,
            "input_tokens": _safe_message_token_count(context.messages),
            "output_tokens": _safe_token_count(reply),
        },
    )
    return ChatResult(
        session=context.state.name,
        model=context.chosen_model,
        reply=reply,
        rag_used=context.rag_used,
        rag_chunks=context.rag_chunks,
        rag_context=context.rag_context,
    )


def chat_stream(
    prompt: str,
    *,
    session: str = "default",
    new: bool = False,
    model: str | None = None,
    system: str | None = None,
    config_path: str | None = None,
    sessions_dir: str | None = None,
    rag_enabled: bool | None = None,
    rag_filters: dict[str, Any] | None = None,
    attachments: list[dict[str, Any]] | None = None,
    on_retry: Callable[[int, float, Exception], None] | None = None,
    output_format: dict[str, Any] | None = None,
) -> Iterator[str]:
    """Yield assistant text chunks for a chat turn while persisting the final reply.

    This function yields incremental text tokens. Callers should print or forward
    chunks as they arrive. Session persistence occurs after streaming completes.

    Args:
        prompt: User prompt text
        session: Session name
        new: Whether to create a new session
        model: Override model name
        system: Override system prompt
        config_path: Path to config file
        sessions_dir: Override sessions directory
        rag_enabled: Enable/disable RAG for this request
        rag_filters: Metadata filters for RAG document selection
        attachments: Optional inline file attachments (images/documents)
        on_retry: Optional callback(attempt, delay, error) for connection retries
    """

    request_start = time.monotonic()
    context = _prepare_chat_context(
        prompt=prompt,
        session=session,
        new=new,
        model=model,
        system=system,
        config_path=config_path,
        sessions_dir=sessions_dir,
        rag_enabled=rag_enabled,
        rag_filters=rag_filters,
        attachments=attachments,
        output_format=output_format,
    )

    logger.info(
        "Chat request started",
        extra={"session": session, "model": context.chosen_model, "streaming": True},
    )

    # Yield RAG metadata as first chunk if RAG was used
    if context.rag_used and context.rag_context:
        import json

        rag_data = {
            "type": "rag_metadata",
            "chunks": [
                {
                    "text": chunk.get("text", ""),
                    "score": chunk.get("score", 0.0),
                    "doc_id": chunk.get("doc_id"),
                    "chunk_id": chunk.get("chunk_id"),
                    "similarity_score": chunk.get("similarity_score"),
                }
                for chunk in context.rag_context
            ],
        }
        yield f"__RAG_START__{json.dumps(rag_data)}__RAG_END__\n"

    # Create retry callback that yields status messages
    def _retry_callback(attempt: int, delay: float, error: Exception) -> None:
        # Yield a special marker for retry status
        import json

        retry_data = {
            "type": "retry_status",
            "attempt": attempt,
            "delay": delay,
            "error": str(error),
        }
        # Store in a list that we'll check and yield
        retry_messages.append(f"__RETRY_START__{json.dumps(retry_data)}__RETRY_END__\n")

        # Call the original callback if provided
        if on_retry:
            on_retry(attempt, delay, error)

    # Store retry messages that occur during connection
    retry_messages: list[str] = []

    # Stream tokens and assemble final reply
    parts: list[str] = []
    ollama_start = time.monotonic()
    try:
        for chunk in ollama_chat_stream_tokens(
            base_url=context.base_url,
            model=context.chosen_model,
            messages=context.messages,
            timeout_s=context.chat_timeout_s,
            on_retry=_retry_callback,
            output_format=context.output_format,
        ):
            # Yield any queued retry messages first
            for retry_msg in retry_messages:
                yield retry_msg
            retry_messages.clear()

            parts.append(chunk)
            yield chunk
    except Exception as e:
        # Log the upstream Ollama/model-runtime failure here too, so callers
        # that don't go through the API's streaming endpoint (CLI, MCP tool,
        # tests) still get actionable detail in the caller's own log file
        # (api.log/cli.log) instead of a bare re-raise.
        logger.error(
            "Chat request failed",
            extra={
                "session": session,
                "model": context.chosen_model,
                "streaming": True,
                "outcome": "error",
                "duration_ms": round((time.monotonic() - request_start) * 1000, 1),
                "ollama_duration_ms": round((time.monotonic() - ollama_start) * 1000, 1),
                "error": str(e),
                "error_type": type(e).__name__,
            },
            exc_info=True,
        )
        # If we have retry messages but the connection ultimately failed,
        # yield them before re-raising
        for retry_msg in retry_messages:
            yield retry_msg
        raise
    ollama_duration_ms = round((time.monotonic() - ollama_start) * 1000, 1)

    reply = "".join(parts)

    cfg = _cfg(config_path)
    _persist_chat_turn(context, prompt, reply, cfg, sessions_dir)

    logger.info(
        "Chat request completed",
        extra={
            "session": session,
            "model": context.chosen_model,
            "streaming": True,
            "outcome": "success",
            "duration_ms": round((time.monotonic() - request_start) * 1000, 1),
            "ollama_duration_ms": ollama_duration_ms,
            "input_tokens": _safe_message_token_count(context.messages),
            "output_tokens": _safe_token_count(reply),
        },
    )
