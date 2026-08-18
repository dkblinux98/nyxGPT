"""FastAPI application for the nyxGPT API server.

Wires together the versioned `/api/v1` router (chat, sessions, RAG,
models, admin/SRE dashboard endpoints), unversioned `/health` and
`/metrics` endpoints, and the middleware stack (CORS, security headers,
request-scoped config loading, request ID propagation, rate limiting,
API key auth, and Prometheus instrumentation). Startup/shutdown behavior
(logging, tracing, error tracking, rate limiter, batch processor,
resource monitor, self-heal watchdog) is handled by the `lifespan`
context manager below.
"""

from __future__ import annotations

import argparse
import inspect
import io
import logging
import os
import secrets
import sys
import tempfile
import threading
import time
import urllib.error
import urllib.request
import uuid
from configparser import ConfigParser
from contextlib import asynccontextmanager, suppress
from dataclasses import dataclass
from pathlib import Path
from types import SimpleNamespace
from typing import Any, cast

from fastapi import (
    APIRouter,
    Body,
    Depends,
    FastAPI,
    File,
    HTTPException,
    Query,
    Request,
    UploadFile,
)
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, Response, StreamingResponse
from starlette.concurrency import run_in_threadpool
from starlette.status import HTTP_401_UNAUTHORIZED

import nyxgpt.config
from nyxgpt import admin_activity as admin_activity_module
from nyxgpt import api_models, aws_credentials_setup, config_wizard, models, secrets_setup, sessions
from nyxgpt import canary as canary_module
from nyxgpt import chat as chat_module
from nyxgpt import cloud_artifact_smoke as cloud_artifact_smoke_module
from nyxgpt import cloud_deploy as cloud_deploy_module
from nyxgpt import cloud_infra as cloud_infra_module
from nyxgpt import cloud_state as cloud_state_module
from nyxgpt import error_tracking as error_tracking_module
from nyxgpt import health as health_module
from nyxgpt import metrics as prom_metrics
from nyxgpt import ops as ops_module
from nyxgpt import portability as portability_module
from nyxgpt import release_candidate as release_candidate_module
from nyxgpt import resource_metrics_store as resource_metrics_store_module
from nyxgpt import restart_state as restart_state_module
from nyxgpt import self_heal as self_heal_module
from nyxgpt import support as support_module
from nyxgpt import tracing as tracing_module
from nyxgpt import usage_analytics as usage_analytics_module
from nyxgpt.api_models import (
    AttachDocumentRequest,
    ChatRequest,
    ChatResponse,
    CollectionClearResponse,
    CollectionDeleteResponse,
    CollectionInfo,
    CollectionSettings,
    CollectionSettingsResponse,
    CollectionsListResponse,
    CreateCollectionRequest,
    CreateCollectionResponse,
    InfoResponse,
    QueryCacheClearResponse,
    QueryCacheStatsResponse,
    RagChunkInfo,
    RagDocumentInfo,
    RagIndexRepoRequest,
    RagIndexRepoResponse,
    RagIngestRequest,
    RagIngestResponse,
    RagMetricsQueryRequest,
    RagMetricsQueryResponse,
    RagQueryRequest,
    RagQueryResponse,
    RagQueryResult,
    ReindexCollectionRequest,
    ReindexCollectionResponse,
    RenameRequest,
    ResourceMetricsResponse,
    SessionDocumentsResponse,
    SessionsListResponse,
    TagsRequest,
    TitleRequest,
)
from nyxgpt.batch_processor import BatchProcessor, RequestPriority
from nyxgpt.chat import chat as run_chat
from nyxgpt.chat import chat_stream
from nyxgpt.cloud import CloudCommandError
from nyxgpt.config import (
    get_auth_api_key,
    get_canary_error_rate_threshold,
    get_canary_latency_p95_threshold_ms,
    get_canary_min_requests,
    get_canary_namespace,
    get_canary_step_percent,
    get_canary_total_replicas,
    get_chat_timeout_seconds,
    get_default_model,
    get_error_tracking_config,
    get_log_aggregation_config,
    get_monitoring_config,
    get_ollama_base_url,
    get_rag_enabled,
    get_rag_good_score_threshold,
    get_rag_medium_score_threshold,
    get_rag_min_score,
    get_rate_limit_config,
    get_rate_limit_enabled,
    get_secrets_provider,
    get_self_heal_backoff_seconds,
    get_self_heal_check_interval_seconds,
    get_self_heal_default_enabled,
    get_self_heal_max_consecutive_restarts,
    get_session_backend,
    get_sessions_dir,
    get_tracing_config,
    load_config,
    log_effective_config,
    validate_bind_security,
)
from nyxgpt.logging import configure_logging, request_id_var
from nyxgpt.ollama_client import ModelRuntimeError, get_json, post_json
from nyxgpt.rag.rag import (
    annotate_chunk_numbering,
    clear_query_cache,
    get_query_cache_stats,
    ingest_document,
    retrieve_context,
)
from nyxgpt.rate_limiter import RateLimiter
from nyxgpt.resource_monitor import ResourceMonitor, get_resource_monitor, init_resource_monitor
from nyxgpt.token_counter import count_tokens as _count_usage_tokens
from nyxgpt.tracing import current_trace_id
from nyxgpt.version import running_version

log = logging.getLogger("nyxgpt.api")

# Global rate limiter instance (initialized at startup if enabled)
_rate_limiter: RateLimiter | None = None

# Global batch processor instance (initialized at startup if enabled)
_batch_processor: BatchProcessor[dict, dict] | None = None

# Global resource monitor instance (initialized at startup)
_resource_monitor: ResourceMonitor | None = None


@dataclass
class ClientCapabilities:
    """Client capability hints for content negotiation.

    Attributes:
        supports_sse: Client supports Server-Sent Events (text/event-stream)
        supports_structured_events: Client supports structured event types (metadata, text, done, error)
        supports_streaming: Client supports streaming responses
        client_version: Optional client version string (e.g., "web-ui/1.0.0")
        max_event_size: Maximum event payload size in bytes (0 = unlimited)
    """

    supports_sse: bool = False
    supports_structured_events: bool = False
    supports_streaming: bool = True
    client_version: str | None = None
    max_event_size: int = 0


def _parse_client_capabilities(request: Request) -> ClientCapabilities:
    """Parse client capability hints from request headers.

    Examines standard and custom headers to determine client capabilities:
    - Accept: text/event-stream indicates SSE support
    - X-Client-Supports-SSE: Explicit SSE capability flag
    - X-Client-Supports-Structured-Events: Structured event support flag
    - X-Client-Version: Client version identifier
    - X-Client-Max-Event-Size: Maximum event size in bytes

    Args:
        request: FastAPI Request object

    Returns:
        ClientCapabilities with detected or default capabilities
    """
    capabilities = ClientCapabilities()

    # Check Accept header for text/event-stream
    accept = request.headers.get("accept", "")
    capabilities.supports_sse = "text/event-stream" in accept.lower()

    # Check explicit capability headers
    supports_sse_header = request.headers.get("x-client-supports-sse", "").lower()
    if supports_sse_header in ("true", "1", "yes"):
        capabilities.supports_sse = True

    supports_structured = request.headers.get("x-client-supports-structured-events", "").lower()
    if supports_structured in ("true", "1", "yes"):
        capabilities.supports_structured_events = True

    # Check for streaming support (default true)
    supports_streaming = request.headers.get("x-client-supports-streaming", "true").lower()
    capabilities.supports_streaming = supports_streaming not in ("false", "0", "no")

    # Get client version
    capabilities.client_version = request.headers.get("x-client-version") or None

    # Get max event size
    max_size_str = request.headers.get("x-client-max-event-size", "0")
    try:
        capabilities.max_event_size = int(max_size_str)
    except ValueError:
        capabilities.max_event_size = 0

    return capabilities


def log_with_context(level, message, request_id=None, **extra):
    """Helper for structured logging with consistent context."""
    extra_fields = {"request_id": request_id, **extra} if request_id else extra
    log.log(level, message, extra=extra_fields)


# ----------------------------
# Startup diagnostics using lifespan
# ----------------------------


@asynccontextmanager
async def lifespan(_app: FastAPI):
    """Run startup and shutdown initialization for the FastAPI app.

    On startup: refuses to start if `[api] host` is non-loopback and
    `[auth] enabled` isn't true (P6-1 hardening gate; skipped inside a
    container, see below), configures centralized logging, initializes
    tracing and error tracking (both no-ops unless enabled in config.ini),
    ensures the sessions directory exists, does a warn-only Ollama
    reachability check, initializes the rate limiter and batch processor if
    enabled, starts the resource monitor, and starts the self-heal watchdog
    (always running so the dashboard toggle takes effect without a restart).
    Beyond the bind-security refusal, initialization failures are logged but
    never prevent the API from starting.

    On shutdown: stops the batch processor and self-heal watchdog.
    """
    global _rate_limiter, _batch_processor

    cfg = load_config(None)

    # P6-1 hardening gate: a non-loopback bind with auth disabled would let
    # anyone who can reach this host/network call the API with no
    # credentials -- refuse to start rather than merely warn. Skipped inside
    # a container (Compose `api` service, Kubernetes pod): both hardcode
    # uvicorn's own `--host 0.0.0.0` for the container's *network namespace*
    # regardless of `[api] host` (see docker/entrypoint.sh) -- real
    # host/cluster exposure there is gated by Docker's port-publish
    # (`NYXGPT_BIND_ADDR`) or the Kubernetes Service type, neither of which
    # is visible to this process (docs/security.md#network-security).
    if not os.environ.get("NYXGPT_CONTAINER_RUNTIME"):
        bind_error = validate_bind_security(cfg)
        if bind_error:
            print(f"ERROR: {bind_error}", file=sys.stderr)
            raise RuntimeError(bind_error)

    # Initialize centralized logging once for the API process
    try:
        configure_logging(cfg, console=False, filename="api.log")
        log.info("Centralized logging initialized", extra={"component": "startup"})
    except Exception as e:
        # Logging should not prevent API startup
        print(f"Logging initialization failed: {e}")

    try:
        log_effective_config(cfg)
    except Exception as e:
        log.warning("Failed to log effective config: %s", e, extra={"component": "startup"})

    # Initialize distributed tracing (no-op unless [tracing] enabled = true)
    try:
        tracing_module.init_tracing(get_tracing_config(cfg))
    except Exception as e:
        log.warning("Tracing initialization failed: %s", e, extra={"component": "startup"})

    # Initialize error tracking (no-op unless [error_tracking] enabled = true
    # and a local DSN is configured)
    try:
        error_tracking_module.init_error_tracking(get_error_tracking_config(cfg))
    except Exception as e:
        log.warning("Error tracking initialization failed: %s", e, extra={"component": "startup"})

    # Ensure sessions directory exists
    sessions_dir = get_sessions_dir(cfg)
    try:
        sessions_dir.mkdir(parents=True, exist_ok=True)
        log.info(
            "Sessions directory ready",
            extra={"component": "startup", "sessions_dir": str(sessions_dir)},
        )
    except Exception as e:
        log.error("Failed to prepare sessions directory %s: %s", sessions_dir, e)

    # One-time (idempotent) import of legacy JSON session files into the DB
    # when the Cassandra session backend is active (#3590). Sessions already
    # present in the DB are never overwritten; the legacy files stay on disk
    # as a read-only archive (see docs/session-storage.md).
    try:
        if get_session_backend(cfg) == "cassandra":
            from nyxgpt import session_db

            report = session_db.migrate_sessions_dir(sessions_dir)
            log.info(
                "Session backend: cassandra (migrated %d legacy file session(s), "
                "%d already in DB, %d invalid, %d errors)",
                len(report["migrated"]),
                len(report["skipped_existing"]),
                len(report["skipped_invalid"]),
                len(report["errors"]),
                extra={"component": "startup"},
            )
    except Exception as e:
        # Never prevent API startup; session endpoints will surface store
        # errors per-request if Cassandra stays unreachable.
        log.error("Legacy session migration failed: %s", e, extra={"component": "startup"})

    # Pre-touch known RAG metric label combinations so legitimate zero
    # states render as 0 on the SPOG panels instead of "No data"
    try:
        prom_metrics.initialize_known_rag_metric_series()
    except Exception as e:
        log.warning("Failed to initialize RAG metric series: %s", e, extra={"component": "startup"})

    # Warn-only Ollama reachability check
    base_url = get_ollama_base_url(cfg).rstrip("/")
    health_url = f"{base_url}/api/tags"
    try:
        req = urllib.request.Request(health_url, method="GET")
        with urllib.request.urlopen(req, timeout=2):
            log.info(
                "Ollama health check passed",
                extra={"component": "startup", "ollama_url": base_url},
            )
    except urllib.error.URLError as e:
        log.warning(
            "Ollama health check failed",
            extra={
                "component": "startup",
                "ollama_url": base_url,
                "error": str(e),
                "note": "API will still start; chat requests may fail until Ollama is available",
            },
        )

    # Initialize rate limiter if enabled
    if get_rate_limit_enabled(cfg):
        rate_cfg = get_rate_limit_config(cfg)
        _rate_limiter = RateLimiter(
            requests_per_second=rate_cfg["requests_per_second"],
            burst_size=rate_cfg["burst_size"],
        )
        log.info(
            "Rate limiting enabled",
            extra={
                "component": "startup",
                "requests_per_second": rate_cfg["requests_per_second"],
                "burst_size": rate_cfg["burst_size"],
            },
        )
    else:
        log.info("Rate limiting disabled", extra={"component": "startup"})

    # Initialize batch processor if enabled
    from nyxgpt.config import get_batch_enabled, get_batch_size, get_batch_wait_time_ms

    if get_batch_enabled(cfg):
        batch_size = get_batch_size(cfg)
        wait_time_ms = get_batch_wait_time_ms(cfg)

        # Define batch processing function
        def _process_chat_batch(requests):
            """Process a batch of chat requests together."""
            results = []
            for batch_req in requests:
                try:
                    # Extract request data
                    req_data = batch_req.data
                    result = run_chat(
                        req_data["prompt"],
                        session=req_data["session"],
                        new=req_data["new"],
                        model=req_data["model"],
                        system=req_data["system"],
                        config_path=req_data["config_path"],
                        sessions_dir=req_data["sessions_dir"],
                        rag_enabled=req_data.get("rag_enabled"),
                        rag_filters=req_data.get("rag_filters"),
                    )
                    # Convert ChatResult to dict
                    result_dict = {
                        "session": result.session,
                        "model": result.model,
                        "reply": result.reply,
                        "rag_used": result.rag_used,
                        "rag_chunks": result.rag_chunks,
                        "rag_context": result.rag_context,
                    }
                    results.append(result_dict)
                except Exception as e:
                    # On error, append error dict
                    results.append({"error": str(e), "error_type": type(e).__name__})

            return results

        _batch_processor = BatchProcessor(
            batch_size=batch_size,
            wait_time_ms=wait_time_ms,
            process_fn=_process_chat_batch,
        )
        _batch_processor.start()

        log.info(
            "Request batching enabled",
            extra={
                "component": "startup",
                "batch_size": batch_size,
                "wait_time_ms": wait_time_ms,
            },
        )
    else:
        log.info("Request batching disabled", extra={"component": "startup"})

    # Initialize resource monitor
    global _resource_monitor
    _resource_monitor = init_resource_monitor(batch_processor=_batch_processor)

    # Start the resource metrics history sampler. It's always running (like
    # the self-heal watchdog below) so server-side history accumulates from
    # process start, independent of whether the Settings page is ever open.
    resource_metrics_store_module.get_sampler().start()
    log.info("Resource metrics sampler initialized", extra={"component": "startup"})

    # Start the self-heal watchdog. It's always running so the SRE/admin
    # dashboard's toggle takes effect immediately without an API restart --
    # it only takes action when enabled (seeded from config.ini on a fresh
    # install, then controlled entirely by the dashboard toggle).
    self_heal_module.seed_enabled_default(get_self_heal_default_enabled(cfg))
    watchdog = self_heal_module.get_watchdog()
    watchdog.interval_seconds = get_self_heal_check_interval_seconds(cfg)
    watchdog.max_consecutive_restarts = get_self_heal_max_consecutive_restarts(cfg)
    watchdog.backoff_seconds = get_self_heal_backoff_seconds(cfg)
    watchdog.start()
    log.info("Self-heal watchdog initialized", extra={"component": "startup"})

    yield

    # Cleanup: stop batch processor
    if _batch_processor:
        _batch_processor.stop()
        log.info("Batch processor stopped", extra={"component": "shutdown"})

    self_heal_module.get_watchdog().stop()
    log.info("Self-heal watchdog stopped", extra={"component": "shutdown"})

    resource_metrics_store_module.get_sampler().stop()
    log.info("Resource metrics sampler stopped", extra={"component": "shutdown"})


# Versioned API router
app = FastAPI(title="nyxGPT", version="1.0.0.md", lifespan=lifespan)

# Must happen here, before uvicorn ever calls `app` -- including for the
# lifespan startup event itself. Starlette freezes `app.middleware_stack` on
# that first call, so instrumenting from inside `lifespan()`/`init_tracing()`
# (as this used to do) silently produced zero HTTP request spans in Jaeger
# even with tracing "enabled" -- see `tracing.instrument_fastapi_app`.
tracing_module.instrument_fastapi_app(app)

api = APIRouter(prefix="/api/v1")


# CORS: default to local-only origins (configurable via NYXGPT_CORS_ORIGINS)
# Example: export NYXGPT_CORS_ORIGINS="http://127.0.0.1:3000,http://localhost:3000"
_origins_env = os.environ.get("NYXGPT_CORS_ORIGINS", "").strip()
if _origins_env:
    allow_origins = [o.strip() for o in _origins_env.split(",") if o.strip()]
else:
    allow_origins = [
        "http://127.0.0.1:3000",
        "http://localhost:3000",
        "http://127.0.0.1:5173",
        "http://localhost:5173",
    ]

app.add_middleware(
    CORSMiddleware,
    allow_origins=allow_origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Security headers middleware
@app.middleware("http")
async def security_headers_middleware(request: Request, call_next):
    """Add security headers to all responses.

    Headers added:
    - Content-Security-Policy: Restrict resource loading
    - X-Content-Type-Options: Prevent MIME sniffing
    - X-Frame-Options: Prevent clickjacking
    - Strict-Transport-Security: Force HTTPS (only for HTTPS requests)
    """
    response = await call_next(request)

    # Content Security Policy: Restrict to same origin and specific trusted sources
    # This prevents XSS by limiting where resources can be loaded from
    response.headers["Content-Security-Policy"] = (
        "default-src 'self'; "
        "script-src 'self' 'unsafe-inline'; "  # Allow inline scripts for development
        "style-src 'self' 'unsafe-inline'; "  # Allow inline styles
        "img-src 'self' data:; "  # Allow data URIs for images
        "connect-src 'self'; "  # API calls to same origin only
        "frame-ancestors 'none'; "  # Prevent embedding (redundant with X-Frame-Options)
        "form-action 'self'; "  # Forms can only submit to same origin
        "base-uri 'self'"  # Prevent base tag injection
    )

    # Prevent MIME type sniffing
    response.headers["X-Content-Type-Options"] = "nosniff"

    # Prevent clickjacking by disallowing framing
    response.headers["X-Frame-Options"] = "DENY"

    # Strict-Transport-Security (HSTS) - only for HTTPS
    # Force browser to use HTTPS for 1 year
    if request.url.scheme == "https":
        response.headers["Strict-Transport-Security"] = "max-age=31536000; includeSubDomains"

    return response


MAX_BODY_BYTES = int(os.environ.get("NYXGPT_MAX_BODY_BYTES", "1048576"))  # 1 MiB default


# Middleware to load config and hot-apply logging on every request
@app.middleware("http")
async def load_cfg_and_refresh_logging(request: Request, call_next):
    """Load config for this request and hot-apply logging level.

    We want edits to ~/.nyxGPT/config.ini (model, rag enabled, log level, auth, etc.)
    to take effect without restarting the API process.

    The loaded config is stored on request.state.cfg for reuse by downstream
    middleware/handlers.
    """

    cfg = load_config(None)
    request.state.cfg = cfg

    # Hot-apply logging config (especially level) on every request.
    # configure_logging() is expected to be idempotent and cheap.
    # Never block request handling on logging reconfiguration.
    with suppress(Exception):
        configure_logging(cfg, console=False, filename="api.log")

    return await call_next(request)


@app.middleware("http")
async def add_request_id_and_limits(request: Request, call_next):
    """Reject oversized request bodies and stamp every request with an ID.

    Returns a `413 payload_too_large` JSON error when the `Content-Length`
    header exceeds `MAX_BODY_BYTES` (malformed/missing `Content-Length` is
    ignored rather than rejected). Otherwise reuses the client-supplied
    `X-Request-Id` header; failing that, derives one from the active OTel
    span's trace id (set by `FastAPIInstrumentor`, which wraps the whole
    ASGI app and so runs before this middleware -- see `tracing.init_tracing`)
    so the human-facing id and the Jaeger trace are the same value; failing
    that (tracing disabled/unreachable), falls back to a new UUID4. Stores
    the result on `request.state.request_id` and the `request_id_var`
    context variable for structured logging, and echoes it back as
    `X-Request-Id` on the response.
    """
    # Request size guard based on Content-Length when available
    cl = request.headers.get("content-length")
    if cl:
        try:
            if int(cl) > MAX_BODY_BYTES:
                return JSONResponse(
                    status_code=413,
                    content={
                        "error": {
                            "code": "payload_too_large",
                            "message": "Request body too large",
                        }
                    },
                )
        except Exception:
            # Ignore malformed content-length
            pass

    # Accept client-provided request ID; else derive from the active trace,
    # so the id printed in logs is the same one Jaeger has; else generate one.
    req_id = request.headers.get("x-request-id") or current_trace_id() or str(uuid.uuid4())
    request.state.request_id = req_id

    # Set request ID in context variable for automatic logging
    request_id_var.set(req_id)

    response = await call_next(request)
    response.headers["X-Request-Id"] = req_id
    return response


@app.middleware("http")
async def rate_limit_middleware(request: Request, call_next):
    """Apply rate limiting to API endpoints.

    Uses token bucket algorithm to limit requests per IP address.
    Returns 429 Too Many Requests when limit exceeded.
    """
    # Skip if rate limiting not enabled
    if _rate_limiter is None:
        return await call_next(request)

    # Get client ID (IP address)
    client_id = _rate_limiter.get_client_ip(request)

    # Check if allowed
    allowed, headers = _rate_limiter.is_allowed(client_id)

    if not allowed:
        # Log rate limit violation
        req_id = getattr(request.state, "request_id", None)
        log.warning(
            "Rate limit exceeded for %s on %s (request_id=%s)",
            client_id,
            request.url.path,
            req_id,
        )
        prom_metrics.RATE_LIMIT_REJECTIONS_TOTAL.labels(path=request.url.path).inc()

        # Return 429 error with rate limit headers
        return JSONResponse(
            status_code=429,
            content={
                "error": {
                    "code": "rate_limit_exceeded",
                    "message": "Too many requests. Please try again later.",
                    "request_id": req_id,
                }
            },
            headers=headers,
        )

    # Add rate limit headers to response
    response = await call_next(request)
    for key, value in headers.items():
        response.headers[key] = value

    return response


@app.middleware("http")
async def request_latency_middleware(request: Request, call_next):
    """Track request latency for resource monitoring.

    Measures total request processing time and records it in the resource monitor.
    """
    start_time = time.time()

    # Process request
    response = await call_next(request)

    # Calculate latency
    latency_ms = (time.time() - start_time) * 1000

    # Record in resource monitor
    monitor = get_resource_monitor()
    if monitor is not None:
        monitor.record_request_latency(latency_ms, is_error=response.status_code >= 500)

    return response


@app.middleware("http")
async def api_key_auth(request: Request, call_next):
    """Enforce API key authentication on versioned `/api/v1` routes.

    `/health`, `/docs`, `/openapi`, and `/redoc` are always exempt, and
    anything outside `/api/v1` is passed through unauthenticated. When
    `[auth] enabled = true` in config.ini, the configured header (default
    `X-API-Key`) must match the configured `api_key` using a constant-time
    comparison (to avoid timing attacks); a missing/incorrect key returns a
    `401 unauthorized` JSON error. When auth is disabled, requests pass
    through unchecked.
    """
    path = request.url.path

    # Allow unauthenticated access to health and docs
    if (
        path == "/health"
        or path.startswith("/docs")
        or path.startswith("/openapi")
        or path.startswith("/redoc")
    ):
        return await call_next(request)

    # Only protect versioned API
    if not path.startswith("/api/v1"):
        return await call_next(request)

    cfg = getattr(request.state, "cfg", None)
    cfg = cfg or load_config(None)
    req_id = getattr(request.state, "request_id", None)

    # Check `enabled` on its own first -- it's a plain config.ini read, no
    # AWS call. Only resolve the (possibly cloud-sourced) api_key when auth
    # is actually enabled, and do that resolution in a thread pool: this
    # middleware is `async def`, so Starlette runs it directly on the event
    # loop rather than dispatching it to a worker thread the way a plain
    # `def` route handler would -- a synchronous boto3 call here would
    # otherwise block the event loop for every concurrent request.
    enabled = cfg.getboolean("auth", "enabled", fallback=False)
    log.debug("auth check (request_id=%s) enabled=%s", req_id, enabled)
    if not enabled:
        return await call_next(request)

    auth = await run_in_threadpool(_auth_cfg, cfg)
    header = auth.get("header", "X-API-Key")
    expected = auth.get("api_key")
    provided = request.headers.get(header)

    # Use constant-time comparison to prevent timing attacks
    # This prevents attackers from determining the correct API key
    # by measuring response times
    auth_valid = False
    if expected and provided:
        auth_valid = secrets.compare_digest(expected, provided)

    if not auth_valid:
        req_id = getattr(request.state, "request_id", None)
        return JSONResponse(
            status_code=HTTP_401_UNAUTHORIZED,
            content={
                "error": {
                    "code": "unauthorized",
                    "message": "Invalid or missing API key",
                    "request_id": req_id,
                }
            },
        )

    return await call_next(request)


@app.middleware("http")
async def prometheus_metrics_middleware(request: Request, call_next):
    """Record Prometheus HTTP metrics for every request.

    Registered after api_key_auth so it becomes the outermost middleware,
    ensuring it observes the final response status even for requests
    rejected by auth, rate limiting, or the body-size guard.
    """
    start_time = time.time()
    response = await call_next(request)
    duration_s = time.time() - start_time

    route = request.scope.get("route")
    path_label = route.path if route is not None else request.url.path
    status_label = str(response.status_code)

    prom_metrics.HTTP_REQUESTS_TOTAL.labels(
        method=request.method, path=path_label, status=status_label
    ).inc()
    prom_metrics.HTTP_REQUEST_DURATION_SECONDS.labels(
        method=request.method, path=path_label
    ).observe(duration_s)
    if response.status_code >= 500:
        prom_metrics.HTTP_ERRORS_TOTAL.labels(method=request.method, path=path_label).inc()

    return response


@app.exception_handler(HTTPException)
async def http_exception_handler(request: Request, exc: HTTPException):
    """Render raised `HTTPException`s in the API's standard error envelope.

    Preserves the exception's status code and puts a string `detail` in
    `error.message`; a non-string `detail` (e.g. a dict/list) is instead
    placed in `error.details` with a generic `error.message`. Includes the
    request's `request_id` (from `add_request_id_and_limits`) for
    correlation with server logs.
    """
    req_id = getattr(request.state, "request_id", None)
    detail = exc.detail
    return JSONResponse(
        status_code=exc.status_code,
        content={
            "error": {
                "code": "http_error",
                "message": detail if isinstance(detail, str) else "Request failed",
                "details": None if isinstance(detail, str) else detail,
                "request_id": req_id,
            }
        },
    )


@app.exception_handler(Exception)
async def unhandled_exception_handler(request: Request, _exc: Exception):
    """Catch-all handler for exceptions not raised as an `HTTPException`.

    Logs the full traceback, forwards the exception to the error tracker
    (GlitchTip, when enabled), and always returns a generic
    `500 internal_error` JSON response so unexpected failures never leak
    internal exception details to the client.
    """
    req_id = getattr(request.state, "request_id", None)
    log.exception("Unhandled API error (request_id=%s)", req_id)
    error_tracking_module.capture_exception(
        _exc,
        request_id=req_id or "N/A",
        path=request.url.path,
        method=request.method,
    )
    return JSONResponse(
        status_code=500,
        content={
            "error": {
                "code": "internal_error",
                "message": "Internal server error",
                "request_id": req_id,
            }
        },
    )


# ----------------------------
# Helpers
# ----------------------------


# Request-scoped config helper
def _req_cfg(request: Request) -> ConfigParser:
    """Return the config loaded for this request, loading it if missing.

    Normally `request.state.cfg` is already populated by the
    `load_cfg_and_refresh_logging` middleware; this falls back to loading
    config.ini directly (and caching it on `request.state`) for code paths
    that run outside that middleware, e.g. exception handlers.
    """
    cfg = getattr(request.state, "cfg", None)
    if cfg is None:
        cfg = load_config(None)
        request.state.cfg = cfg
    return cfg


# Auth config is read on each request so ~/.nyxGPT/config.ini edits
# take effect without restarting the API.
def _auth_cfg(cfg: ConfigParser | None = None) -> dict[str, Any]:
    """Read the `[auth]` section into a plain dict.

    Args:
        cfg: Config to read from; loads config.ini fresh if omitted.

    Returns:
        Dict with `enabled` (bool), `api_key` (str, empty if unset), and
        `header` (str, defaults to `X-API-Key`).
    """
    cfg = cfg or load_config(None)
    enabled = cfg.getboolean("auth", "enabled", fallback=False)
    api_key = get_auth_api_key(cfg)
    header = cfg.get("auth", "header", fallback="X-API-Key").strip() or "X-API-Key"
    return {
        "enabled": enabled,
        "api_key": api_key,
        "header": header,
    }


# --- Config file helpers and hot-update endpoints ---


def _config_file_path() -> Path:
    """Return the canonical per-user config.ini path (`~/.nyxGPT/config.ini`)."""
    return Path.home() / ".nyxGPT" / "config.ini"


def _apply_hot_config_updates(updates: dict[str, Any]) -> dict[str, Any]:
    """Apply a small set of hot config updates to ~/.nyxGPT/config.ini.

    Supported updates:
    - default_model (str) -> [ollama] default_model
    - rag_enabled (bool)  -> [rag] enable_chat_context
    - log_level (str)     -> [logging] level

    After writing, reload config and re-configure logging so log level changes apply immediately.
    Model and RAG are read per-request by chat endpoints, so they apply immediately as well.
    """

    cfg_path = _config_file_path()
    cfg_path.parent.mkdir(parents=True, exist_ok=True)

    parser = ConfigParser()
    # Preserve key case (ConfigParser lowercases by default)
    parser.optionxform = str  # type: ignore[assignment]
    if cfg_path.exists():
        parser.read(cfg_path)

    def ensure_section(name: str) -> None:
        """Add section `name` to `parser` if it doesn't already exist."""
        if not parser.has_section(name):
            parser.add_section(name)

    out: dict[str, Any] = {}

    if "default_model" in updates and isinstance(updates.get("default_model"), str):
        ensure_section("nyxgpt")
        parser.set("nyxgpt", "default_model", updates["default_model"].strip())
        out["default_model"] = updates["default_model"].strip()

    if "rag_enabled" in updates:
        ensure_section("rag")
        val = bool(updates["rag_enabled"])
        parser.set("rag", "enable_chat_context", "true" if val else "false")
        out["rag_enabled"] = val

    if "log_level" in updates and isinstance(updates.get("log_level"), str):
        ensure_section("logging")
        lvl = updates["log_level"].strip().upper()
        parser.set("logging", "level", lvl)
        out["log_level"] = lvl

    # Persist changes
    with cfg_path.open("w", encoding="utf-8") as f:
        parser.write(f)
    os.chmod(cfg_path, 0o600)

    # Invalidate config cache to force reload on next access
    # This ensures mtime-based caching works even for rapid writes/reads
    nyxgpt.config._CACHED_CFG = None
    nyxgpt.config._CACHED_PATH = None
    nyxgpt.config._CACHED_MTIME_NS = None

    # Hot-apply logging changes immediately
    try:
        cfg = load_config(None)
        configure_logging(cfg, console=False, filename="api.log")
    except Exception:
        # Do not fail the request if logging reconfig fails
        pass

    return out


def _apply_auth_config_updates(updates: dict[str, Any]) -> dict[str, Any]:
    """Apply auth-section updates to ~/.nyxGPT/config.ini.

    Supported updates:
    - enabled (bool) -> [auth] enabled
    - header (str)   -> [auth] header
    - api_key (str)  -> [auth] api_key

    The **api** tier reads auth config fresh on every request (see
    `_auth_cfg`), so it honours a change here immediately. The **web** tier
    does not: its service wrapper reads `[auth]` once at process start and
    exports `NYXGPT_AUTH_API_KEY` into a Node process that cannot observe a
    later edit, so rotating the key from this endpoint 401s every proxied
    call until `web` restarts (#3806).

    That divergence is recorded, not left to be discovered: this is the third
    writer of a restart-required key (with the Configuration Wizard's
    `POST /config/sections` and `nyxgpt secrets setup`), and all three go
    through the same `config_wizard` classification and the same on-disk
    `restart_state`, so all three raise the same persistent notice and offer
    the same restart.
    """

    cfg_path = _config_file_path()
    cfg_path.parent.mkdir(parents=True, exist_ok=True)

    parser = ConfigParser()
    parser.optionxform = str  # type: ignore[assignment]
    if cfg_path.exists():
        parser.read(cfg_path)

    out: dict[str, Any] = {}

    if "enabled" in updates:
        out["enabled"] = bool(updates["enabled"])

    if "header" in updates and isinstance(updates.get("header"), str):
        out["header"] = updates["header"].strip()

    if "api_key" in updates and isinstance(updates.get("api_key"), str):
        out["api_key"] = updates["api_key"]

    # Computed against the *pre-write* parser: `restart_required_detail` needs
    # the value the still-running service loaded, which is what is on disk
    # right now. Keys the wizard schema does not declare are skipped rather
    # than assumed hot -- an undeclared key has no classification to consult.
    classified = config_wizard.activation_classification()
    tracked = {key: value for key, value in out.items() if f"auth.{key}" in classified}
    restart_detail = config_wizard.restart_required_detail({"auth": tracked}, parser)

    if not parser.has_section("auth"):
        parser.add_section("auth")

    if "enabled" in out:
        parser.set("auth", "enabled", "true" if out["enabled"] else "false")
    if "header" in out:
        parser.set("auth", "header", out["header"])
    if "api_key" in out:
        parser.set("auth", "api_key", out["api_key"])

    with cfg_path.open("w", encoding="utf-8") as f:
        parser.write(f)
    os.chmod(cfg_path, 0o600)

    # Mark first, reconcile second -- the same load-bearing order as
    # `config_sections_update`: putting a key back to the value the running
    # service still holds arrives here as a change, and only the reconcile
    # pass (comparing against the originally recorded running value) can
    # retire it without a restart.
    for component, changes in restart_detail.items():
        restart_state_module.mark_pending(component, changes)
    for component, saved in config_wizard.restart_activation_saved({"auth": tracked}).items():
        restart_state_module.reconcile_saved(component, saved)

    nyxgpt.config._CACHED_CFG = None
    nyxgpt.config._CACHED_PATH = None
    nyxgpt.config._CACHED_MTIME_NS = None

    return out


def _mask_api_key(api_key: str) -> str:
    """Mask an API key for display, keeping only a few edge characters."""
    if not api_key:
        return ""
    if len(api_key) <= 8:
        return "*" * len(api_key)
    return f"{api_key[:4]}{'*' * (len(api_key) - 8)}{api_key[-4:]}"


# Ollama model management helpers
def _ollama_url(cfg: ConfigParser, path: str) -> str:
    """Join the configured Ollama base URL with `path` into a full URL."""
    base = get_ollama_base_url(cfg).rstrip("/")
    if not path.startswith("/"):
        path = "/" + path
    return base + path


def _cfg(cfg_path: Path | None = None):
    """Load config.ini from `cfg_path`, or the default location if omitted."""
    return load_config(cfg_path)


def _chat_runtime_defaults(cfg: ConfigParser | None = None) -> dict[str, Any]:
    """Read chat-related defaults from config.ini.

    If a request-scoped config is available, reuse it so a single request does
    not re-read config.ini multiple times.

    Config edits (model, rag enabled, etc.) should take effect without restart.
    """

    cfg = cfg or load_config(None)
    rag_enabled = get_rag_enabled(cfg)
    return {
        "cfg": cfg,
        "default_model": get_default_model(cfg),
        "rag_enabled": rag_enabled,
        "chat_timeout_seconds": get_chat_timeout_seconds(cfg),
    }


def _maybe_kw(fn, name: str) -> bool:
    """Return True if function `fn` accepts keyword argument `name`."""

    try:
        return name in inspect.signature(fn).parameters
    except Exception:
        return False


def _sessions_dir_from_str(s: str | None) -> Path | None:
    """Convert an optional query-param string into an expanded `Path`.

    Returns None for a falsy input (so callers can fall back to the
    config-derived sessions directory) instead of a bogus `Path("")`.

    Security (CodeQL #8, py/path-injection): the sessions-dir override is a
    client-controlled string. In normal operation the web UI never sends it --
    it relies on the server-configured directory -- so honouring an arbitrary
    absolute path here would only ever help an attacker read/write session
    files outside the intended data area. We therefore accept the override only
    when it resolves inside a known-safe root (the user's home directory, which
    holds the default ~/.nyxGPT data area, or the system temp directory used by
    tests and ephemeral runs). Anything else is refused (returns None -> caller
    falls back to the configured default). The `resolve()` + `relative_to()`
    containment check also neutralises `..` traversal.
    """
    if not s:
        return None
    # Normalise the client string with os.path.realpath, then require the
    # result to sit inside a fixed safe root via a string-prefix check. This is
    # the canonical CodeQL-recognised path-injection barrier (py/path-injection)
    # -- Path.relative_to() is NOT modelled as a sanitizer, so the guarded value
    # must be produced by realpath and gated by a single `startswith(root+sep)`
    # check before it is ever turned back into a Path. realpath also collapses `..`
    # traversal. The web UI never sends this override (it uses the configured
    # directory); accepting an arbitrary absolute path would only help an
    # attacker read/write session files outside the intended data area.
    try:
        real = os.path.realpath(os.path.expanduser(s))
    except (OSError, ValueError):
        log.warning("Refused unresolvable sessions-dir override: %r", s)
        return None
    _home = os.path.realpath(os.path.expanduser("~"))
    _tmp = os.path.realpath(tempfile.gettempdir())
    # Each return below is controlled by exactly one condition: CodeQL's
    # barrier-guard analysis only credits a guard whose branch is dominated
    # by a single sanitizing comparison, never a disjunction or loop of them.
    if real.startswith(_home + os.sep):
        return Path(real)
    if real.startswith(_tmp + os.sep):
        return Path(real)
    log.warning("Refused sessions-dir override outside allowed roots: %r", s)
    return None


# ----------------------------
# Routes
# ----------------------------


@app.get("/health")
def health() -> dict[str, str]:
    """Liveness check. Always returns `{"status": "ok"}` if the process is up.

    Unauthenticated and unversioned so load balancers/orchestrators can
    probe it without an API key.
    """
    return {"status": "ok"}


@app.get("/metrics")
def prometheus_metrics_endpoint() -> Response:
    """Expose metrics in Prometheus text exposition format for scraping.

    Includes HTTP request counts, latency histograms, and error rates
    (all endpoints) plus business metrics (chat requests, RAG queries) and
    a live snapshot of process resource usage (memory, CPU, queue depth).
    Left unauthenticated like /health since Prometheus scrapers typically
    don't send an API key.
    """
    monitor = get_resource_monitor()
    if monitor is not None:
        snapshot = monitor.get_metrics()
        prom_metrics.update_resource_gauges(
            rss_mb=snapshot.memory_rss_mb,
            cpu_percent=snapshot.cpu_percent_process,
            queue_depth=snapshot.queue_depth,
            disk_percent=snapshot.disk_percent,
            memory_percent=snapshot.memory_percent,
        )
    body, content_type = prom_metrics.render_metrics()
    return Response(content=body, media_type=content_type)


@api.get("/info", response_model=InfoResponse)
def info(request: Request) -> InfoResponse:
    """Return basic server info: Ollama URL, default model, sessions dir, and running version.

    `release_version` is the version of the installed `nyxgpt` package -- what
    is actually running -- not a configuration value. The agent tooling's
    `[github] RELEASE_BRANCH` setting is reported separately as
    `release_branch` so the two can never be confused (#3716).
    """
    cfg = _req_cfg(request)
    return InfoResponse(
        ollama_base_url=get_ollama_base_url(cfg),
        default_model=get_default_model(cfg),
        sessions_dir=str(get_sessions_dir(cfg)),
        release_version=running_version(),
        release_branch=cfg.get("github", "RELEASE_BRANCH", fallback=None),
    )


@api.get("/batch/metrics")
def batch_metrics() -> dict[str, Any]:
    """Get batch processing metrics.

    Returns:
        Dictionary with batch metrics including total requests, avg batch size,
        throughput, wait times, etc. Returns empty metrics if batching disabled.
    """
    if _batch_processor is None:
        return {
            "enabled": False,
            "message": "Request batching is not enabled",
        }

    metrics = _batch_processor.get_metrics()
    return {
        "enabled": True,
        **metrics.to_dict(),
    }


@api.get("/metrics", response_model=ResourceMetricsResponse)
def resource_metrics() -> ResourceMetricsResponse:
    """Get system resource usage metrics.

    Returns comprehensive resource monitoring including:
    - Memory usage (RSS, VMS, available)
    - CPU utilization (process and system)
    - Request latency (avg, p50, p95, p99)
    - Queue depth (current and total requests)

    Returns:
        ResourceMetricsResponse with current resource metrics
    """
    monitor = get_resource_monitor()
    if monitor is None:
        # Return zero metrics if monitor not initialized
        return ResourceMetricsResponse(
            memory={"rss_mb": 0, "vms_mb": 0, "percent": 0, "available_mb": 0},
            cpu={"process_percent": 0, "system_percent": 0},
            latency={"avg_ms": 0, "p50_ms": 0, "p95_ms": 0, "p99_ms": 0},
            queue={"depth": 0, "total_requests": 0},
        )

    metrics = monitor.get_metrics()
    return ResourceMetricsResponse(**metrics.to_dict())


@api.get("/metrics/history")
def resource_metrics_history(
    range: str = Query(  # noqa: A002 - "range" is the public query param name
        "1h", pattern="^(1h|24h|7d)$", description="Time window: 1h, 24h, or 7d"
    ),
) -> dict[str, Any]:
    """Get server-side historical resource usage metrics for a time window.

    Unlike `/api/v1/metrics` (a current-snapshot endpoint), this returns a
    downsampled time series sampled server-side once a minute and persisted
    to disk, so history exists independent of any open browser tab and
    survives an API restart. `history_available_seconds` tells the caller
    honestly how much of the requested window is actually backed by data,
    so the UI can avoid rendering a misleadingly full-width chart on a
    fresh install.

    Args:
        range: Requested window -- "1h", "24h", or "7d".

    Returns:
        Dict with the requested range, downsampled points, sample cadence,
        and available-history accounting.
    """
    return resource_metrics_store_module.query_history(range)


def _jaeger_curated_views(jaeger_ui_url: str, service_name: str) -> list[dict[str, str]]:
    """Build curated Jaeger search links for the main request flows.

    Each entry deep-links to the service's trace search (a query Jaeger
    always supports) and names the operation(s) to pick from the dropdown
    for that flow, rather than guessing an exact `operation` URL param --
    the auto-instrumented span name for a given FastAPI route can vary by
    library version, so a wrong guess there would silently show "no
    results" instead of failing loudly.
    """
    search_url = f"{jaeger_ui_url}/search?service={service_name}&lookback=1h"
    return [
        {
            "label": "Chat requests",
            "hint": "Filter by operation: POST /api/v1/chat or POST /api/v1/chat/stream",
            "url": search_url,
        },
        {
            "label": "RAG query",
            "hint": "Filter by operation: rag.retrieve (the retrieval pipeline span)",
            "url": search_url,
        },
        {
            "label": "RAG ingest",
            "hint": "Filter by operation: POST /api/v1/rag/ingest, /rag/upload, or /rag/index-repo",
            "url": search_url,
        },
        {
            "label": "Ollama backend calls",
            "hint": "Filter by operation: ollama.request, ollama.request.stream, or ollama.embeddings",
            "url": search_url,
        },
    ]


@api.get("/tracing")
def tracing_status(request: Request) -> dict[str, Any]:
    """Get distributed tracing status and how to reach the local Jaeger UI.

    Tracing is opt-in and local-only: enable it with `[tracing] enabled =
    true` in config.ini and start the API alongside the `tracing` Compose
    profile (local OTel collector + Jaeger all-in-one). No spans are ever
    exported outside the machine.
    """
    cfg = _req_cfg(request)
    tracing_config = get_tracing_config(cfg)
    active = tracing_module.is_tracing_enabled()
    return {
        **tracing_config,
        "active": active,
        # `active` only means init_tracing() ran -- it says nothing about
        # whether the OTLP collector is actually reachable (see #3350, where
        # otel-collector published no host port and every span was silently
        # dropped while this still reported "active"). Only probed when
        # active, and with a short timeout, so a healthy install's panel
        # load isn't delayed and a disabled install never pays for a connect.
        "reachable": (
            tracing_module.otlp_endpoint_reachable(tracing_config["otlp_endpoint"], timeout=0.5)
            if active
            else None
        ),
        "curated_views": _jaeger_curated_views(
            tracing_config["jaeger_ui_url"], tracing_config["service_name"]
        ),
    }


@api.get("/error-tracking")
def error_tracking_status(request: Request) -> dict[str, Any]:
    """Get error tracking status and how to reach the local GlitchTip UI.

    Error tracking is opt-in and local-only: enable it with
    `[error_tracking] enabled = true` and a `dsn` pointed at a self-hosted
    tracker in config.ini, then start the API alongside the `errors`
    Compose profile (local GlitchTip instance). No exception data is ever
    sent to Sentry's own SaaS.
    """
    cfg = _req_cfg(request)
    error_tracking_config = get_error_tracking_config(cfg)
    return {
        **error_tracking_config,
        "active": error_tracking_module.is_error_tracking_enabled(),
    }


@api.post("/error-tracking/report")
def error_tracking_report(request: Request, body: api_models.ClientErrorReportRequest) -> Response:
    """Forward a web UI client-side error to the local error tracker.

    Returns `503 {"status": "inactive"}` when error tracking isn't actually
    initialized (disabled, or enabled with no valid DSN), rather than a
    blanket `202` -- a misconfigured DSN or a tracker nobody enabled must not
    look like a successfully delivered test event. The web UI's fire-and-forget
    client error reporter (`ClientErrorReporter.tsx`) ignores the response
    either way, so this stays safe to call whether or not tracking is active.
    """
    if not error_tracking_module.is_error_tracking_enabled():
        return JSONResponse(
            status_code=503,
            content={
                "status": "inactive",
                "detail": "Error tracking is disabled or has no valid DSN configured; "
                "this event was not sent.",
            },
        )

    req_id = getattr(request.state, "request_id", None)
    error_tracking_module.capture_message(
        body.message,
        request_id=req_id or "N/A",
        source="web-client",
        url=body.url or "N/A",
        stack=body.stack or "N/A",
    )
    return JSONResponse(status_code=202, content={"status": "accepted"})


@api.get("/monitoring")
def monitoring_status(request: Request) -> dict[str, Any]:
    """Get monitoring stack status and how to reach the local Grafana UI.

    The Grafana/Prometheus stack is opt-in and local-only: enable it with
    `[monitoring] enabled = true` in config.ini and start the API alongside
    the `monitoring` Compose profile (local Prometheus + Grafana). No
    metrics are ever exported outside the machine.
    """
    cfg = _req_cfg(request)
    monitoring_config = get_monitoring_config(cfg)
    return {
        **monitoring_config,
        "active": monitoring_config["enabled"],
    }


def _loki_curated_queries() -> list[dict[str, str]]:
    """Curated LogQL saved queries for the nyxGPT job's log stream.

    Mirrors the per-component panels already provisioned in the "Operational
    Logs" Grafana dashboard (docker/grafana/dashboards/operational-logs.json)
    so the same queries are also copy-pasteable into Grafana Explore for ad
    hoc search, not just fixed dashboard panels. Query text only, not a
    Grafana Explore deep-link URL -- like `_jaeger_curated_views` above,
    guessing a URL-encoded state param that's changed shape across Grafana
    major versions risks silently opening Explore without the query loaded;
    the dashboard panels are the reliable path, this is a documented
    fallback for ad hoc search.
    """
    return [
        {
            "label": "Self-heal events",
            "hint": "Restart / recovery / backoff activity",
            "query": '{job="nyxgpt", logger="nyxgpt.self_heal"}',
        },
        {
            "label": "Canary events",
            "hint": "Canary deploy, rollout, evaluation, and promotion activity",
            "query": '{job="nyxgpt", logger="nyxgpt.canary"}',
        },
        {
            "label": "Chat errors",
            "hint": "ERROR/CRITICAL log lines from the chat pipeline",
            "query": '{job="nyxgpt", logger="nyxgpt.chat", level=~"ERROR|CRITICAL"}',
        },
        {
            "label": "RAG pipeline",
            "hint": "Retrieval and ingest activity",
            "query": '{job="nyxgpt", logger="nyxgpt.rag.rag"}',
        },
    ]


@api.get("/log-aggregation")
def log_aggregation_status(request: Request) -> dict[str, Any]:
    """Get log aggregation status and how to reach the local Loki search UI.

    Log aggregation is opt-in and local-only: enable it with
    `[log_aggregation] enabled = true` in config.ini and start the API
    alongside the `logging` Compose profile (local Loki + promtail).
    Promtail ships the API's log files under ~/.nyxGPT/logs into Loki, which
    is searched via Grafana Explore (the `monitoring` profile). No logs are
    ever exported outside the machine.
    """
    cfg = _req_cfg(request)
    log_aggregation_config = get_log_aggregation_config(cfg)
    return {
        **log_aggregation_config,
        "active": log_aggregation_config["enabled"],
        "curated_queries": _loki_curated_queries(),
    }


# --- Config get/set endpoints ---


@api.get("/config")
def config_get(request: Request) -> dict[str, Any]:
    """Return the current hot-configurable settings: Ollama URL, default model, RAG enabled, log level."""
    cfg = _req_cfg(request)
    return {
        "ollama_base_url": get_ollama_base_url(cfg),
        "default_model": get_default_model(cfg),
        "rag_enabled": get_rag_enabled(cfg),
        "log_level": cfg.get("logging", "level", fallback="INFO").strip().upper(),
    }


@api.post("/config")
def config_update(request: Request, payload: dict[str, Any] = Body(...)) -> dict[str, Any]:
    """Update `default_model`, `rag_enabled`, and/or `log_level` in config.ini.

    Unknown keys in `payload` are silently ignored. Applied changes take
    effect immediately (no restart) and are recorded to the admin activity
    log. Returns the set of fields that were changed plus the resulting
    effective config values.
    """
    # Only apply known keys; ignore the rest.
    updates: dict[str, Any] = {}
    if "default_model" in payload:
        updates["default_model"] = payload.get("default_model")
    if "rag_enabled" in payload:
        updates["rag_enabled"] = payload.get("rag_enabled")
    if "log_level" in payload:
        updates["log_level"] = payload.get("log_level")

    changed = _apply_hot_config_updates(updates)
    if changed:
        admin_activity_module.record(
            "config.updated", ", ".join(f"{k}={v}" for k, v in changed.items())
        )

    # After applying, reload config and update request.state.cfg
    request.state.cfg = load_config(None)
    cfg = _req_cfg(request)
    return {
        "updated": changed,
        "effective": {
            "ollama_base_url": get_ollama_base_url(cfg),
            "default_model": get_default_model(cfg),
            "rag_enabled": get_rag_enabled(cfg),
            "log_level": cfg.get("logging", "level", fallback="INFO").strip().upper(),
        },
    }


# PATCH endpoint for config updates
@api.patch("/config")
def config_patch(request: Request, payload: dict[str, Any] = Body(...)) -> dict[str, Any]:
    """PATCH alias for `POST /config` (see `config_update`)."""
    result: dict[str, Any] = config_update(request, payload)
    return result


# --- Full config wizard endpoints (#3354) ---


def _reconcile_observability(cfg: ConfigParser) -> dict[str, Any]:
    """Bring the observability Compose stack in line with config.ini's enabled flags.

    Called after a wizard save touches any of the
    monitoring/tracing/error_tracking/log_aggregation `enabled` fields, so
    the wizard results in a working feature rather than a dangling flag.
    Wraps `ops.reconcile_observability`, which itself mirrors `nyxgpt ops
    observability` / `nyxgpt ops stop --target observability` -- this never
    shells out to `docker compose` directly.
    """
    any_enabled = (
        get_monitoring_config(cfg)["enabled"]
        or get_tracing_config(cfg)["enabled"]
        or get_error_tracking_config(cfg)["enabled"]
        or get_log_aggregation_config(cfg)["enabled"]
    )
    try:
        results = ops_module.reconcile_observability(any_enabled)
        return {"ok": all(r.ok for r in results), "messages": [r.message for r in results]}
    except Exception:
        # The full exception is logged server-side; the response carries only a
        # generic message so no internal detail reaches the API client
        # (CodeQL py/stack-trace-exposure, alert #24).
        log.exception("Observability reconciliation failed")
        return {"ok": False, "messages": ["Observability reconciliation failed"]}


@api.get("/config/sections")
def config_sections_get(request: Request) -> dict[str, Any]:
    """Return every wizard-editable config field, grouped by section.

    Backs the full Configuration Wizard (#3354): current values (secrets
    masked) plus schema metadata (which fields are secret, and which need a
    service restart or observability reconciliation to take effect) so the
    web UI can render every section without hardcoding that knowledge twice.
    `stale_keys` surfaces config.ini drift (#3388): keys present on disk but
    no longer declared in `example.config.ini`, for the wizard to offer
    removing -- never removed automatically.
    """
    cfg = _req_cfg(request)
    return {
        "sections": config_wizard.read_sections(cfg),
        "schema": config_wizard.schema_summary(),
        "field_defaults": config_wizard.field_defaults(cfg),
        "stale_keys": config_wizard.find_stale_keys(cfg),
    }


@api.post("/config/sections/stale-keys/remove")
def config_stale_keys_remove(
    request: Request, payload: dict[str, Any] = Body(...)
) -> dict[str, Any]:
    """Delete specific stale keys the wizard reported, on explicit user confirmation (#3388).

    Payload shape: `{"remove": {section: [key, ...]}}`. Only keys the caller
    names are ever touched -- this is the sole path that can delete anything
    from config.ini; a regular wizard save (`POST /config/sections`) never
    does. Only keys `config_wizard.find_stale_keys` currently reports are
    accepted, so this can't be used to delete an arbitrary managed field.
    """
    cfg = _req_cfg(request)
    requested = payload.get("remove")
    if not isinstance(requested, dict):
        raise HTTPException(
            status_code=400, detail="'remove' must be an object of {section: [keys]}"
        )

    stale = config_wizard.find_stale_keys(cfg)
    to_remove: dict[str, list[str]] = {}
    for section, keys in requested.items():
        if not isinstance(keys, list):
            raise HTTPException(status_code=400, detail=f"remove.{section} must be a list of keys")
        allowed = set(stale.get(section, []))
        matched = [k for k in keys if k in allowed]
        if matched:
            to_remove[section] = matched

    removed = config_wizard.remove_keys(_config_file_path(), to_remove)

    if removed:
        nyxgpt.config._CACHED_CFG = None
        nyxgpt.config._CACHED_PATH = None
        nyxgpt.config._CACHED_MTIME_NS = None
        request.state.cfg = load_config(None)
        cfg = _req_cfg(request)

        removed_fields = [f"{section}.{key}" for section, keys in removed.items() for key in keys]
        admin_activity_module.record("config.stale_keys_removed", ", ".join(removed_fields))

    return {
        "removed": removed,
        "sections": config_wizard.read_sections(cfg),
        "stale_keys": config_wizard.find_stale_keys(cfg),
    }


@api.post("/config/sections")
def config_sections_update(request: Request, payload: dict[str, Any] = Body(...)) -> dict[str, Any]:
    """Validate, apply, and reload a full-config wizard payload (#3354).

    Payload shape: `{section: {key: value}}`, validated against
    `config_wizard.WIZARD_SCHEMA`. On success, writes config.ini (the single
    source of truth, #3194), invalidates the config cache so hot-reloaded
    settings apply immediately, re-applies the log level, and -- if any
    observability `enabled` flag changed -- reconciles the Compose stack.
    Returns which touched fields still need `POST /config/restart` (or
    `nyxgpt ops restart`, offered by the same mechanism) to fully apply.
    """
    cfg = _req_cfg(request)
    validated, errors = config_wizard.validate_updates(payload)
    if errors:
        raise HTTPException(status_code=422, detail={"errors": errors})
    if not validated:
        raise HTTPException(status_code=400, detail="No valid fields to update")

    restart_detail = config_wizard.restart_required_detail(validated, cfg)
    restart_needed = sorted(restart_detail)
    needs_observability = config_wizard.observability_changed(validated, cfg)

    applied = config_wizard.apply_updates(_config_file_path(), validated)

    # Mark first, reconcile second, and the order is load-bearing (#3806).
    # Reverting key K to the value the running service still holds arrives
    # here as a *change* (disk B -> A), so `restart_detail` lists it and
    # `mark_pending` would re-add it; `mark_pending` keeps the originally
    # recorded running value (A), and `reconcile_saved` then sees the newly
    # saved A match it and retires the entry. Reconciling first would leave
    # the stale flag standing forever.
    for component, changes in restart_detail.items():
        restart_state_module.mark_pending(component, changes)
    for component, saved in config_wizard.restart_activation_saved(validated).items():
        restart_state_module.reconcile_saved(component, saved)

    nyxgpt.config._CACHED_CFG = None
    nyxgpt.config._CACHED_PATH = None
    nyxgpt.config._CACHED_MTIME_NS = None

    request.state.cfg = load_config(None)
    cfg = _req_cfg(request)

    with suppress(Exception):
        configure_logging(cfg, console=False, filename="api.log")

    observability_result = _reconcile_observability(cfg) if needs_observability else None

    changed_fields = [f"{section}.{key}" for section, fields in applied.items() for key in fields]
    admin_activity_module.record(
        "config.wizard_updated", ", ".join(changed_fields) or "config updated"
    )

    return {
        "applied": config_wizard.mask_applied(applied),
        "sections": config_wizard.read_sections(cfg),
        "field_defaults": config_wizard.field_defaults(cfg),
        "restart_required": restart_needed,
        # The full pending set, not just what this save added: the wizard's
        # notice must show everything still awaiting a restart (including a
        # key rotated earlier from the CLI), and must *disappear* when this
        # save reverted the last pending key (#3806).
        "restart_pending": restart_state_module.snapshot(),
        "observability_reconciled": needs_observability,
        "observability_result": observability_result,
    }


# --- Guided secrets / AWS credentials status endpoints (#3505, #3512) ---
#
# Read-only by owner decision (#3805). The matching write endpoints
# (`POST /config/secrets`, `/config/secrets/sync`, `/config/aws-credentials`,
# `/config/aws-credentials/secret-store`) and the two `/admin` screens that
# drove them are gone: a credential typed into a browser crosses an HTTP
# request and the page's process on its way to disk -- and over a cloud
# access tunnel it would cross that path too -- while `nyxgpt secrets setup`,
# `nyxgpt ops secrets-sync` and `nyxgpt cloud credentials-setup` take masked
# input and write straight to config.ini, ~/.aws/credentials or the OS
# keychain. By the time this API is answering, reaching it already required
# the secrets those screens collected. What remains here reports *whether*
# something is configured, never a cleartext value, and never accepts one.
# Do not re-add a write path: the Configuration Wizard
# (`/config/sections`, which excludes `[cloud]` and the `openai`/`github`
# agent-system sections) is the sanctioned in-product configuration surface.


@api.get("/config/secrets")
def config_secrets_get(request: Request) -> dict[str, Any]:
    """Return the guided secrets' metadata plus each one's current set/masked state.

    The machine-readable counterpart of `nyxgpt secrets setup` (same
    `secrets_setup` module, so the two can't drift about what the closed
    `GUIDED_SECRETS` set is or which of them are set). Never returns
    cleartext; setting a value is a CLI operation (#3805).
    """
    cfg = _req_cfg(request)
    return {"secrets": secrets_setup.secret_status(cfg)}


@api.get("/config/aws-credentials")
def config_aws_credentials_get(request: Request) -> dict[str, Any]:
    """Return the AWS credentials status: fields, `[cloud]` reference, and where the
    access key pair is currently stored (masked). Never returns cleartext.

    The machine-readable counterpart of `nyxgpt cloud credentials-setup`,
    which is where the key pair is entered (#3805).
    """
    cfg = _req_cfg(request)
    return aws_credentials_setup.aws_credentials_status(cfg)


_RESTART_TARGETS = {"api", "web", "ollama", "cassandra", "observability", "all"}


@api.post("/config/restart")
def config_restart(_request: Request, payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Restart a service so wizard changes that need a process bounce take effect (#3354).

    Wraps `nyxgpt ops restart` -- the wizard offers this after a save
    instead of ever telling the user to run a raw restart command
    themselves (the wrapper rule). The restart is scheduled a moment after
    this response is sent, since the target may be this very API process.
    """
    target = payload.get("target") or "all"
    if target not in _RESTART_TARGETS:
        raise HTTPException(
            status_code=400, detail=f"target must be one of {sorted(_RESTART_TARGETS)}"
        )

    admin_activity_module.record("config.restart_requested", f"target={target}")

    def _do_restart() -> None:
        try:
            ops_module.restart(SimpleNamespace(target=target))
        except Exception:
            log.exception("ops restart failed (target=%s)", target)

    threading.Timer(0.75, _do_restart).start()

    return {"target": target, "status": "scheduled"}


@api.get("/infra/restart-status")
def infra_restart_status() -> dict[str, Any]:
    """Config changes saved but not yet in effect, and what has to restart (#3407, #3806).

    Backs the persistent pending-restart notice shown on both the Admin
    Dashboard and the Configuration Wizard: `pending` maps each `nyxgpt ops
    restart` target to the `section.key` fields whose saved value differs from
    the value that service is still running with, and when that divergence
    started. Empty once every flagged component has actually been restarted --
    via `POST /infra/restart-required`, via `nyxgpt ops restart`, or because
    the value was reverted to what the service already had.

    The state is read from disk (`restart_state`), so it is the same set the
    CLI reports and it survives this process restarting -- which matters
    precisely for the `web` entries, whose restart does not touch this
    process.

    `restart_command` is the wrapped command that clears the whole set, so the
    UI can show the user the CLI equivalent of its own button rather than a
    raw `brew services`/`docker` invocation.
    """
    pending = restart_state_module.snapshot()
    return {
        "pending": pending,
        "restart_command": (
            restart_state_module.restart_command(sorted(pending)) if pending else None
        ),
        # Restarting `web` from a page served by `web` drops the browser's
        # connection to the very server rendering it. The UI states this
        # before it happens (IntelliJ-style) instead of appearing to hang, so
        # the backend is the one place that decides when it applies.
        "session_disrupting": sorted(c for c in pending if c == "web"),
    }


def _do_restart_required(targets: list[str]) -> None:
    """Mode-aware restart of each of `targets`, run off the request thread (#3407).

    Reuses `self_heal.heal_now(service=...)`, the same dispatcher the
    dashboard's manual "Heal Now" button already uses (#3193/#3344): it picks
    native/Compose/Terraform/Kubernetes restart behavior to match whatever's
    actually running, bypassing the health check so a healthy-but-stale
    component still restarts. Restarting `api` kills this very process once
    the underlying `brew services restart`/`docker restart`/`kubectl delete
    pod` command lands -- deferred a moment so the triggering HTTP response
    can be sent first, same as `config_restart` above.
    """
    for component in targets:
        try:
            result = self_heal_module.heal_now(service=component)
        except Exception:
            log.exception("mode-aware restart-required failed (component=%s)", component)
            continue
        healed = result.get("healed", [])
        for event in healed:
            admin_activity_module.record(
                "infra.restart_required", f"{event['service']}: {event['message']}"
            )
            # Recorded as an ops lifecycle action too (#3390), same as a manual
            # Heal Now click -- this is equally an operator-initiated restart.
            ops_module.record_manual_restart(event["service"], event["ok"], event["message"])
        if healed and all(e["ok"] for e in healed):
            restart_state_module.clear_pending(component)


@api.post("/infra/restart-required")
def infra_restart_required(payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Restart whichever component(s) a wizard save flagged as pending (#3407).

    With no `target`, restarts every currently pending component. Mode-aware
    (native/Compose/Terraform/Kubernetes) via `_do_restart_required` -- the
    caller never needs to know or type which raw command applies. Runs
    off-thread, so the response reports "running"; the caller (the Admin
    Dashboard restart button) polls `GET /infra/restart-status` to learn
    when the pending flag clears.
    """
    pending = restart_state_module.snapshot()
    target = payload.get("target")
    if target is not None:
        if target not in pending:
            raise HTTPException(
                status_code=400, detail=f"No restart is currently pending for '{target}'"
            )
        targets = [target]
    else:
        targets = sorted(pending)
        if not targets:
            raise HTTPException(status_code=400, detail="No restart is currently pending")

    admin_activity_module.record("infra.restart_required_requested", ", ".join(targets))
    threading.Timer(0.5, _do_restart_required, args=(targets,)).start()

    return {"targets": targets, "status": "running"}


# --- Admin dashboard endpoints (system status overview, activity log, access) ---


@api.get("/admin/overview")
def admin_overview(request: Request) -> dict[str, Any]:
    """Aggregate system status for the admin dashboard overview panel.

    Combines app info, resource metrics, canary status, and the
    enabled/disabled state of opt-in observability stacks into a single
    response so the dashboard can render a status summary in one request.
    Individual sub-sections degrade to an `{"error": ...}` payload instead
    of failing the whole request when a backing service (e.g. a local K8s
    cluster for canary) is unavailable.
    """
    cfg = _req_cfg(request)

    def _safe(fn, *args, **kwargs) -> dict[str, Any]:
        """Call `fn`, degrading to a generic `{"error": ...}` payload on failure.

        The real exception is logged server-side; the response never carries
        the raw exception detail (CodeQL py/stack-trace-exposure, alert #25).
        """
        try:
            result: dict[str, Any] = fn(*args, **kwargs)
            return result
        except Exception:
            log.exception("admin/overview sub-section failed: %s", getattr(fn, "__name__", fn))
            return {"error": "unavailable"}

    monitor = get_resource_monitor()
    resource_metrics_summary = monitor.get_metrics().to_dict() if monitor is not None else None

    return {
        "info": {
            "ollama_base_url": get_ollama_base_url(cfg),
            "default_model": get_default_model(cfg),
            "rag_enabled": get_rag_enabled(cfg),
        },
        "resource_metrics": resource_metrics_summary,
        "canary": _safe(canary_module.status, get_canary_namespace(cfg)),
        "self_heal": _safe(self_heal_module.status),
        "observability": {
            "monitoring": get_monitoring_config(cfg)["enabled"],
            "tracing": get_tracing_config(cfg)["enabled"],
            "error_tracking": get_error_tracking_config(cfg)["enabled"],
            "log_aggregation": get_log_aggregation_config(cfg)["enabled"],
        },
        "auth_enabled": _auth_cfg(cfg)["enabled"],
    }


@api.get("/admin/health")
def admin_health(request: Request) -> dict[str, Any]:
    """Aggregate system health for the admin health dashboard.

    Combines service uptime, dependency reachability checks (Ollama, and
    Cassandra when RAG is enabled), resource utilization, and
    threshold-based alert indicators into a single response.
    """
    cfg = _req_cfg(request)
    rag_enabled = get_rag_enabled(cfg)

    monitor = get_resource_monitor()
    resource_metrics_summary = monitor.get_metrics().to_dict() if monitor is not None else None

    dependencies = [
        health_module.check_ollama(get_ollama_base_url(cfg)),
        health_module.check_cassandra(rag_enabled),
    ]
    grafana_alerts = health_module.fetch_grafana_alerts(cfg)
    if grafana_alerts is not None:
        alerts = grafana_alerts
        alerts_source = "grafana"
    else:
        alerts = health_module.compute_alerts(resource_metrics_summary, dependencies)
        alerts_source = "local"

    return {
        "service": {"status": "ok", "uptime_s": round(health_module.uptime_seconds(), 1)},
        "dependencies": [d.to_dict() for d in dependencies],
        "resource_metrics": resource_metrics_summary,
        "alerts_source": alerts_source,
        "alerts": [a.to_dict() for a in alerts],
    }


@api.get("/admin/activity")
def admin_activity_list(request: Request, limit: int = 50) -> dict[str, Any]:
    """Return recent admin dashboard activity (audit trail)."""
    cfg = _req_cfg(request)
    bounded_limit = max(1, min(limit, 500))
    return {"events": admin_activity_module.recent(bounded_limit, cfg=cfg)}


@api.get("/admin/access")
def admin_access_get(request: Request) -> dict[str, Any]:
    """Return the current API-key access configuration (key masked)."""
    cfg = _req_cfg(request)
    auth = _auth_cfg(cfg)
    return {
        "enabled": auth["enabled"],
        "header": auth["header"],
        "api_key_set": bool(auth["api_key"]),
        "api_key_masked": _mask_api_key(auth["api_key"]) if auth["api_key"] else None,
    }


@api.post("/admin/access")
def admin_access_update(
    request: Request, payload: dict[str, Any] = Body(default={})
) -> dict[str, Any]:
    """Update API-key access configuration: enable/disable auth, change the
    header name, or rotate the API key.

    On rotation, the newly generated key is returned once in the response
    body (`api_key`) so the operator can copy it; it is never returned
    again by `GET /admin/access`, which only shows a masked value.
    """
    cfg = _req_cfg(request)
    auth = _auth_cfg(cfg)

    updates: dict[str, Any] = {}
    if "enabled" in payload:
        updates["enabled"] = bool(payload["enabled"])
    if "header" in payload and isinstance(payload.get("header"), str) and payload["header"].strip():
        updates["header"] = payload["header"]

    rotate = bool(payload.get("rotate"))
    new_key: str | None = None
    if rotate:
        # When a cloud secrets provider is configured, `get_auth_api_key`
        # always prefers the AWS-resolved value over config.ini (see
        # `_resolve_cloud_secret`), so writing a new key to config.ini here
        # would be inert: the middleware would keep enforcing the old
        # cloud-stored key while this endpoint reported the new one as
        # active. Reject explicitly instead of silently no-op'ing.
        if get_secrets_provider(cfg):
            raise HTTPException(
                status_code=400,
                detail=(
                    "API key rotation via this dashboard is disabled because a cloud "
                    "secrets provider ([secrets] provider) is configured. Rotate the "
                    "value in AWS SSM Parameter Store or Secrets Manager instead -- "
                    "see docs/cloud.md for the rotation procedure."
                ),
            )
        new_key = secrets.token_urlsafe(32)
        updates["api_key"] = new_key

    if not updates:
        raise HTTPException(status_code=400, detail="No valid fields to update")

    _apply_auth_config_updates(updates)

    request.state.cfg = load_config(None)
    cfg = _req_cfg(request)
    auth = _auth_cfg(cfg)

    changes = []
    if "enabled" in updates:
        changes.append("enabled" if updates["enabled"] else "disabled")
    if "header" in updates:
        changes.append("header changed")
    if rotate:
        changes.append("API key rotated")
    admin_activity_module.record("access.updated", ", ".join(changes) or "access settings updated")

    result: dict[str, Any] = {
        "enabled": auth["enabled"],
        "header": auth["header"],
        "api_key_set": bool(auth["api_key"]),
        "api_key_masked": _mask_api_key(auth["api_key"]) if auth["api_key"] else None,
    }
    if new_key is not None:
        result["api_key"] = new_key
    return result


@api.get("/analytics/usage")
def analytics_usage(request: Request) -> dict[str, Any]:
    """Return aggregated usage analytics for the admin dashboard.

    Summarizes recorded chat requests: total requests/tokens, distinct
    session count, and per-model and per-day breakdowns.
    """
    cfg = _req_cfg(request)
    return usage_analytics_module.summary(cfg=cfg)


@api.get("/analytics/export")
def analytics_export(request: Request, format: str = "json") -> Response:
    """Export recorded usage analytics as a downloadable JSON or CSV report."""
    cfg = _req_cfg(request)
    try:
        content, content_type, filename = usage_analytics_module.export_report(format, cfg=cfg)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e)) from e
    return Response(
        content=content,
        media_type=content_type,
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


# --- Local canary deployment endpoints (SRE/admin dashboard) ---
# Blue/green (deploy.py, `/api/v1/deploy/*`) was retired in favor of canary,
# a strict superset for traffic purposes (0%/100% reproduces a cutover) plus
# metrics-gated gradual shift and auto-rollback -- see #3409.
@api.get("/canary/status")
def canary_status(request: Request, component: str = "api") -> dict[str, Any]:
    """Return canary rollout status: active flag, traffic weight, stable/canary health/version,
    the currently detected deployment mode, and a metrics snapshot.

    `component` (query param, default `api`) selects which component's pair
    to report on -- `api` or `web` (see canary.py's `COMPONENTS`, #3419).
    """
    cfg = _req_cfg(request)
    return canary_module.status(get_canary_namespace(cfg), component=component)


@api.post("/canary/deploy")
def canary_deploy(request: Request, payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Build the current checkout into a versioned image and deploy it to canary only.

    Body: `{"component": str}` (default `api`; `api` or `web`, #3419). Never
    touches stable, even on failure. Traffic weighting is a separate,
    deliberate action (`/canary/start` and `/canary/promote`). Returns `409`
    if the build, image patch, or rollout wait fails. Records a
    `canary.deploy` admin activity event on success.
    """
    cfg = _req_cfg(request)
    component = payload.get("component", "api")
    result = canary_module.deploy(namespace=get_canary_namespace(cfg), component=component)
    if not result.ok:
        raise HTTPException(status_code=409, detail=result.message)
    admin_activity_module.record("canary.deploy", result.message)
    return {"ok": result.ok, "message": result.message}


@api.post("/canary/start")
def canary_start(request: Request, payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Start a canary rollout: scale up the canary Deployment to `weight_percent` of traffic.

    Body: `{"weight_percent": int, "component": str}` (weight default 10;
    component default `api`, #3409/#3419). Returns `409` if a rollout is
    already in progress or the scale operation fails. Records a
    `canary.start` admin activity event on success.
    """
    cfg = _req_cfg(request)
    weight_percent = int(payload.get("weight_percent", 10))
    component = payload.get("component", "api")
    result = canary_module.start(
        namespace=get_canary_namespace(cfg),
        weight_percent=weight_percent,
        total_replicas=get_canary_total_replicas(cfg),
        component=component,
    )
    if not result.ok:
        raise HTTPException(status_code=409, detail=result.message)
    admin_activity_module.record("canary.start", result.message)
    return {"ok": result.ok, "message": result.message}


@api.post("/canary/evaluate")
def canary_evaluate(request: Request, payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Compare live error-rate/latency metrics against configured thresholds.

    Body: `{"component": str}` (default `api`, #3419). Automatically rolls
    back the canary if either threshold is breached. Returns `409` if no
    rollout is active. Returns ok=True with an "insufficient data" note
    (rather than failing) when too few requests have been observed yet to
    judge the canary. Records a `canary.evaluate` admin activity event.
    """
    cfg = _req_cfg(request)
    component = payload.get("component", "api")
    result = canary_module.evaluate(
        get_canary_namespace(cfg),
        error_rate_threshold_percent=get_canary_error_rate_threshold(cfg),
        latency_p95_threshold_ms=get_canary_latency_p95_threshold_ms(cfg),
        min_requests=get_canary_min_requests(cfg),
        component=component,
    )
    if not result.ok:
        raise HTTPException(status_code=409, detail=result.message)
    admin_activity_module.record("canary.evaluate", result.message)
    return {"ok": result.ok, "message": result.message}


@api.post("/canary/promote")
def canary_promote(request: Request, payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Increase the canary's traffic share by `step_percent`, finalizing at 100%.

    Body: `{"step_percent": int, "component": str}` (step defaults to the
    configured `canary_step_percent` when omitted; component default `api`,
    #3419). Returns `409` if no rollout is active or the scale operation
    fails. Records a `canary.promote` admin activity event on success.
    """
    cfg = _req_cfg(request)
    step_percent = payload.get("step_percent")
    component = payload.get("component", "api")
    result = canary_module.promote(
        namespace=get_canary_namespace(cfg),
        step_percent=(
            int(step_percent) if step_percent is not None else get_canary_step_percent(cfg)
        ),
        total_replicas=get_canary_total_replicas(cfg),
        component=component,
    )
    if not result.ok:
        raise HTTPException(status_code=409, detail=result.message)
    admin_activity_module.record("canary.promote", result.message)
    return {"ok": result.ok, "message": result.message}


@api.post("/canary/rollback")
def canary_rollback(request: Request, payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Cut all traffic back to the stable Deployment, the emergency escape hatch for a bad canary.

    Body: `{"component": str}` (default `api`, #3419). Scales the canary
    Deployment to 0 first (removing it from the Service's endpoints) before
    restoring the stable Deployment's full replica count. Returns `409` if
    no rollout is active or the scale operation fails. Records a
    `canary.rollback` admin activity event on success.
    """
    cfg = _req_cfg(request)
    component = payload.get("component", "api")
    result = canary_module.rollback(
        namespace=get_canary_namespace(cfg),
        total_replicas=get_canary_total_replicas(cfg),
        component=component,
    )
    if not result.ok:
        raise HTTPException(status_code=409, detail=result.message)
    admin_activity_module.record("canary.rollback", result.message)
    return {"ok": result.ok, "message": result.message}


# --- Self-heal watchdog endpoints (SRE/admin dashboard) ---
@api.get("/self-heal/status")
def self_heal_status(_request: Request) -> dict[str, Any]:
    """Per-component health of the Docker Compose stack, plus recent heal events."""
    return self_heal_module.status()


@api.post("/self-heal/toggle")
def self_heal_toggle(payload: dict[str, Any] = Body(...)) -> dict[str, Any]:
    """Enable/disable the self-heal watchdog from the SRE/admin dashboard."""
    if "enabled" not in payload:
        raise HTTPException(status_code=400, detail="Missing 'enabled' field")
    enabled = self_heal_module.set_enabled(bool(payload["enabled"]))
    admin_activity_module.record("self_heal.toggle", "enabled" if enabled else "disabled")
    return {"enabled": enabled}


@api.post("/self-heal/heal")
def self_heal_heal(payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Manually trigger a heal pass -- one component (payload.service) or all of them."""
    service = payload.get("service")
    result = self_heal_module.heal_now(service=service)
    if result.get("error"):
        raise HTTPException(status_code=404, detail=result["error"])
    for event in result.get("healed", []):
        admin_activity_module.record(
            "self_heal.restart",
            f"{event['service']}: {event['message']} ({event['reason']})",
        )
        # Recorded as an ops lifecycle action too (#3390) -- a dashboard-triggered
        # heal is an operator-initiated restart, same as `nyxgpt ops restart`, distinct
        # from nyxgpt_selfheal_restarts_total (which self_heal.heal_now() already
        # incremented above for this same restart).
        ops_module.record_manual_restart(event["service"], event["ok"], event["message"])
    return result


@api.get("/self-heal/logs")
def self_heal_logs(service: str, tail: int = Query(default=200, ge=1, le=2000)) -> dict[str, Any]:
    """Recent logs for one component, from the SRE/admin dashboard.

    Dispatched by the component's actual deployment mode (Compose/native/
    Terraform/Kubernetes -- see `self_heal.component_logs`), so this reads
    the real source (e.g. a native API's own log file) rather than only ever
    checking Docker Compose. Lets an operator read a component's output
    (e.g. the GlitchTip container's first-account registration confirmation
    link, printed there by its console email backend) without running a raw
    `docker`/`docker compose`/`kubectl` command themselves.
    """
    # CodeQL #4 (py/command-line-injection): `tail` reaches subprocess argv
    # as `str(tail)` in the log dispatchers. FastAPI already coerces and
    # bounds it (int, 1..2000), but CodeQL does not model that validation.
    # Select the equal value from a trusted range so every downstream argv
    # receives an untainted int -- the same equality-selection idiom as the
    # compose-service resolution in self_heal (#3661). The fallback is
    # unreachable: Query(ge=1, le=2000) guarantees membership.
    tail = next((t for t in range(1, 2001) if t == tail), 200)
    result = self_heal_module.component_logs(service, tail=tail)
    if not result.ok:
        raise HTTPException(status_code=502, detail=result.message)
    return {"service": service, "tail": tail, "logs": result.details}


# --- Infrastructure status (SRE/admin dashboard) ---
#
# Read-only: which deployment mode is actually running (native/compose/
# terraform/kubernetes) and each mode's component state -- see
# ops.infra_status() and #3410. Terraform/Kubernetes install/destroy are
# `nyxgpt ops install|down --terraform|--kubernetes --local` CLI-only (see
# docs/terraform.md / docs/kubernetes.md); this page never mutates
# infrastructure, so there is no install/down endpoint here for the web UI
# to reach.


@api.get("/infra/status")
def infra_status(_request: Request) -> dict[str, Any]:
    """Deployment status: detected mode plus per-mode component/pod state."""
    return ops_module.infra_status()


# --- Cloud substrate endpoints (AWS, P6-8/#3509) ---
#
# The SRE/admin dashboard's counterpart to `nyxgpt cloud infra`. Per CLAUDE.md's
# Definition of Done the dashboard *observes* the cloud substrate; operating it
# is a CLI job (owner decision, 2026-08-16, #3804), because a UI served by the
# instance cannot safely change the substrate it is running on and there is no
# practical second nyxGPT to drive it from.
#
# So there is one browser-reachable substrate endpoint and it is a read.
# `apply`/`destroy` remain as API operations for a client that is not the
# instance's own dashboard -- both shell out to Terraform, can take minutes, and
# are deliberately synchronous -- but no web surface posts to them; `plan` was
# removed with the dashboard control that was its only caller (`nyxgpt cloud
# infra plan` is the surviving surface).


def _cloud_infra_args(payload: dict[str, Any]) -> argparse.Namespace:
    """Build the argparse-shaped namespace `nyxgpt.cloud_infra` expects from a JSON body.

    The dashboard posts the same inputs the CLI flags carry; anything omitted
    falls back to the settings a previous run saved, exactly as on the CLI.
    """
    return argparse.Namespace(
        region=payload.get("region") or None,
        profile=payload.get("profile") or None,
        owner_ip=payload.get("owner_ip") or None,
        ssh_public_key=payload.get("ssh_public_key") or None,
        ssh_key_name=payload.get("ssh_key_name") or None,
        instance_type=payload.get("instance_type") or None,
        root_volume_size=payload.get("root_volume_size") or None,
    )


@api.get("/cloud/infra")
def cloud_infra_status(_request: Request) -> dict[str, Any]:
    """What AWS substrate is provisioned, and how it is reachable.

    Cheap and side-effect free -- answers from instance metadata when this
    process runs on the instance, and from the recorded Terraform outputs
    otherwise -- so the dashboard can poll it.
    """
    return cloud_infra_module.infra_status()


@api.post("/cloud/infra/apply")
def cloud_infra_apply(payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Provision (or reconcile) the AWS substrate and record its ids."""
    try:
        result = cloud_infra_module.apply_infra(_cloud_infra_args(payload))
    except CloudCommandError as e:
        raise HTTPException(status_code=409, detail=str(e)) from e
    outputs = result.get("outputs") or {}
    admin_activity_module.record(
        "cloud_infra.apply",
        f"instance={outputs.get('instance_id', 'unknown')} region={result['settings']['aws_region']}",
    )
    return result


@api.post("/cloud/infra/destroy")
def cloud_infra_destroy(payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Tear the AWS substrate down.

    Requires `{"confirm": true}`: this deletes the instance and its root
    volume, and anything living only on that box goes with it.
    """
    if not payload.get("confirm"):
        raise HTTPException(
            status_code=400,
            detail=(
                "Destroying the AWS substrate deletes the instance and its root volume. "
                'Re-send with {"confirm": true} to proceed.'
            ),
        )
    try:
        result = cloud_infra_module.destroy_infra(_cloud_infra_args(payload))
    except CloudCommandError as e:
        raise HTTPException(status_code=409, detail=str(e)) from e
    admin_activity_module.record("cloud_infra.destroy", result["settings"]["aws_region"])
    return result


# --- Terraform remote state endpoint (P6-9, #3510) ---
#
# The substrate's state starts as one local file, which stops being correct as
# soon as a second operator or a CI runner applies the same substrate.
# Migrating it to S3 with a DynamoDB lock, listing what is stored there,
# restoring a version and breaking a stuck lock are all `nyxgpt cloud state`
# subcommands and nothing else (owner decision, 2026-08-16, #3804): they act on
# the record of the substrate the dashboard itself may be running on, which is
# precisely the surface a UI must not offer. What remains here is the read the
# Infrastructure page's information panel needs.


@api.get("/cloud/state")
def cloud_state_status(_request: Request, verify: bool = False) -> dict[str, Any]:
    """Where the substrate's Terraform state lives, and how it is locked.

    Offline by default so the dashboard can poll it. `?verify=true`
    additionally confirms against AWS that the bucket and lock table exist and
    that versioning -- the whole recovery story -- is actually on.
    """
    try:
        return cloud_state_module.state_status(verify=verify)
    except CloudCommandError as e:
        raise HTTPException(status_code=409, detail=str(e)) from e


# --- Cloud deploy endpoints (P6-11, #3513) ---
#
# `GET` reports what release is on the instance, whether the access tunnel is
# open, and the deploy history -- the information the Infrastructure page's AWS
# section renders. Opening and closing that tunnel is `nyxgpt cloud tunnel`
# (owner decision, 2026-08-16, #3804).
#
# `deploy`/`destroy` remain as API operations for a client that is not the
# instance's own dashboard; no web surface posts to them. Both run Terraform
# and a remote install and take minutes, and are deliberately synchronous.


def _cloud_deploy_args(payload: dict[str, Any]) -> argparse.Namespace:
    """Build the namespace `nyxgpt.cloud_deploy` expects from a JSON body.

    A superset of `_cloud_infra_args` -- `deploy` applies the substrate first,
    so it takes every provisioning input as well as the deploy-specific ones.
    """
    namespace = _cloud_infra_args(payload)
    namespace.host = payload.get("host") or None
    namespace.ssh_user = payload.get("ssh_user") or None
    namespace.identity_file = payload.get("identity_file") or None
    namespace.version = payload.get("version") or None
    namespace.skip_observability = bool(payload.get("skip_observability"))
    namespace.no_tunnel = bool(payload.get("no_tunnel"))
    namespace.health_timeout = payload.get("health_timeout") or None
    namespace.ssh_timeout = payload.get("ssh_timeout") or None
    return namespace


@api.get("/cloud/deploy")
def cloud_deploy_status(_request: Request, probe_health: bool = False) -> dict[str, Any]:
    """What is deployed, at what version, whether the tunnel is open, and its history.

    Side-effect free by default -- reads the recorded deploy/tunnel state and
    the lifecycle history rather than calling AWS or the instance -- so the
    dashboard can poll it. `probe_health=true` adds one short request to the
    tunneled API health endpoint, which is what the Cloud Deployment page asks
    for on an explicit load or refresh.
    """
    return cloud_deploy_module.deploy_status(probe_health=probe_health)


@api.post("/cloud/deploy")
def cloud_deploy_run(payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Provision AWS and deploy the full stack onto it. Idempotent."""
    try:
        result = cloud_deploy_module.deploy(_cloud_deploy_args(payload))
    except CloudCommandError as e:
        raise HTTPException(status_code=409, detail=str(e)) from e
    admin_activity_module.record(
        "cloud_deploy.deploy",
        f"version={result['plan']['version']} instance={result['target']['instance_id'] or 'unknown'}",
    )
    return result


@api.post("/cloud/deploy/destroy")
def cloud_deploy_destroy(payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Close the tunnel and tear the whole cloud deployment down.

    Requires `{"confirm": true}` for the same reason the substrate teardown
    does: the instance and its root volume go with it.
    """
    if not payload.get("confirm"):
        raise HTTPException(
            status_code=400,
            detail=(
                "Destroying the cloud deployment deletes the instance and its root volume. "
                'Re-send with {"confirm": true} to proceed.'
            ),
        )
    try:
        result = cloud_deploy_module.destroy(_cloud_deploy_args(payload))
    except CloudCommandError as e:
        raise HTTPException(status_code=409, detail=str(e)) from e
    admin_activity_module.record("cloud_deploy.destroy", result["settings"].get("aws_region", ""))
    return result


# --- Containerized cloud artifact-install smoke endpoints (#3784) ---
#
# The SRE-surface half of `nyxgpt cloud smoke --container`: start a run of the
# artifact install path on a bare Amazon Linux 2023 container, and read the
# last run's verdict. Both surfaces drive the same
# `nyxgpt.cloud_artifact_smoke` functions and read the same recorded result,
# so there is one smoke rather than two implementations of it. This one stays
# startable from the dashboard under the #3804 rule: it exercises a throwaway
# local container, so it changes nothing the dashboard is running on.
#
# Unlike the AWS cloud endpoints above, this one is *not* synchronous: a run
# builds an image, boots systemd, installs Node, Docker, Ollama and the wheel,
# and waits for services -- tens of minutes, far past any HTTP timeout. The
# POST starts a background run and returns immediately; the panel polls the
# GET. This costs nothing and touches no AWS account, which is why it is a
# button at all where `cloud deploy` is a CLI pointer.


@api.get("/ops/cloud-artifact-smoke")
def ops_cloud_artifact_smoke_status(_request: Request) -> dict[str, Any]:
    """Report the last containerized artifact-install smoke, and whether one is running.

    Cheap enough to poll: it reads the recorded result and makes one local
    `docker info` call, which is what tells the panel whether a run is even
    possible on this machine.
    """
    return cloud_artifact_smoke_module.smoke_status()


@api.post("/ops/cloud-artifact-smoke")
def ops_cloud_artifact_smoke_run(payload: dict[str, Any] = Body(default={})) -> dict[str, Any]:
    """Start a containerized artifact-install smoke in the background.

    `{"version": "3.0.0rc9"}` pins the release the bootstrap installs;
    `{"inject": ["old-python"]}` runs the fault-injected variant, which passes
    only if the smoke fails. A run already in flight is a 409 rather than a
    second container fighting the first for the same name.
    """
    namespace = argparse.Namespace(
        container=True,
        version=payload.get("version") or None,
        wheel=payload.get("wheel") or None,
        image=payload.get("image") or None,
        inject=[str(f) for f in (payload.get("inject") or [])],
        keep=bool(payload.get("keep")),
        bootstrap_timeout=payload.get("bootstrap_timeout") or None,
        build_timeout=payload.get("build_timeout") or None,
        health_timeout=payload.get("health_timeout") or None,
        json=True,
        status=False,
    )
    available, detail = cloud_artifact_smoke_module.docker_available()
    if not available:
        raise HTTPException(
            status_code=409,
            detail=(
                f"A Docker engine is required to run the containerized cloud smoke -- {detail}."
            ),
        )
    try:
        result = cloud_artifact_smoke_module.start_background_run(namespace)
    except CloudCommandError as e:
        raise HTTPException(status_code=409, detail=str(e)) from e
    admin_activity_module.record(
        "cloud_artifact_smoke.run",
        f"version={result.get('version', 'latest')} inject={','.join(namespace.inject) or 'none'}",
    )
    return result


# --- Portability matrix endpoint (P6-16, #3516) ---
#
# The machine-readable half of `nyxgpt ops portability`: which deployment
# targets install and operate without a repo checkout, what evidence backs
# each one, what gaps are still open, and the clean-machine acceptance
# sequence.
#
# Read-only by construction. The matrix is a property of the *product* (which
# artifacts are published, which commands are wrapped), not of this machine's
# state, so there is nothing here to act on. That is also why it has no
# dashboard screen (#3803, which removed the one #3516 added): a page that
# restates the product's own documentation is not an ops surface. `commands`
# in the payload carries the wrapped CLI invocations a consumer points at.


@api.get("/ops/portability")
def ops_portability(_request: Request) -> dict[str, Any]:
    """Report the repo-less portability matrix and the capstone acceptance sequence.

    Side-effect free and cheap (no subprocesses, no network, a handful of
    `Path.exists` calls against the checkout when there is one), so a caller
    can poll it freely.
    """
    return portability_module.check_matrix()


# --- Release publish endpoint (#3727) ---
#
# The SRE-surface half of `nyxgpt release publish`: which release line the
# tip is on, which dev builds and RCs PyPI already serves, what the channel
# would publish next, and the exact pinned commands that point an acceptance
# install at it.
#
# Read-only, following the same #3514 status-plus-CLI-pointers decision the
# portability surface uses -- and here for a second reason: publishing to
# PyPI is an owner action carrying repo credentials, so it belongs to the
# dispatch-only workflow and the owner's `nyxgpt release publish --publish`,
# never to a button in a browser session.


@api.get("/ops/release-candidate")
def ops_release_candidate(
    _request: Request, branch: str | None = None, channel: str | None = None
) -> dict[str, Any]:
    """Report the publish plan for `branch` (default: the configured release branch).

    `channel` selects rc (default) or stable -- the same two channels the
    single publish pipeline understands (#3735 retired the nightly `dev`
    channel).

    Makes one outbound call, to PyPI's JSON API, to learn which versions
    already exist. A failed lookup is reported in `pypi_lookup_error` and
    clears `publishable` rather than failing the request -- the dashboard
    should still render the line's state when PyPI is unreachable.
    """
    target = (branch or release_candidate_module.default_branch()).strip()
    requested_channel = (channel or "rc").strip().lower()
    if requested_channel not in release_candidate_module.CHANNELS:
        raise HTTPException(
            status_code=400,
            detail=(
                f"Unknown channel '{requested_channel}' -- expected one of "
                + ", ".join(release_candidate_module.CHANNELS)
            ),
        )
    try:
        return release_candidate_module.plan(target, requested_channel)
    except release_candidate_module.ReleaseCandidateError as e:
        raise HTTPException(status_code=502, detail=str(e)) from e


# --- Support endpoints (#3745) ---
#
# The web UI's Support menu, which is the only documentation surface a user
# who installed from PyPI or Homebrew has: they never checked the repo out,
# so the product docs ship in the wheel and are served from the package here.
# Only product documentation ships -- the agent/CI process and contributor
# docs are not in the artifact at all (#3809, `support.DOC_SECTIONS`).
#
# All three are read-only. "File an Issue" is deliberately a *link* the UI
# opens, not an endpoint that posts on the user's behalf -- see
# `nyxgpt.support`.


@api.get("/support/docs")
def support_docs_index(_request: Request) -> dict[str, Any]:
    """List the packaged documentation, grouped and flat.

    `{sections: [{title, documents: [...]}], documents: [{slug, title,
    summary}]}` -- the viewer renders `sections` (#3809); `documents` is the
    same set flattened, in the same order.
    """
    sections = support_module.list_sections()
    return {
        "sections": sections,
        "documents": [doc for section in sections for doc in section["documents"]],
    }


@api.get("/support/docs/{slug}")
def support_document(_request: Request, slug: str) -> dict[str, str]:
    """Render one packaged document to HTML for the Support -> Docs viewer.

    A slug naming no packaged document is a 404, not a 500: it is ordinary
    user input from the URL bar, and `nyxgpt.support` rejects anything that
    could reach outside the packaged docs directory.
    """
    try:
        return support_module.render_document(slug)
    except support_module.DocumentNotFoundError as e:
        raise HTTPException(status_code=404, detail=str(e)) from e


@api.get("/support/context")
def support_context(_request: Request) -> dict[str, Any]:
    """Report the running environment and the prefilled issue-form link."""
    return support_module.support_context()


# --- Model management endpoints ---
@api.get("/models")
def models_list(request: Request) -> dict[str, Any]:
    """List model names currently available in Ollama.

    Returns `{"models": [name, ...]}`. Raises a `502` if the Ollama
    `/api/tags` call fails (e.g. Ollama unreachable).
    """
    cfg = _req_cfg(request)
    try:
        data = get_json(_ollama_url(cfg, "/api/tags"), timeout_s=10.0)
        models = data.get("models", []) if isinstance(data, dict) else []
        # Normalize to a list of model names
        names: list[str] = []
        for m in models:
            if isinstance(m, dict) and isinstance(m.get("name"), str):
                names.append(m["name"])
        return {"models": names}
    except Exception as e:
        raise HTTPException(
            status_code=502, detail=f"Failed to list models from Ollama: {e}"
        ) from e


@api.post("/models/pull")
def models_pull(request: Request, payload: dict[str, Any] = Body(...)) -> Response:
    """Pull a model from Ollama.

    When ``stream=true`` in the request body, returns a Server-Sent Events (SSE)
    stream with progress events (``status``, ``completed``, ``total``, ``percent``).
    Each event is a JSON object sent as ``data: {...}\\n\\n``.

    When ``stream`` is omitted or ``false``, returns a single JSON response on
    completion.
    """
    import json as _json

    cfg = _req_cfg(request)
    model = payload.get("model")
    if not isinstance(model, str) or not model.strip():
        raise HTTPException(status_code=400, detail="Missing 'model'")
    model = model.strip()
    stream = bool(payload.get("stream", False))

    if stream:
        base_url = get_ollama_base_url(cfg)

        def _progress_generator():
            from nyxgpt.ollama_client import post_json_lines

            url = base_url.rstrip("/") + "/api/pull"
            pull_payload = {"name": model, "stream": True}
            try:
                for event in post_json_lines(url, pull_payload, timeout_s=600.0):
                    status = event.get("status", "")
                    completed = event.get("completed", 0)
                    total = event.get("total", 0)
                    percent = round(completed / total * 100, 1) if total > 0 else 0.0
                    data = {
                        "status": status,
                        "completed": completed,
                        "total": total,
                        "percent": percent,
                    }
                    yield f"data: {_json.dumps(data)}\n\n"
                admin_activity_module.record("model.pull", model)
                yield f"data: {_json.dumps({'status': 'success', 'ok': True, 'model': model})}\n\n"
            except Exception as exc:
                # Log the detail server-side; return a generic message so raw
                # exception text isn't exposed to the client (CodeQL #26).
                log.exception("model.pull failed")
                del exc
                yield f"data: {_json.dumps({'error': 'Model pull failed'})}\n\n"

        return StreamingResponse(
            _progress_generator(),
            media_type="text/event-stream",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
        )

    try:
        # Non-streaming pull; Ollama may take a while.
        data = post_json(
            _ollama_url(cfg, "/api/pull"),
            {"name": model, "stream": False},
            timeout_s=600.0,
        )
        admin_activity_module.record("model.pull", model)
        return JSONResponse({"ok": True, "model": model, "result": data})
    except Exception as e:
        raise HTTPException(status_code=502, detail=f"Failed to pull model via Ollama: {e}") from e


@api.delete("/models/{model_name}")
def models_delete(request: Request, model_name: str) -> dict[str, Any]:
    """Delete a model from Ollama."""
    cfg = _req_cfg(request)
    try:
        models.delete_model(model_name, base_url=get_ollama_base_url(cfg))
        admin_activity_module.record("model.delete", model_name)
        return {"ok": True, "model": model_name}
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e)) from e
    except Exception as e:
        raise HTTPException(
            status_code=502, detail=f"Failed to delete model via Ollama: {e}"
        ) from e


@api.get("/models/{model_name}/info")
def models_info(request: Request, model_name: str) -> dict[str, Any]:
    """Get detailed information about a model."""
    cfg = _req_cfg(request)
    try:
        info = models.show_model(model_name, base_url=get_ollama_base_url(cfg))
        return {"ok": True, "model": model_name, "info": info}
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e)) from e
    except Exception as e:
        raise HTTPException(
            status_code=502, detail=f"Failed to get model info via Ollama: {e}"
        ) from e


@api.get("/sessions", response_model=SessionsListResponse)
def sessions_list(sessions_dir: str | None = None) -> SessionsListResponse:
    """List all chat sessions with summary metadata (message count, pinned, tags, title, model, etc.)."""
    cfg = _cfg(None)
    effective_dir = _sessions_dir_from_str(sessions_dir) or get_sessions_dir(cfg)
    rows = sessions.list_sessions(effective_dir)
    out: list[dict[str, Any]] = []
    for r in rows:
        meta = r.get("meta") or {}
        out.append(
            {
                "name": r.get("name"),
                "messages": r.get("messages"),
                "modified": r.get("modified"),
                "pinned": bool(meta.get("pinned")),
                "tags": meta.get("tags") if isinstance(meta.get("tags"), list) else [],
                "title": meta.get("title") if isinstance(meta.get("title"), str) else "",
                "summary": meta.get("summary") if isinstance(meta.get("summary"), str) else "",
                "token_estimate": (
                    meta.get("token_estimate")
                    if isinstance(meta.get("token_estimate"), int)
                    else None
                ),
                "model": meta.get("model") if isinstance(meta.get("model"), str) else "",
            }
        )
    return SessionsListResponse(sessions=out)


def search_params(
    query: str = Query(..., min_length=1, description="Text to search for in messages"),
    case_sensitive: bool = Query(False, description="Whether to perform case-sensitive search"),
    role_filter: api_models.MessageRole | None = Query(
        None, description="Filter by message role (user, assistant, system)"
    ),
    session_filter: str | None = Query(None, description="Filter to specific session name"),
    limit: int = Query(50, ge=1, le=500, description="Maximum number of results to return"),
) -> api_models.SearchRequest:
    """Dependency function to validate search parameters using SearchRequest model.

    This ensures proper validation of all search parameters according to the
    SearchRequest Pydantic model definition, including enum validation for role_filter.
    """
    return api_models.SearchRequest(
        query=query,
        case_sensitive=case_sensitive,
        role_filter=role_filter,
        session_filter=session_filter,
        limit=limit,
    )


@api.get("/sessions/search", response_model=api_models.SearchResponse)
def sessions_search(
    params: api_models.SearchRequest = Depends(search_params),
    sessions_dir: str | None = None,
) -> api_models.SearchResponse:
    """Search for messages across all sessions or within a specific session.

    Args:
        params: Validated search parameters (query, filters, limit)
        sessions_dir: Optional sessions directory override

    Returns:
        SearchResponse containing matching messages with context
    """
    cfg = _cfg(None)
    effective_dir = _sessions_dir_from_str(sessions_dir) or get_sessions_dir(cfg)

    # Perform search
    # Convert enum to string for sessions module
    role_filter_str = params.role_filter.value if params.role_filter else None
    results = sessions.search_messages(
        query=params.query,
        sessions_dir=effective_dir,
        case_sensitive=params.case_sensitive,
        role_filter=role_filter_str,
        session_filter=params.session_filter,
        limit=params.limit,
    )

    # Convert to API models
    result_items = [
        api_models.SearchResultItem(
            session_name=r["session_name"],
            session_title=r.get("session_title"),
            message_index=r["message_index"],
            role=r["role"],
            content=r["content"],
            content_preview=r["content_preview"],
            timestamp=r.get("timestamp"),
            matches=r["matches"],
        )
        for r in results
    ]

    return api_models.SearchResponse(
        query=params.query,
        total_results=len(result_items),
        results=result_items,
    )


@api.get("/sessions/{name}")
def sessions_show(
    name: str,
    sessions_dir: str | None = None,
    offset: int | None = None,
    limit: int | None = None,
) -> dict[str, Any]:
    """Return a session's messages and metadata, optionally paginated.

    When `offset`/`limit` are provided, messages are loaded page-by-page
    rather than reading the whole session file into memory. Raises `404`
    if the session doesn't exist, or `400` for a negative `offset`/`limit`.
    """
    sd = _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))
    sf = sessions.session_file_for(name, sd or sessions.default_sessions_dir())
    mf = sessions.meta_file_for(sf)
    if not sessions.session_file_exists(sf):
        raise HTTPException(status_code=404, detail="No such session")

    # Validate pagination parameters (Medium Issue 4)
    if offset is not None and offset < 0:
        raise HTTPException(status_code=400, detail="offset must be non-negative")
    if limit is not None and limit < 0:
        raise HTTPException(status_code=400, detail="limit must be non-negative")

    # Use paginated loading to avoid loading all messages into memory (Critical Issue 1)
    if offset is not None or limit is not None:
        start = offset if offset is not None else 0
        msgs, total_count = sessions.load_session_messages_paginated(sf, start, limit)
    else:
        # No pagination requested - use regular load for backward compatibility
        all_msgs = sessions.load_session_messages(sf)
        msgs = all_msgs
        total_count = len(all_msgs)

    meta = sessions.load_session_meta(mf)
    return {
        "name": name,
        "messages": msgs,
        "meta": meta,
        "total": total_count,
        "offset": offset if offset is not None else 0,
        "limit": limit if limit is not None else total_count,
    }


# Lightweight session initialization endpoint (does NOT invoke the model)


@api.post("/sessions/init")
def sessions_init(req: dict[str, Any] = Body(...)) -> dict[str, Any]:
    """Create a session file (and its metadata) without invoking the model.

    Idempotent: if a session with `name` already exists, returns
    `{"ok": True, "existed": True}` without modifying it. `system`/`model`
    fall back to config defaults when omitted. Raises `422` for an invalid
    session `name`, `400` for a missing `name` or if session initialization
    otherwise fails, and `500` on an unexpected internal error resolving the
    session path.
    """
    name = req.get("name")
    if not name or not isinstance(name, str):
        raise HTTPException(status_code=400, detail="Session name is required")

    cfg = _cfg(None)
    sd = get_sessions_dir(cfg)

    # Ensure sessions directory exists
    sd.mkdir(parents=True, exist_ok=True)

    # Idempotent behavior: if the session already exists, succeed
    try:
        sf = sessions.session_file_for(name, sd)
    except ValueError as e:
        # Validation errors should return 422 (Unprocessable Entity)
        log.warning("Invalid session name for init: %s", e)
        raise HTTPException(status_code=422, detail=str(e)) from e
    except Exception as e:
        # Other errors are internal server errors
        log.error("Failed to get session file path: %s", e)
        raise HTTPException(status_code=500, detail="Internal server error") from e

    if sessions.session_file_exists(sf):
        return {"ok": True, "name": name, "existed": True}

    system = req.get("system")
    if not isinstance(system, str) or not system:
        # Use config default (empty string if not set)
        system = cfg.get("nyxgpt", "system_prompt", fallback="")

    model = req.get("model")
    if not isinstance(model, str) or not model:
        model = get_default_model(cfg)

    try:
        _sf, _mf, _msgs, _meta = sessions.init_session(
            name,
            sd,
            new_session=True,
            model=model,
            system=system,
        )
    except Exception as e:
        log.exception("sessions.init_session failed")
        raise HTTPException(status_code=400, detail=str(e)) from e

    return {"ok": True, "name": name, "existed": False}


@api.delete("/sessions/{name}")
def sessions_delete(name: str, sessions_dir: str | None = None) -> dict[str, Any]:
    """Delete a session's file and its metadata. Raises `404` if the session doesn't exist."""
    ok = sessions.delete_session(
        name, _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))
    )
    if not ok:
        raise HTTPException(status_code=404, detail="No such session")
    return {"ok": True}


@api.post("/sessions/{name}/summarize")
def sessions_summarize(name: str, sessions_dir: str | None = None) -> dict[str, Any]:
    """Generate and store a summary in the session's metadata. Raises `400` on failure (e.g. no such session)."""
    ok, msg = sessions.summarize_session(
        name, _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))
    )
    if not ok:
        raise HTTPException(status_code=400, detail=msg)
    return {"ok": True}


@api.post("/sessions/{name}/pin")
def sessions_pin(name: str, sessions_dir: str | None = None) -> dict[str, Any]:
    """Mark a session as pinned. Raises `400` on failure (e.g. no such session)."""
    ok, msg = sessions.set_pinned(
        name, True, _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))
    )
    if not ok:
        raise HTTPException(status_code=400, detail=msg)
    return {"ok": True}


@api.post("/sessions/{name}/unpin")
def sessions_unpin(name: str, sessions_dir: str | None = None) -> dict[str, Any]:
    """Remove a session's pinned flag. Raises `400` on failure (e.g. no such session)."""
    ok, msg = sessions.set_pinned(
        name,
        False,
        _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None)),
    )
    if not ok:
        raise HTTPException(status_code=400, detail=msg)
    return {"ok": True}


@api.post("/sessions/{name}/title")
def sessions_title(name: str, req: TitleRequest, sessions_dir: str | None = None) -> dict[str, Any]:
    """Set a session's display title (stored in metadata; does not rename the underlying file).

    Raises `400` on failure (e.g. no such session).
    """
    ok, msg = sessions.set_title(
        name,
        req.title,
        _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None)),
    )
    if not ok:
        raise HTTPException(status_code=400, detail=msg)
    return {"ok": True}


@api.post("/sessions/{name}/tags/add")
def sessions_tags_add(
    name: str, req: TagsRequest, sessions_dir: str | None = None
) -> dict[str, Any]:
    """Add one or more tags to a session's metadata.

    Raises `400` if `req.tags` is empty or the underlying update fails
    (e.g. no such session).
    """
    if not req.tags:
        raise HTTPException(status_code=400, detail="At least one tag is required")
    ok, msg = sessions.add_tags(
        name,
        req.tags,
        _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None)),
    )
    if not ok:
        raise HTTPException(status_code=400, detail=msg)
    return {"ok": True}


@api.post("/sessions/{name}/tags/remove")
def sessions_tags_remove(
    name: str, req: TagsRequest, sessions_dir: str | None = None
) -> dict[str, Any]:
    """Remove one or more tags from a session's metadata.

    Raises `400` if `req.tags` is empty or the underlying update fails
    (e.g. no such session).
    """
    if not req.tags:
        raise HTTPException(status_code=400, detail="At least one tag is required")
    ok, msg = sessions.remove_tags(
        name,
        req.tags,
        _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None)),
    )
    if not ok:
        raise HTTPException(status_code=400, detail=msg)
    return {"ok": True}


@api.post("/sessions/{name}/rename")
def sessions_rename(
    name: str, req: RenameRequest, sessions_dir: str | None = None
) -> dict[str, Any]:
    """Rename a session with optional title update and filename sync.

    This endpoint supports two modes:
    1. Direct rename: Provide a valid session name to rename files directly
    2. Title-based rename: Provide a title, which will be sanitized for use as filename

    If sync_filename=True (default), the session title will be updated and the
    filename will be synced to match the sanitized title.
    """
    _sessions_dir = _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))

    # Check if current session exists
    sf = sessions.session_file_for(name, _sessions_dir)
    if not sessions.session_file_exists(sf):
        raise HTTPException(status_code=404, detail=f"Session '{name}' not found")

    if req.sync_filename:
        # Mode 1: Update title and sync filename
        # Set the title first
        ok, msg = sessions.set_title(name, req.new_name, _sessions_dir)
        if not ok:
            raise HTTPException(status_code=400, detail=msg)

        # Then sync filename based on title
        success, status, new_name = sessions.sync_filename_with_title(
            name, _sessions_dir, force=True
        )
        if not success:
            raise HTTPException(
                status_code=500,
                detail=f"Title updated but filename sync failed: {status}",
            )

        return {
            "ok": True,
            "old_name": name,
            "new_name": new_name,
            "message": "Session renamed and filename synced",
        }
    else:
        # Mode 2: Direct rename (validate new name first)
        try:
            validated_name = sessions.validate_session_name(req.new_name)
        except ValueError as e:
            raise HTTPException(status_code=400, detail=str(e)) from e

        # Use existing rename function
        ok, msg = sessions.rename_session(name, validated_name, _sessions_dir)
        if not ok:
            raise HTTPException(status_code=400, detail=msg)

        return {
            "ok": True,
            "old_name": name,
            "new_name": validated_name,
            "message": "Session renamed",
        }


@api.post("/sessions/{name}/sync-filename")
def sessions_sync_filename(name: str, sessions_dir: str | None = None) -> dict[str, Any]:
    """Force filename sync for a session based on its current title.

    This endpoint is useful when:
    - A session was created before auto-sync was enabled
    - Manual title changes were made without filename sync
    - You want to clean up session filenames to match their titles
    """
    _sessions_dir = _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))

    # Check if session exists
    sf = sessions.session_file_for(name, _sessions_dir)
    if not sessions.session_file_exists(sf):
        raise HTTPException(status_code=404, detail=f"Session '{name}' not found")

    # Force filename sync
    success, status, new_name = sessions.sync_filename_with_title(name, _sessions_dir, force=True)

    if not success:
        raise HTTPException(status_code=500, detail=f"Filename sync failed: {status}")

    responses = {
        "no_title": {"ok": True, "message": "No title set, filename unchanged", "name": name},
        "no_change": {"ok": True, "message": "Filename already matches title", "name": name},
        "renamed": {
            "ok": True,
            "old_name": name,
            "new_name": new_name,
            "message": "Filename synced with title",
        },
    }
    return responses.get(status, {"ok": True, "message": status, "name": new_name})


@api.patch("/sessions/{name}/messages/{message_index}")
def edit_message(
    name: str,
    message_index: int,
    req: api_models.EditMessageRequest,
    sessions_dir: str | None = None,
) -> dict[str, Any]:
    """Edit a message in a session.

    By default, forks the conversation (truncates messages after the edited one).
    Set fork=false to edit without truncating.
    """
    _sessions_dir = _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))

    ok, msg = sessions.edit_message(
        session_name=name,
        message_index=message_index,
        new_content=req.content,
        sessions_dir=_sessions_dir,
        fork=req.fork,
    )

    if not ok:
        raise HTTPException(status_code=400, detail=msg)

    return {"ok": True, "message": msg}


@api.get("/sessions/{name}/messages/{message_index}/rag")
def get_message_rag_chunks(
    name: str, message_index: int, sessions_dir: str | None = None
) -> dict[str, Any]:
    """Get RAG chunks for a specific message.

    Returns the RAG chunks associated with a message if available.
    This endpoint enables lazy loading of RAG citation data.
    """
    _sessions_dir = _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))

    # Load session messages
    sf = sessions.session_file_for(name, _sessions_dir)
    if not sessions.session_file_exists(sf):
        raise HTTPException(status_code=404, detail=f"Session '{name}' not found")

    msgs = sessions.load_session_messages(sf)

    # Validate message index
    if message_index < 0 or message_index >= len(msgs):
        raise HTTPException(
            status_code=400,
            detail=f"Invalid message_index: {message_index} (session has {len(msgs)} messages)",
        )

    message = msgs[message_index]

    # Extract RAG chunks if present
    rag_chunks: list[Any] = cast(list[Any], message.get("rag_chunks", []))

    return {
        "message_index": message_index,
        "has_rag": len(rag_chunks) > 0,
        "chunks": rag_chunks,
    }


def _escape_markdown(text: Any) -> str:
    """Escape markdown special characters in text.

    Citation fields come from stored session data and may be missing
    (`None`, if a chunk was persisted without a `doc_id`/`text`) or, for
    malformed data, a non-string shape (e.g. a nested dict). Coerce to
    string instead of crashing with AttributeError on `.replace()`.
    """
    if text is None:
        return ""
    result: str = text if isinstance(text, str) else str(text)
    # Escape common markdown special characters that could break formatting
    replacements = {
        "\\": "\\\\",  # Backslash must be first
        "`": "\\`",
        "*": "\\*",
        "_": "\\_",
        "{": "\\{",
        "}": "\\}",
        "[": "\\[",
        "]": "\\]",
        "(": "\\(",
        ")": "\\)",
        "#": "\\#",
        "+": "\\+",
        "-": "\\-",
        ".": "\\.",
        "!": "\\!",
        "|": "\\|",
    }
    for char, escaped in replacements.items():
        result = result.replace(char, escaped)
    return result


@api.get("/sessions/{name}/citations/export", response_model=None)
def export_session_citations(
    name: str, format: str = "json", sessions_dir: str | None = None
) -> dict[str, Any] | Response:
    """Export all RAG citations from a session.

    Returns all RAG citations from assistant messages in the session.
    Useful for generating bibliographies or reference lists.

    Supported formats: json, markdown
    """
    # Validate session name to prevent path traversal attacks
    if not name or ".." in name or "/" in name or "\\" in name:
        raise HTTPException(
            status_code=400,
            detail="Invalid session name. Must not contain path separators or navigation characters.",
        )

    format_lower = format.lower()
    if format_lower not in ("json", "markdown"):
        raise HTTPException(
            status_code=400,
            detail="Invalid format. Must be one of: json, markdown",
        )

    _sessions_dir = _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))

    # Load session messages
    sf = sessions.session_file_for(name, _sessions_dir)
    if not sessions.session_file_exists(sf):
        raise HTTPException(status_code=404, detail=f"Session '{name}' not found")

    # Use file locking to prevent race conditions during read
    with sessions.session_lock(sf, timeout=5.0):
        msgs = sessions.load_session_messages(sf)

    # Extract all citations from assistant messages
    citations: list[dict[str, Any]] = []
    for msg_idx, msg in enumerate(msgs):
        if msg.get("role") == "assistant":
            rag_chunks: list[Any] = cast(list[Any], msg.get("rag_chunks", []))
            if rag_chunks and isinstance(rag_chunks, list):
                for chunk_idx, chunk in enumerate(rag_chunks):
                    citations.append(
                        {
                            "message_index": msg_idx,
                            "citation_index": chunk_idx,
                            "doc_id": chunk.get("doc_id"),
                            "chunk_id": chunk.get("chunk_id"),
                            "text": chunk.get("text"),
                            "score": chunk.get("score"),
                            "similarity_score": chunk.get("similarity_score"),
                            "collection": chunk.get("collection"),
                            "chunk_number": chunk.get("chunk_number"),
                            "total_chunks": chunk.get("total_chunks"),
                        }
                    )

    if format_lower == "json":
        return {
            "session": name,
            "total_citations": len(citations),
            "citations": citations,
        }
    else:  # markdown
        lines: list[str] = []
        lines.append(f"# Citations for {name}\n")
        lines.append(f"Total sources: {len(citations)}\n")
        lines.append("---\n")

        for idx, citation in enumerate(citations, 1):
            doc_id = _escape_markdown(citation.get("doc_id", "Unknown"))
            # Use explicit None checking to avoid treating 0.0 as falsy
            score = citation.get("similarity_score")
            if score is None:
                score = citation.get("score", 0.0)
            text = _escape_markdown(citation.get("text", ""))

            chunk_ref = sessions.format_chunk_ref(citation)
            collection = citation.get("collection")
            collection_suffix = f", {collection}" if collection and collection != "default" else ""
            lines.append(f"## [{idx}] {doc_id} ({chunk_ref}{collection_suffix})\n")
            lines.append(f"**Confidence:** {score:.3f}")
            lines.append(f"**Message:** {citation['message_index']}\n")

            if text:
                lines.append(f"**Source text:**\n> {text}\n")

        return Response(
            content="\n".join(lines),
            media_type="text/markdown; charset=utf-8",
            headers={"Content-Disposition": f'attachment; filename="{name}-citations.md"'},
        )


@api.post("/sessions/{name}/messages/{message_index}/regenerate")
def regenerate_response(
    name: str,
    message_index: int,
    req: api_models.RegenerateRequest,
    sessions_dir: str | None = None,
) -> dict[str, Any]:
    """Regenerate response from a specific message.

    The message at message_index should be a user message.
    This will:
    1. Optionally replace the user message with new prompt (if provided)
    2. Truncate conversation after that message
    3. Generate a new response using the chat endpoint

    Returns the new response.
    """
    _sessions_dir = _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))

    # Load session to validate message index
    sf = sessions.session_file_for(name, _sessions_dir)
    if not sessions.session_file_exists(sf):
        raise HTTPException(status_code=404, detail="No such session")

    msgs = sessions.load_session_messages(sf)
    if message_index < 0 or message_index >= len(msgs):
        raise HTTPException(status_code=400, detail=f"Invalid message index: {message_index}")

    message = msgs[message_index]
    if message.get("role") != "user":
        raise HTTPException(status_code=400, detail="Can only regenerate from user messages")

    # If new prompt provided, edit the message first
    prompt = req.prompt if req.prompt else message.get("content", "")
    if req.prompt:
        ok, msg = sessions.edit_message(
            session_name=name,
            message_index=message_index,
            new_content=req.prompt,
            sessions_dir=_sessions_dir,
            fork=True,  # Always fork when regenerating
        )
        if not ok:
            raise HTTPException(status_code=400, detail=msg)
    else:
        # Just truncate after this message
        ok, msg = sessions.truncate_after_message(
            session_name=name,
            message_index=message_index,
            sessions_dir=_sessions_dir,
        )
        if not ok:
            raise HTTPException(status_code=400, detail=msg)

    # Generate new response using existing chat endpoint logic
    try:
        result = chat_module.chat(
            prompt=prompt,
            session=name,
            new=False,
            model=req.model,
            config_path=None,
            sessions_dir=str(_sessions_dir) if _sessions_dir else None,
            rag_enabled=req.rag_enabled,
        )
        return {
            "ok": True,
            "session": result.session,
            "model": result.model,
            "reply": result.reply,
            "rag_used": result.rag_used,
        }
    except Exception as e:
        log.error(f"Failed to regenerate response: {e}")
        raise HTTPException(status_code=500, detail=f"Regeneration failed: {e}") from e


@api.get("/sessions/{name}/export")
def sessions_export(
    name: str, format: str = "markdown", sessions_dir: str | None = None
) -> Response:
    """Export session to markdown, JSON, or HTML format."""
    format_lower = format.lower()
    if format_lower not in ("markdown", "json", "html"):
        raise HTTPException(
            status_code=400,
            detail="Invalid format. Must be one of: markdown, json, html",
        )

    # Validate the session name at this same response boundary -- the raw
    # `name` path parameter must never reach the HTML branch's <title>/body
    # or the Content-Disposition filename unvalidated (CodeQL py/reflected-xss,
    # alert #14). `validate_session_name` only accepts
    # `^[a-zA-Z0-9_-]{1,64}$`, which also rules out header-injection
    # characters (quotes, CR/LF) in the filename.
    try:
        safe_name = sessions.validate_session_name(name)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e)) from e

    sd = _sessions_dir_from_str(sessions_dir) or get_sessions_dir(_cfg(None))

    # nosniff guards the markdown/json branches against content-sniffing and
    # is set explicitly here (rather than relying solely on the global
    # security-headers middleware) so it's visible on the export response
    # itself.
    headers = {"X-Content-Type-Options": "nosniff"}

    if format_lower == "markdown":
        ok, content = sessions.export_session_markdown(safe_name, sd)
        if not ok:
            raise HTTPException(status_code=404, detail=content)
        return Response(
            content=content,
            media_type="text/markdown; charset=utf-8",
            headers={
                "Content-Disposition": f'attachment; filename="{safe_name}.md"',
                **headers,
            },
        )

    if format_lower == "json":
        ok, content = sessions.export_session_json(safe_name, sd)
        if not ok:
            raise HTTPException(status_code=404, detail=content)
        return Response(
            content=content,
            media_type="application/json; charset=utf-8",
            headers={
                "Content-Disposition": f'attachment; filename="{safe_name}.json"',
                **headers,
            },
        )

    ok, content = sessions.export_session_html(safe_name, sd)
    if not ok:
        raise HTTPException(status_code=404, detail=content)
    return Response(
        content=content,
        media_type="text/html; charset=utf-8",
        headers={
            "Content-Disposition": f'attachment; filename="{safe_name}.html"',
            **headers,
        },
    )


@api.post("/chat", response_model=ChatResponse)
def chat(request: Request, req: ChatRequest) -> ChatResponse:
    """Send a prompt to the model and return the full (non-streaming) reply.

    Creates or appends to the named session, resolves the model and RAG
    settings from the request (falling back to config defaults), and
    routes through the batch processor when request batching is enabled.
    Records usage analytics and Prometheus chat/RAG metrics as a side
    effect. Returns the session name, model used, reply text, and any RAG
    chunks that were retrieved.

    Raises:
        HTTPException 422: Invalid request (e.g. bad session name).
        HTTPException 502: Underlying model runtime failure (Ollama crash/
            OOM/timeout).
        HTTPException 504: Timed out waiting for the batch processor.
        HTTPException 500: Any other unexpected failure.
    """
    req_id = getattr(request.state, "request_id", None)
    usage_start = time.monotonic()

    try:
        d = _chat_runtime_defaults(_req_cfg(request))
        chosen_model = req.model or d["default_model"]

        log.info(
            "Chat request received",
            extra={
                "request_id": req_id,
                "session": req.session,
                "model": chosen_model,
                "new_session": req.new,
                "prompt_length": len(req.prompt),
                "rag_enabled": d.get("rag_enabled", False),
                "batching_enabled": _batch_processor is not None,
            },
        )

        # Determine RAG settings
        rag_val = getattr(req, "rag_enabled", None)
        rag_enabled = d["rag_enabled"] if rag_val is None else bool(rag_val)
        rag_filters = getattr(req, "rag_filters", None)
        rag_filters_dict = rag_filters.model_dump() if rag_filters else None
        if rag_filters_dict and rag_filters_dict.get("collection"):
            rag_filters_dict["collection"] = _resolve_and_validate_collection(
                rag_filters_dict["collection"]
            )

        # If batching is enabled, route through batch processor
        if _batch_processor is not None:
            # Prepare request data for batch processing
            _resolved_sessions_dir = _sessions_dir_from_str(req.sessions_dir)
            batch_req_data = {
                "prompt": req.prompt,
                "session": req.session,
                "new": req.new,
                "model": chosen_model,
                "system": req.system,
                "config_path": None,
                "sessions_dir": (str(_resolved_sessions_dir) if _resolved_sessions_dir else None),
                "rag_enabled": rag_enabled,
                "rag_filters": rag_filters_dict,
            }

            # Determine priority - interactive requests get higher priority
            # Could be extended to check request headers or user preference
            priority = RequestPriority.INTERACTIVE

            # Submit to batch processor. The wait here must cover however long
            # the underlying ollama_chat() call itself is allowed to take
            # (chat_timeout_seconds) plus queueing/scheduling overhead --
            # otherwise a slow cold model load times out the batch submit
            # before the actual chat request has had a chance to finish.
            batch_submit_timeout = d["chat_timeout_seconds"] + 10
            try:
                result_dict = _batch_processor.submit(
                    batch_req_data, priority=priority, timeout=batch_submit_timeout
                )
            except TimeoutError as e:
                log.error(
                    "Batch processing timed out",
                    extra={"request_id": req_id, "error": str(e)},
                )
                raise HTTPException(
                    status_code=504,
                    detail=(
                        "Chat request timed out waiting for the batch processor "
                        f"after {batch_submit_timeout:.0f}s. The model may require "
                        "more memory than is available, or chat_timeout_seconds "
                        "may need to be increased."
                    ),
                ) from e

            # Check for errors in result
            if "error" in result_dict:
                error_msg = result_dict.get("error", "Unknown error")
                error_type = result_dict.get("error_type", "Exception")
                log.error(
                    "Batch processing failed",
                    extra={
                        "request_id": req_id,
                        "error": error_msg,
                        "error_type": error_type,
                    },
                )
                status_code = 502 if error_type == "ModelRuntimeError" else 500
                raise HTTPException(status_code=status_code, detail=error_msg)

            # Convert result dict to ChatResult-like structure
            session_name = result_dict["session"]
            model_name = result_dict["model"]
            reply_text = result_dict["reply"]
            rag_used = result_dict["rag_used"]
            rag_context_data = result_dict.get("rag_context")

        else:
            # No batching - process directly
            _resolved_sessions_dir = _sessions_dir_from_str(req.sessions_dir)
            kwargs: dict[str, Any] = {
                "session": req.session,
                "new": req.new,
                "model": chosen_model,
                "system": req.system,
                "config_path": None,
                "sessions_dir": (str(_resolved_sessions_dir) if _resolved_sessions_dir else None),
            }

            # Optional runtime override: only pass if chat implementation supports it.
            if _maybe_kw(run_chat, "rag_enabled"):
                kwargs["rag_enabled"] = rag_enabled

            if _maybe_kw(run_chat, "rag_filters"):
                kwargs["rag_filters"] = rag_filters_dict

            if _maybe_kw(run_chat, "attachments"):
                attachments_val = getattr(req, "attachments", None)
                if attachments_val:
                    kwargs["attachments"] = [a.model_dump() for a in attachments_val]

            if _maybe_kw(run_chat, "output_format"):
                output_format_val = getattr(req, "output_format", None)
                if output_format_val is not None:
                    kwargs["output_format"] = output_format_val

            result = run_chat(req.prompt, **kwargs)
            session_name = result.session
            model_name = result.model
            reply_text = result.reply
            rag_used = result.rag_used
            rag_context_data = result.rag_context

        log.info(
            "Chat request completed",
            extra={
                "request_id": req_id,
                "session": session_name,
                "model": model_name,
                "reply_length": len(reply_text),
            },
        )

        prom_metrics.CHAT_REQUESTS_TOTAL.labels(model=model_name, streaming="false").inc()
        if rag_used:
            prom_metrics.RAG_QUERIES_TOTAL.labels(source="chat").inc()

        try:
            usage_analytics_module.record(
                session=session_name,
                model=model_name,
                prompt_tokens=_count_usage_tokens(req.prompt),
                completion_tokens=_count_usage_tokens(reply_text),
                duration_s=time.monotonic() - usage_start,
                cfg=_req_cfg(request),
            )
        except Exception:
            log.debug("Usage analytics recording failed", exc_info=True)

        # Convert RAG context to RagChunkInfo objects
        rag_chunks = []
        if rag_context_data:
            for chunk_data in rag_context_data:
                rag_chunks.append(
                    RagChunkInfo(
                        text=chunk_data.get("text", ""),
                        score=chunk_data.get("score", 0.0),
                        doc_id=chunk_data.get("doc_id"),
                        chunk_id=chunk_data.get("chunk_id"),
                        similarity_score=chunk_data.get("similarity_score"),
                        collection=chunk_data.get("collection"),
                        chunk_number=chunk_data.get("chunk_number"),
                        total_chunks=chunk_data.get("total_chunks"),
                    )
                )

        return ChatResponse(
            session=session_name,
            model=model_name,
            reply=reply_text,
            rag_used=rag_used,
            rag_chunks=rag_chunks,
        )
    except ValueError as e:
        # Validation errors (e.g., invalid session name)
        log.warning("Chat validation error", extra={"request_id": req_id, "error": str(e)})
        raise HTTPException(status_code=422, detail=str(e)) from e
    except ModelRuntimeError as e:
        # Upstream Ollama/model-runtime failure (crash, OOM, timeout) -- give
        # the operator/user the actionable detail instead of a bare 500.
        log.error(
            "Chat request failed: model runtime error",
            extra={
                "request_id": req_id,
                "error": str(e),
                "error_type": type(e).__name__,
            },
            exc_info=True,
        )
        raise HTTPException(status_code=502, detail=str(e)) from e
    except HTTPException:
        # Already a well-formed HTTP error (e.g. the batch-processing branch
        # above) -- re-raise as-is instead of letting it fall into the
        # catch-all below, which would otherwise flatten it into a generic
        # 500 "Internal server error" and discard its real status/detail.
        raise
    except Exception as e:
        log.error(
            "Chat request failed",
            extra={
                "request_id": req_id,
                "error": str(e),
                "error_type": type(e).__name__,
            },
            exc_info=True,
        )
        raise HTTPException(status_code=500, detail="Internal server error") from e


def _create_streaming_response(request: Request, req: ChatRequest) -> StreamingResponse:
    """Create a streaming chat response with request ID context propagation.

    This helper consolidates the streaming logic used by both versioned and legacy endpoints.
    It handles request ID capture and context setting for proper log traceability.

    Now supports Server-Sent Events (SSE) framing with proper event types, IDs, and heartbeats.
    Adapts response format based on client capability hints for backwards compatibility.

    Args:
        request: FastAPI Request object containing state and configuration
        req: Chat request parameters

    Returns:
        StreamingResponse configured for text/event-stream (SSE) streaming or plain text

    Raises:
        HTTPException: 422 for validation errors, 500 for server errors
    """
    try:
        # Validate session name early, before creating generator
        # This ensures validation errors are caught and return 422 instead of failing during streaming
        from nyxgpt.sessions import validate_session_name

        try:
            validate_session_name(req.session)
        except ValueError as e:
            raise HTTPException(status_code=422, detail=str(e)) from e

        # Validate the requested RAG collection (if any) before streaming starts,
        # so an unknown collection surfaces as a clean 400 instead of a mid-stream error.
        if req.rag_filters and req.rag_filters.collection:
            req.rag_filters.collection = _resolve_and_validate_collection(
                req.rag_filters.collection
            )

        # Parse client capabilities for content negotiation
        capabilities = _parse_client_capabilities(request)

        # Log client capabilities for debugging
        log.debug(
            "Client capabilities detected",
            extra={
                "supports_sse": capabilities.supports_sse,
                "supports_structured_events": capabilities.supports_structured_events,
                "client_version": capabilities.client_version,
            },
        )

        # Capture request ID before entering generator (context may not propagate)
        req_id = request.state.request_id

        def _stream_sse():
            """Generate SSE-formatted events from chat stream.

            Adapts output format based on client capabilities:
            - Legacy clients (no SSE): Plain text streaming
            - SSE clients (no structured events): Simple SSE text events
            - Modern clients: Full structured SSE events (metadata, text, done, error)
            """
            # Explicitly set request ID in context for streaming generator
            try:
                request_id_var.set(req_id)
            except Exception as e:
                log.warning(f"Failed to set request ID in streaming context: {e}")
            # Continue regardless - streaming should work even if request ID fails

            import json
            import time

            event_id = 0
            start_time = time.time()
            total_tokens = 0
            rag_used_in_stream = False

            d = _chat_runtime_defaults(_req_cfg(request))
            chosen_model = req.model or d["default_model"]
            prom_metrics.CHAT_REQUESTS_TOTAL.labels(model=chosen_model, streaming="true").inc()

            # Only send structured events if client supports them
            if capabilities.supports_sse and capabilities.supports_structured_events:
                # Send an immediate heartbeat to establish connection
                event_id += 1
                yield f"event: heartbeat\ndata: {json.dumps({'timestamp': start_time})}\nid: {event_id}\n\n"

                # Send metadata event with session/model info
                event_id += 1
                metadata = {
                    "session": req.session,
                    "model": chosen_model,
                    "timestamp": start_time,
                }
                yield f"event: metadata\ndata: {json.dumps(metadata)}\nid: {event_id}\n\n"

            _resolved_sessions_dir = _sessions_dir_from_str(req.sessions_dir)
            kwargs: dict[str, Any] = {
                "session": req.session,
                "new": req.new,
                "model": chosen_model,
                "system": req.system,
                "config_path": None,
                "sessions_dir": (str(_resolved_sessions_dir) if _resolved_sessions_dir else None),
            }

            if _maybe_kw(chat_stream, "rag_enabled"):
                rag_val = getattr(req, "rag_enabled", None)
                kwargs["rag_enabled"] = d["rag_enabled"] if rag_val is None else bool(rag_val)

            if _maybe_kw(chat_stream, "rag_filters"):
                rag_filters_val = getattr(req, "rag_filters", None)
                if rag_filters_val:
                    # Convert RagFilters model to dict
                    kwargs["rag_filters"] = rag_filters_val.model_dump()

            if _maybe_kw(chat_stream, "attachments"):
                attachments_val = getattr(req, "attachments", None)
                if attachments_val:
                    # Convert AttachmentBlock models to dicts
                    kwargs["attachments"] = [a.model_dump() for a in attachments_val]

            if _maybe_kw(chat_stream, "output_format"):
                output_format_val = getattr(req, "output_format", None)
                if output_format_val is not None:
                    kwargs["output_format"] = output_format_val

            # Process the chat stream and wrap chunks in appropriate format
            try:
                for chunk in chat_stream(req.prompt, **kwargs):
                    # Strip RAG/retry markers for legacy clients
                    clean_chunk = chunk
                    if "__RAG_START__" in chunk and "__RAG_END__" in chunk:
                        rag_used_in_stream = True
                        rag_start = chunk.find("__RAG_START__")
                        rag_end = chunk.find("__RAG_END__") + len("__RAG_END__")
                        clean_chunk = chunk[:rag_start] + chunk[rag_end:]
                    if "__RETRY_START__" in clean_chunk and "__RETRY_END__" in clean_chunk:
                        retry_start = clean_chunk.find("__RETRY_START__")
                        retry_end = clean_chunk.find("__RETRY_END__") + len("__RETRY_END__")
                        clean_chunk = clean_chunk[:retry_start] + clean_chunk[retry_end:]

                    # Legacy client: plain text streaming (no SSE)
                    if not capabilities.supports_sse:
                        if clean_chunk:
                            yield clean_chunk
                        continue

                    event_id += 1

                    # SSE client without structured events: simple data streaming
                    if not capabilities.supports_structured_events:
                        if clean_chunk:
                            yield f"data: {clean_chunk}\n\n"
                        continue

                    # Modern client: full structured SSE events
                    # Check if chunk contains RAG metadata markers
                    if "__RAG_START__" in chunk and "__RAG_END__" in chunk:
                        # Extract RAG metadata and send as separate event
                        rag_start = chunk.find("__RAG_START__") + len("__RAG_START__")
                        rag_end = chunk.find("__RAG_END__")
                        rag_json = chunk[rag_start:rag_end]

                        # Send RAG metadata as rag_context event
                        yield f"event: rag_context\ndata: {rag_json}\nid: {event_id}\n\n"

                        # Send any text before/after RAG markers as text events
                        text_before = chunk[: chunk.find("__RAG_START__")]
                        text_after = chunk[rag_end + len("__RAG_END__") :]
                        remaining_text = text_before + text_after

                        if remaining_text:
                            event_id += 1
                            # Count tokens for this chunk
                            try:
                                from nyxgpt.token_counter import count_tokens

                                total_tokens += count_tokens(remaining_text)
                            except Exception as e:
                                log.debug(
                                    "Token counting failed for chunk, continuing without it: %s",
                                    e,
                                    extra={"component": "chat"},
                                )

                            elapsed = time.time() - start_time
                            yield f"event: text\ndata: {json.dumps({'content': remaining_text, 'tokens': total_tokens, 'elapsed': elapsed})}\nid: {event_id}\n\n"
                    elif "__RETRY_START__" in chunk and "__RETRY_END__" in chunk:
                        # Handle retry status messages (already formatted as JSON)
                        retry_start = chunk.find("__RETRY_START__") + len("__RETRY_START__")
                        retry_end = chunk.find("__RETRY_END__")
                        retry_json = chunk[retry_start:retry_end]
                        yield f"event: retry\ndata: {retry_json}\nid: {event_id}\n\n"
                    else:
                        # Regular text content
                        # Count tokens for this chunk
                        try:
                            from nyxgpt.token_counter import count_tokens

                            total_tokens += count_tokens(chunk)
                        except Exception as e:
                            log.debug(
                                "Token counting failed for chunk, continuing without it: %s",
                                e,
                                extra={"component": "chat"},
                            )

                        elapsed = time.time() - start_time
                        yield f"event: text\ndata: {json.dumps({'content': chunk, 'tokens': total_tokens, 'elapsed': elapsed})}\nid: {event_id}\n\n"

                if rag_used_in_stream:
                    prom_metrics.RAG_QUERIES_TOTAL.labels(source="chat").inc()

                try:
                    usage_analytics_module.record(
                        session=req.session,
                        model=chosen_model,
                        prompt_tokens=_count_usage_tokens(req.prompt),
                        completion_tokens=total_tokens,
                        duration_s=time.time() - start_time,
                        cfg=_req_cfg(request),
                    )
                except Exception:
                    log.debug("Usage analytics recording failed", exc_info=True)

                # Send done event with final stats (only for modern clients)
                if capabilities.supports_sse and capabilities.supports_structured_events:
                    event_id += 1
                    elapsed = time.time() - start_time
                    done_data = {
                        "event_id": event_id,
                        "total_tokens": total_tokens,
                        "elapsed": elapsed,
                    }
                    yield f"event: done\ndata: {json.dumps(done_data)}\nid: {event_id}\n\n"
            except Exception as e:
                # This is the operator's only server-side record of a streaming chat
                # failure: once the response has started, FastAPI/uvicorn have no
                # request context left to log with, so if we don't log here the
                # upstream Ollama/model-runtime detail (status, model, message) is
                # never written to api.log and is invisible in the log viewer.
                log.error(
                    "Chat stream failed",
                    extra={
                        "request_id": req_id,
                        "session": req.session,
                        "model": chosen_model,
                        "error": str(e),
                        "error_type": type(e).__name__,
                    },
                    exc_info=True,
                )
                # Send error event (only for SSE clients). The full exception
                # (message, type, traceback) is captured in the server-side log
                # above; the client only ever sees a generic message so no
                # internal detail is exposed (CodeQL py/stack-trace-exposure).
                if capabilities.supports_sse and capabilities.supports_structured_events:
                    event_id += 1
                    elapsed = time.time() - start_time
                    error_data = {
                        "error": "The model request failed. Please try again.",
                        "elapsed": elapsed,
                    }
                    yield f"event: error\ndata: {json.dumps(error_data)}\nid: {event_id}\n\n"
                raise

        # Adapt response headers based on client capabilities
        if capabilities.supports_sse:
            media_type = "text/event-stream; charset=utf-8"
        else:
            media_type = "text/plain; charset=utf-8"

        response_headers = {
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",  # Disable proxy buffering
            # Inform client of server capabilities
            "X-Server-Supports-SSE": "true",
            "X-Server-Supports-Structured-Events": "true",
            "X-Server-Supports-Streaming": "true",
            "X-Server-Version": "nyxgpt/1.0.0",
        }

        return StreamingResponse(
            _stream_sse(),
            media_type=media_type,
            headers=response_headers,
        )
    except HTTPException:
        # Re-raise HTTP exceptions (e.g., 422 from validation above)
        raise
    except Exception as e:
        log.error(f"Streaming setup failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail="Internal server error") from e


# Streaming chat endpoint
@api.post("/chat/stream")
def chat_stream_api(request: Request, req: ChatRequest):
    """Streaming chat endpoint (API v1).

    Delegates to shared streaming response helper.
    """
    return _create_streaming_response(request, req)


# Non-versioned alias for streaming chat endpoint
@app.post("/api/chat/stream")
def chat_stream_api_legacy(request: Request, req: ChatRequest):
    """Streaming chat endpoint (legacy, non-versioned).

    Delegates to shared streaming response helper.
    """
    return _create_streaming_response(request, req)


@api.post("/rag/ingest", response_model=RagIngestResponse)
def rag_ingest(_request: Request, req: RagIngestRequest) -> RagIngestResponse:
    """Chunk, embed, and store a document's text in the RAG vector store.

    Re-ingesting the same `doc_id` with changed content replaces the
    previous version (the response's `previous_hash` reflects any prior
    version). Raises `400` if ingestion fails (e.g. embedding or vector
    store errors).
    """
    try:
        result = ingest_document(
            doc_id=req.doc_id,
            text=req.text,
            metadata=req.metadata,
            ensure_schema=req.ensure_schema,
            collection=req.collection,
        )
        prom_metrics.RAG_INGESTS_TOTAL.labels(source="document", result="success").inc()
        return RagIngestResponse(
            doc_id=req.doc_id,
            chunks_ingested=result["chunks_ingested"],
            status=result["status"],
            doc_hash=result["doc_hash"],
            previous_hash=result["previous_hash"],
            collection=req.collection,
        )
    except Exception as e:
        prom_metrics.RAG_INGESTS_TOTAL.labels(source="document", result="failure").inc()
        raise HTTPException(status_code=400, detail=str(e)) from e


@api.get("/rag/documents/{doc_id}", response_model=RagDocumentInfo)
def rag_document_info(
    _request: Request, doc_id: str, collection: str = "default"
) -> RagDocumentInfo:
    """Get document version and metadata information."""
    from nyxgpt.rag.vectorstore_cassandra import CassandraVectorStore

    store = CassandraVectorStore(collection=collection)
    try:
        info = store.get_document_info(doc_id)
        if not info:
            raise HTTPException(
                status_code=404,
                detail=f"Document '{doc_id}' not found in collection '{collection}'",
            )
        return RagDocumentInfo(**info)
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e)) from e
    finally:
        store.close()


@api.get("/rag/config")
def rag_config(request: Request) -> dict[str, Any]:
    """Get RAG configuration including score thresholds for visual indicators.

    Returns:
        Dictionary containing:
        - min_score: Minimum score threshold for retrieval
        - good_score_threshold: Threshold for high-confidence results (green)
        - medium_score_threshold: Threshold for medium-confidence results (yellow)
    """
    cfg = _req_cfg(request)
    return {
        "min_score": get_rag_min_score(cfg),
        "good_score_threshold": get_rag_good_score_threshold(cfg),
        "medium_score_threshold": get_rag_medium_score_threshold(cfg),
    }


@api.get("/rag/cache/stats", response_model=QueryCacheStatsResponse)
def rag_cache_stats(_request: Request) -> QueryCacheStatsResponse:
    """Get hit rate, size, and configuration details for the RAG query result cache.

    Returns `enabled=False` with zeroed stats if query result caching is
    disabled (`[cache] query_cache_enabled = false`).
    """
    stats = get_query_cache_stats()
    max_size = stats["max_size"]
    ttl_seconds = stats["ttl_seconds"]
    return QueryCacheStatsResponse(
        hits=int(stats["hits"]),  # type: ignore[arg-type]
        misses=int(stats["misses"]),  # type: ignore[arg-type]
        hit_rate=float(stats["hit_rate"]),  # type: ignore[arg-type]
        size=int(stats["size"]),  # type: ignore[arg-type]
        enabled=bool(stats["enabled"]),
        backend=str(stats["backend"]),
        max_size=int(max_size) if max_size is not None else None,
        ttl_seconds=int(ttl_seconds) if ttl_seconds is not None else None,
        rag_enabled=bool(stats["rag_enabled"]),
    )


@api.post("/rag/cache/clear", response_model=QueryCacheClearResponse)
def rag_cache_clear(_request: Request) -> QueryCacheClearResponse:
    """Manually clear the RAG query result cache.

    The cache is also cleared automatically on document ingestion/update,
    collection deletion, and collection re-indexing, so this is mainly
    useful for manual troubleshooting.
    """
    clear_query_cache()
    return QueryCacheClearResponse(status="Query result cache cleared")


def _configured_embedding_model(
    store: Any, stored_settings: dict[str, Any], docs: list[dict[str, Any]] | None = None
) -> str | None:
    """Resolve the collection's *configured* embedding model.

    Prefers the explicitly stored per-collection setting; falls back to the
    single observed model across ingested chunks when unambiguous. Returns
    None if nothing is configured and chunks use more than one model.

    `docs` may be passed in when the caller already fetched `list_docs()`, to
    avoid an extra round-trip; otherwise it's fetched lazily only if needed.
    """
    embedding_model_raw = stored_settings.get("embedding_model")
    if isinstance(embedding_model_raw, str):
        return embedding_model_raw

    if docs is None:
        docs = store.list_docs()

    observed_models = list({d["embedding_model"] for d in docs if d.get("embedding_model")})
    return observed_models[0] if len(observed_models) == 1 else None


@api.get("/rag/collections", response_model=CollectionsListResponse)
def rag_collections_list(_request: Request) -> CollectionsListResponse:
    """List all RAG collections with their statistics.

    Returns information about each collection including:
    - Collection name
    - Number of documents
    - Total number of chunks
    - Configured embedding model (from stored settings, or derived from
      ingested chunks if unambiguous)
    - Embedding models observed in ingested chunks
    """
    from nyxgpt.api_models import CollectionsListResponse
    from nyxgpt.rag.vectorstore_cassandra import CassandraVectorStore

    # Get list of all collections
    temp_store = CassandraVectorStore()
    try:
        collection_names = temp_store.list_collections()
    except Exception as e:
        log.error(f"Failed to list collections: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to list collections: {str(e)}") from e
    finally:
        temp_store.close()

    # Gather stats for each collection
    collections_info = []
    for coll_name in collection_names:
        store = CassandraVectorStore(collection=coll_name)
        try:
            docs = store.list_docs()
            doc_count = len(docs)
            chunk_count = sum(d["chunks"] for d in docs)

            # Get unique embedding models observed in ingested chunks
            embedding_models = list(
                {d["embedding_model"] for d in docs if d.get("embedding_model")}
            )
            embedding_models.sort()

            collections_info.append(
                CollectionInfo(
                    name=coll_name,
                    doc_count=doc_count,
                    chunk_count=chunk_count,
                    embedding_model=_configured_embedding_model(
                        store, store.get_collection_settings(), docs
                    ),
                    embedding_models=embedding_models,
                )
            )
        except Exception as e:
            log.error(f"Failed to get stats for collection '{coll_name}': {e}", exc_info=True)
            # Include collection with zero stats on error
            collections_info.append(
                CollectionInfo(
                    name=coll_name,
                    doc_count=0,
                    chunk_count=0,
                    embedding_model=None,
                    embedding_models=[],
                )
            )
        finally:
            store.close()

    return CollectionsListResponse(collections=collections_info)


@api.post("/rag/collections", response_model=CreateCollectionResponse, status_code=201)
def rag_collection_create(
    _request: Request, body: CreateCollectionRequest
) -> CreateCollectionResponse:
    """Create a new RAG collection.

    Creates a new empty collection with the specified embedding dimension.
    Collection names must be alphanumeric with underscores only (no hyphens),
    and are capped in length by Cassandra's identifier limit (see
    `max_collection_name_length`).

    Args:
        body: Collection creation request with name and embedding_dim

    Returns:
        CreateCollectionResponse with collection name and status
    """
    from nyxgpt.rag.vectorstore_cassandra import (
        COLLECTION_NAME_PATTERN,
        CassandraVectorStore,
        max_collection_name_length,
    )

    # Validate collection name
    collection_name = body.name.strip()

    # Prevent creating 'default' collection manually
    if collection_name == "default":
        raise HTTPException(
            status_code=400,
            detail="Cannot manually create 'default' collection. It is automatically managed.",
        )

    # Validate name format (alphanumeric and underscores only - no hyphens for Cassandra compatibility)
    if not COLLECTION_NAME_PATTERN.match(collection_name):
        raise HTTPException(
            status_code=400,
            detail=(
                "Collection name must contain only letters, numbers, and underscores "
                "(no hyphens, spaces, or other characters)."
            ),
        )

    max_name_len = max_collection_name_length()
    if len(collection_name) > max_name_len:
        raise HTTPException(
            status_code=400,
            detail=f"Collection name must be at most {max_name_len} characters long.",
        )

    # Validate embedding dimension
    if body.embedding_dim <= 0 or body.embedding_dim > 10000:
        raise HTTPException(
            status_code=400,
            detail="Embedding dimension must be between 1 and 10000.",
        )

    store = CassandraVectorStore(collection=collection_name)
    try:
        # Check if collection already exists
        existing_collections = store.list_collections()
        if collection_name in existing_collections:
            raise HTTPException(
                status_code=409,
                detail=f"Collection '{collection_name}' already exists.",
            )

        # Create the collection schema
        store.ensure_schema(embedding_dim=body.embedding_dim, collection=collection_name)

        # Persist the submitted embedding model as this collection's configured
        # setting, so it's used by ingestion and shown on the collection card
        # instead of being silently discarded.
        if body.embedding_model:
            store.update_collection_settings(embedding_model=body.embedding_model)

        return CreateCollectionResponse(
            collection=collection_name,
            status=f"Collection '{collection_name}' created successfully",
            embedding_dim=body.embedding_dim,
        )
    except HTTPException:
        # Re-raise HTTP exceptions
        raise
    except ImportError as e:
        log.error(f"Cassandra driver import error: {e}", exc_info=True)
        raise HTTPException(
            status_code=503, detail="RAG service unavailable: Cassandra driver not found"
        ) from e
    except Exception as e:
        log.error(f"Failed to create collection '{collection_name}': {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to create collection: {str(e)}") from e
    finally:
        store.close()


@api.post("/rag/collections/{collection_name}/clear", response_model=CollectionClearResponse)
def rag_collection_clear(_request: Request, collection_name: str) -> CollectionClearResponse:
    """Clear all documents and chunks from a RAG collection.

    This removes every document/chunk from the collection's vector table via
    a Cassandra `TRUNCATE`. The collection itself and its stored settings
    (embedding model, chunk size/overlap) are left in place -- only its
    contents are removed. This operation cannot be undone.

    To remove the collection entirely (data, settings, and the backing
    Cassandra table), use `DELETE /rag/collections/{collection_name}`.
    """
    from nyxgpt.rag.vectorstore_cassandra import CassandraVectorStore

    # Prevent clearing the default collection via a more descriptive error
    if collection_name == "default":
        raise HTTPException(
            status_code=400,
            detail="Cannot clear the 'default' collection. This collection is protected.",
        )

    store = CassandraVectorStore(collection=collection_name)
    try:
        existing_collections = store.list_collections()
        if collection_name not in existing_collections:
            raise HTTPException(
                status_code=404,
                detail=f"Collection '{collection_name}' not found.",
            )

        # Truncate the collection (remove all data, keep the collection/settings)
        store.truncate()

        # Invalidate query result cache: the document set changed, so any
        # cached retrieval results may now be stale.
        clear_query_cache()

        return CollectionClearResponse(
            collection=collection_name,
            status=(
                f"Collection '{collection_name}' has been cleared. All documents and "
                "chunks were removed; the collection and its settings remain."
            ),
            doc_count=0,
            chunk_count=0,
        )
    except HTTPException:
        raise
    except ImportError as e:
        # Cassandra driver not available
        log.error(f"Cassandra driver import error: {e}", exc_info=True)
        raise HTTPException(
            status_code=503, detail="RAG service unavailable: Cassandra driver not found"
        ) from e
    except Exception as e:
        # Catch database errors and other issues
        log.error(f"Failed to clear collection '{collection_name}': {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to clear collection: {str(e)}") from e
    finally:
        store.close()


@api.delete("/rag/collections/{collection_name}", response_model=CollectionDeleteResponse)
def rag_collection_delete(_request: Request, collection_name: str) -> CollectionDeleteResponse:
    """Permanently delete a RAG collection.

    Removes the collection's backing Cassandra table (all documents and
    chunks) along with its stored settings/metadata. This cannot be undone.

    The 'default' collection cannot be deleted: chat and ingestion fall back
    to it by default. To empty a collection's documents while keeping the
    collection and its settings, use
    `POST /rag/collections/{collection_name}/clear` instead.
    """
    from nyxgpt.rag.vectorstore_cassandra import CassandraVectorStore

    # Prevent deletion of the default collection via a more descriptive error
    if collection_name == "default":
        raise HTTPException(
            status_code=400,
            detail=(
                "Cannot delete the 'default' collection. It is protected because chat "
                "and ingestion fall back to it by default."
            ),
        )

    store = CassandraVectorStore(collection=collection_name)
    try:
        existing_collections = store.list_collections()
        if collection_name not in existing_collections:
            raise HTTPException(
                status_code=404,
                detail=f"Collection '{collection_name}' not found.",
            )

        # Drop the backing table and remove any stored settings/metadata.
        store.drop_collection()
        store.delete_collection_settings()

        # Invalidate query result cache: the collection no longer exists, so
        # any cached retrieval results referencing it may now be stale.
        clear_query_cache()

        return CollectionDeleteResponse(
            collection=collection_name,
            status=f"Collection '{collection_name}' has been permanently deleted.",
        )
    except HTTPException:
        raise
    except ValueError as e:
        # drop_collection() also guards against dropping the default collection
        raise HTTPException(status_code=400, detail=str(e)) from e
    except ImportError as e:
        # Cassandra driver not available
        log.error(f"Cassandra driver import error: {e}", exc_info=True)
        raise HTTPException(
            status_code=503, detail="RAG service unavailable: Cassandra driver not found"
        ) from e
    except Exception as e:
        # Catch database errors and other issues
        log.error(f"Failed to delete collection '{collection_name}': {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=f"Failed to delete collection: {str(e)}") from e
    finally:
        store.close()


@api.post("/rag/collections/{collection_name}/reindex", response_model=ReindexCollectionResponse)
def rag_collection_reindex(
    _request: Request, collection_name: str, body: ReindexCollectionRequest
) -> ReindexCollectionResponse:
    """Re-index a collection with a different embedding model.

    This operation regenerates embeddings for all chunks in the collection
    using the specified target embedding model.

    WARNING: This is a long-running operation. For large collections, this may
    take several minutes.

    Args:
        collection_name: Name of collection to re-index
        body: Re-index request with target_embedding_model and embedding_dim

    Returns:
        ReindexCollectionResponse with status and progress
    """
    from nyxgpt.rag.embeddings import embed_texts
    from nyxgpt.rag.vectorstore_cassandra import CassandraVectorStore

    # Prevent re-indexing default collection to avoid accidents
    if collection_name == "default":
        raise HTTPException(
            status_code=400,
            detail="Cannot re-index the 'default' collection. Create a new collection instead.",
        )

    # Validate embedding dimension
    if body.embedding_dim <= 0 or body.embedding_dim > 10000:
        raise HTTPException(
            status_code=400,
            detail="Embedding dimension must be between 1 and 10000.",
        )

    store = CassandraVectorStore(collection=collection_name)
    try:
        # Verify collection exists
        existing_collections = store.list_collections()
        if collection_name not in existing_collections:
            raise HTTPException(
                status_code=404,
                detail=f"Collection '{collection_name}' not found.",
            )

        # Get all chunks from the collection
        log.info(
            f"Re-indexing collection '{collection_name}' with model '{body.target_embedding_model}'"
        )
        chunks = store.get_all_chunks()

        if not chunks:
            return ReindexCollectionResponse(
                collection=collection_name,
                status="Collection is empty, nothing to re-index",
                chunks_processed=0,
                chunks_total=0,
            )

        chunks_total = len(chunks)
        log.info(f"Found {chunks_total} chunks to re-index in collection '{collection_name}'")

        # Extract text from each chunk for re-embedding
        texts = [chunk["text"] for chunk in chunks]

        # Re-generate embeddings with new model
        log.info(f"Generating new embeddings with model '{body.target_embedding_model}'")
        try:
            embed_texts(
                texts,
                model=body.target_embedding_model,
                dimension=body.embedding_dim,
            )
        except Exception as e:
            log.error(f"Failed to generate embeddings: {e}", exc_info=True)
            raise HTTPException(
                status_code=500,
                detail=f"Failed to generate embeddings with model '{body.target_embedding_model}': {str(e)}",
            ) from e

        # Invalidate query result cache: cached results may reflect the
        # collection's pre-reindex embeddings.
        clear_query_cache()

        # Return success response
        return ReindexCollectionResponse(
            collection=collection_name,
            status="Successfully re-indexed collection",
            chunks_processed=chunks_total,
            chunks_total=chunks_total,
        )
    except HTTPException:
        raise
    except Exception as e:
        log.error(f"Failed to re-index collection '{collection_name}': {e}", exc_info=True)
        raise HTTPException(
            status_code=500, detail=f"Failed to re-index collection: {str(e)}"
        ) from e
    finally:
        store.close()


@api.get("/rag/collections/{collection_name}/settings", response_model=CollectionSettingsResponse)
def rag_collection_get_settings(
    _request: Request, collection_name: str
) -> CollectionSettingsResponse:
    """Get settings for a collection.

    Returns the current configuration for a collection, including:
    - Embedding model (from stored settings or derived from documents)
    - Default chunk size (from stored settings or global config)
    - Default chunk overlap (from stored settings or global config)

    Args:
        collection_name: Name of collection

    Returns:
        CollectionSettingsResponse with current settings
    """
    from nyxgpt.rag.vectorstore_cassandra import CassandraVectorStore

    store = CassandraVectorStore(collection=collection_name)
    try:
        # Verify collection exists
        existing_collections = store.list_collections()
        if collection_name not in existing_collections:
            raise HTTPException(
                status_code=404,
                detail=f"Collection '{collection_name}' not found.",
            )

        # Get stored settings (if any)
        stored_settings = store.get_collection_settings()

        # Get global defaults from config
        cfg = load_config(None)
        global_chunk_size = cfg.getint("rag", "chunk_size", fallback=1000)
        global_chunk_overlap = cfg.getint("rag", "chunk_overlap", fallback=200)

        # If no stored settings, derive embedding_model from documents
        embedding_model = _configured_embedding_model(store, stored_settings)

        # Use stored settings if available, otherwise fall back to global config
        chunk_size_raw = stored_settings.get("chunk_size")
        chunk_size: int = chunk_size_raw if isinstance(chunk_size_raw, int) else global_chunk_size

        chunk_overlap_raw = stored_settings.get("chunk_overlap")
        chunk_overlap: int = (
            chunk_overlap_raw if isinstance(chunk_overlap_raw, int) else global_chunk_overlap
        )

        return CollectionSettingsResponse(
            collection=collection_name,
            settings=CollectionSettings(
                embedding_model=embedding_model,
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
            ),
        )

    except HTTPException:
        raise
    except ImportError as e:
        log.error(f"Cassandra driver import error: {e}", exc_info=True)
        raise HTTPException(
            status_code=503, detail="RAG service unavailable: Cassandra driver not found"
        ) from e
    except Exception as e:
        log.error(f"Failed to get settings for collection '{collection_name}': {e}", exc_info=True)
        raise HTTPException(
            status_code=500, detail=f"Failed to get collection settings: {str(e)}"
        ) from e
    finally:
        store.close()


@api.put("/rag/collections/{collection_name}/settings", response_model=CollectionSettingsResponse)
def rag_collection_update_settings(
    _request: Request, collection_name: str, body: CollectionSettings
) -> CollectionSettingsResponse:
    """Update settings for a collection.

    This allows configuring per-collection settings like:
    - Preferred embedding model
    - Default chunk size
    - Default chunk overlap

    Note: Settings are stored separately from documents and are used as defaults
    when ingesting new documents to this collection.

    Args:
        collection_name: Name of collection
        body: New settings to apply

    Returns:
        CollectionSettingsResponse with updated settings
    """
    from nyxgpt.rag.vectorstore_cassandra import CassandraVectorStore

    # Verify collection exists first
    store = CassandraVectorStore(collection=collection_name)
    try:
        existing_collections = store.list_collections()
        if collection_name not in existing_collections:
            raise HTTPException(
                status_code=404,
                detail=f"Collection '{collection_name}' not found.",
            )

        # Update collection settings
        store.update_collection_settings(
            embedding_model=body.embedding_model,
            chunk_size=body.chunk_size,
            chunk_overlap=body.chunk_overlap,
        )

        # Return updated settings
        return CollectionSettingsResponse(
            collection=collection_name,
            settings=body,
        )

    except HTTPException:
        raise
    except ImportError as e:
        log.error(f"Cassandra driver import error: {e}", exc_info=True)
        raise HTTPException(
            status_code=503, detail="RAG service unavailable: Cassandra driver not found"
        ) from e
    except Exception as e:
        log.error(
            f"Failed to update settings for collection '{collection_name}': {e}", exc_info=True
        )
        raise HTTPException(
            status_code=500, detail=f"Failed to update collection settings: {str(e)}"
        ) from e
    finally:
        store.close()


@api.get("/rag/documents")
def rag_documents_list(
    _request: Request,
    collection: str = Query("default", description="Vector store collection name"),
) -> dict[str, Any]:
    """List all documents in the RAG vector store.

    Returns list of documents with metadata including:
    - doc_id: Document identifier
    - chunks: Number of chunks stored
    - embedding_model: Model used for embeddings
    - filename: Filename (from metadata if available)
    - tags: Document tags (from metadata if available)
    - ingested_at: Ingestion timestamp (from metadata if available)
    """
    from nyxgpt.rag.vectorstore_cassandra import CassandraVectorStore, parse_metadata

    store = CassandraVectorStore(collection=collection)
    try:
        docs = store.list_docs()

        # Enrich with metadata from individual chunks (get first chunk per doc)
        enriched_docs = []
        for doc in docs:
            doc_id = doc["doc_id"]
            # Query for one chunk to get metadata
            try:
                from cassandra.query import SimpleStatement

                stmt = SimpleStatement(
                    f"SELECT metadata, ingested_at FROM {store.table_name} WHERE doc_id = %s LIMIT 1"
                )
                rows = store.session.execute(stmt, (doc_id,))
                row = next(iter(rows), None)

                metadata = (
                    parse_metadata(row.metadata, doc_id=doc_id)
                    if row and hasattr(row, "metadata")
                    else {}
                )
                ingested_at = (
                    row.ingested_at.isoformat()
                    if row and hasattr(row, "ingested_at") and row.ingested_at
                    else None
                )

                enriched_docs.append(
                    {
                        "doc_id": doc_id,
                        "chunks": doc["chunks"],
                        "embedding_model": doc.get("embedding_model"),
                        "filename": metadata.get("filename") if metadata else None,
                        "tags": metadata.get("tags") if metadata else None,
                        "ingested_at": ingested_at,
                    }
                )
            except Exception as e:
                log.warning(f"Failed to enrich metadata for doc_id={doc_id}: {e}")
                # Fallback: return basic info without metadata
                enriched_docs.append(
                    {
                        "doc_id": doc_id,
                        "chunks": doc["chunks"],
                        "embedding_model": doc.get("embedding_model"),
                        "filename": None,
                        "tags": None,
                        "ingested_at": None,
                    }
                )

        return {"documents": enriched_docs}
    except Exception as e:
        log.error(f"Failed to list RAG documents: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e)) from e
    finally:
        store.close()


def _resolve_and_validate_collection(collection: str | None) -> str:
    """Normalize a user-supplied collection name and confirm it exists.

    "default" is always valid (it's automatically managed). Any other name
    must already exist (created via the admin collections API) -- this
    catches a typo'd collection name before it silently queries/creates an
    unrelated empty table.

    Raises:
        HTTPException 400: The named collection does not exist.
    """
    name = (collection or "default").strip() or "default"
    if name == "default":
        return name

    from nyxgpt.rag.vectorstore_cassandra import CassandraVectorStore

    store = CassandraVectorStore()
    try:
        if name not in store.list_collections():
            raise HTTPException(status_code=400, detail=f"Unknown collection '{name}'.")
    finally:
        store.close()
    return name


@api.post("/rag/query", response_model=RagQueryResponse)
def rag_query(_request: Request, req: RagQueryRequest) -> RagQueryResponse:
    """Run a RAG similarity search and return the matching chunks.

    Supports optional metadata filtering (`doc_ids`, `filename`, `tags`,
    `date_from`/`date_to`) and an optional `collection` to scope the search
    (defaults to the "default" collection). When `req.debug_mode` is set,
    the response also includes a `debug_info` breakdown of per-stage timings
    and score distribution. Raises `400` on any retrieval failure (e.g.
    malformed date filters, unknown collection, vector store errors).
    """
    try:
        from datetime import datetime

        from nyxgpt.rag.vectorstore_cassandra import MetadataFilter

        collection = _resolve_and_validate_collection(req.collection)

        # Build metadata filter if any filter params are provided
        metadata_filter = None
        if any([req.doc_ids, req.filename, req.tags, req.date_from, req.date_to]):
            # Parse dates
            date_from_dt = datetime.fromisoformat(req.date_from) if req.date_from else None
            date_to_dt = datetime.fromisoformat(req.date_to) if req.date_to else None

            metadata_filter = MetadataFilter(
                doc_ids=req.doc_ids,
                filename=req.filename,
                tags=req.tags,
                date_from=date_from_dt,
                date_to=date_to_dt,
            )

        result = retrieve_context(
            req.query,
            top_k=req.top_k,
            debug_mode=req.debug_mode,
            collection=collection,
            metadata_filter=metadata_filter,
        )

        if req.debug_mode:
            # Type narrowing: debug_mode=True means result is tuple[list[dict], RAGDebugInfo]
            from nyxgpt.rag.rag import RAGDebugInfo

            results, debug_info = cast(tuple[list[dict], RAGDebugInfo], result)
            # Convert RAGDebugInfo to RagDebugInfo (API model)
            from nyxgpt.api_models import RagDebugInfo

            api_debug_info = RagDebugInfo(
                total_time_ms=debug_info.total_time_ms,
                query_expansion_time_ms=debug_info.query_expansion_time_ms,
                embedding_time_ms=debug_info.embedding_time_ms,
                vector_search_time_ms=debug_info.vector_search_time_ms,
                filtering_time_ms=debug_info.filtering_time_ms,
                composition_time_ms=debug_info.composition_time_ms,
                original_query=debug_info.original_query,
                query_variants=debug_info.query_variants,
                num_queries=debug_info.num_queries,
                embedding_model=debug_info.embedding_model,
                embedding_dim=debug_info.embedding_dim,
                num_texts_embedded=debug_info.num_texts_embedded,
                batch_size=debug_info.batch_size,
                raw_results_count=debug_info.raw_results_count,
                score_min=debug_info.score_min,
                score_max=debug_info.score_max,
                score_mean=debug_info.score_mean,
                after_min_score_filter=debug_info.after_min_score_filter,
                after_dedupe_filter=debug_info.after_dedupe_filter,
                after_max_chunks_filter=debug_info.after_max_chunks_filter,
                total_chars_before_truncation=debug_info.total_chars_before_truncation,
                total_chars_after_truncation=debug_info.total_chars_after_truncation,
                chunks_included=debug_info.chunks_included,
                collection=debug_info.collection,
            )
        else:
            # Type narrowing: debug_mode=False means result is list[dict]
            results = cast(list[dict], result)
            api_debug_info = None

        results = annotate_chunk_numbering(results, collection=collection)
        out = [
            RagQueryResult(
                doc_id=str(r.get("doc_id", "")),
                chunk_id=int(r.get("chunk_id", 0)),
                text=str(r.get("text", "")),
                score=float(r.get("score", 0.0)),
                similarity_score=r.get("similarity_score"),
                collection=r.get("collection"),
                chunk_number=r.get("chunk_number"),
                total_chunks=r.get("total_chunks"),
            )
            for r in results
        ]
        prom_metrics.RAG_QUERIES_TOTAL.labels(source="rag_query").inc()
        return RagQueryResponse(results=out, debug_info=api_debug_info)
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e)) from e


@api.post("/rag/metrics/query", response_model=RagMetricsQueryResponse)
def rag_metrics_query(_request: Request, req: RagMetricsQueryRequest) -> RagMetricsQueryResponse:
    """Query RAG with comprehensive evaluation metrics.

    This endpoint extends the standard RAG query with evaluation metrics including:
    - Retrieval accuracy (hit rate, unique docs, score distribution)
    - Latency tracking (per-stage breakdowns)
    - Hit rate analysis (success rate, threshold performance)

    Requires debug_mode=True to collect metrics.
    """
    try:
        from nyxgpt.config import get_rag_min_score

        collection = _resolve_and_validate_collection(req.collection)

        # Force debug mode to collect metrics
        result = retrieve_context(
            req.query, top_k=req.top_k, debug_mode=True, collection=collection
        )

        # Type narrowing: debug_mode=True means result is tuple[list[dict], RAGDebugInfo]
        from nyxgpt.rag.rag import RAGDebugInfo, compute_evaluation_metrics

        results, debug_info = cast(tuple[list[dict], RAGDebugInfo], result)

        # Convert RAGDebugInfo to RagDebugInfo (API model)
        from nyxgpt.api_models import (
            HitRateMetrics as ApiHitRateMetrics,
        )
        from nyxgpt.api_models import (
            LatencyMetrics as ApiLatencyMetrics,
        )
        from nyxgpt.api_models import (
            RagDebugInfo,
        )
        from nyxgpt.api_models import (
            RagEvaluationMetrics as ApiRagEvaluationMetrics,
        )
        from nyxgpt.api_models import (
            RetrievalAccuracyMetrics as ApiRetrievalAccuracyMetrics,
        )

        api_debug_info = RagDebugInfo(
            total_time_ms=debug_info.total_time_ms,
            query_expansion_time_ms=debug_info.query_expansion_time_ms,
            embedding_time_ms=debug_info.embedding_time_ms,
            vector_search_time_ms=debug_info.vector_search_time_ms,
            filtering_time_ms=debug_info.filtering_time_ms,
            composition_time_ms=debug_info.composition_time_ms,
            original_query=debug_info.original_query,
            query_variants=debug_info.query_variants,
            num_queries=debug_info.num_queries,
            embedding_model=debug_info.embedding_model,
            embedding_dim=debug_info.embedding_dim,
            num_texts_embedded=debug_info.num_texts_embedded,
            batch_size=debug_info.batch_size,
            raw_results_count=debug_info.raw_results_count,
            score_min=debug_info.score_min,
            score_max=debug_info.score_max,
            score_mean=debug_info.score_mean,
            after_min_score_filter=debug_info.after_min_score_filter,
            after_dedupe_filter=debug_info.after_dedupe_filter,
            after_max_chunks_filter=debug_info.after_max_chunks_filter,
            total_chars_before_truncation=debug_info.total_chars_before_truncation,
            total_chars_after_truncation=debug_info.total_chars_after_truncation,
            chunks_included=debug_info.chunks_included,
            collection=debug_info.collection,
        )

        # Compute evaluation metrics if requested
        evaluation_metrics = None
        if req.collect_metrics:
            cfg = load_config(None)
            min_score = get_rag_min_score(cfg)
            eval_metrics = compute_evaluation_metrics(results, debug_info, min_score)

            # Convert to API models
            evaluation_metrics = ApiRagEvaluationMetrics(
                retrieval_accuracy=ApiRetrievalAccuracyMetrics(
                    results_returned=eval_metrics.retrieval_accuracy.results_returned,
                    query_success=eval_metrics.retrieval_accuracy.query_success,
                    unique_docs_retrieved=eval_metrics.retrieval_accuracy.unique_docs_retrieved,
                    total_chunks_retrieved=eval_metrics.retrieval_accuracy.total_chunks_retrieved,
                    score_distribution=eval_metrics.retrieval_accuracy.score_distribution,
                ),
                latency=ApiLatencyMetrics(
                    total_time_ms=eval_metrics.latency.total_time_ms,
                    stage_timings=eval_metrics.latency.stage_timings,
                    percentiles=eval_metrics.latency.percentiles,
                ),
                hit_rate=ApiHitRateMetrics(
                    query_success_rate=eval_metrics.hit_rate.query_success_rate,
                    total_queries=eval_metrics.hit_rate.total_queries,
                    successful_queries=eval_metrics.hit_rate.successful_queries,
                    failed_queries=eval_metrics.hit_rate.failed_queries,
                    avg_top_score=eval_metrics.hit_rate.avg_top_score,
                    score_above_threshold_rate=eval_metrics.hit_rate.score_above_threshold_rate,
                ),
                query_id=eval_metrics.query_id,
                timestamp=eval_metrics.timestamp,
            )

        results = annotate_chunk_numbering(results, collection=collection)
        out = [
            RagQueryResult(
                doc_id=str(r.get("doc_id", "")),
                chunk_id=int(r.get("chunk_id", 0)),
                text=str(r.get("text", "")),
                score=float(r.get("score", 0.0)),
                similarity_score=r.get("similarity_score"),
                collection=r.get("collection"),
                chunk_number=r.get("chunk_number"),
                total_chunks=r.get("total_chunks"),
            )
            for r in results
        ]
        prom_metrics.RAG_QUERIES_TOTAL.labels(source="rag_query").inc()

        return RagMetricsQueryResponse(
            results=out,
            debug_info=api_debug_info,
            evaluation_metrics=evaluation_metrics,
        )
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e)) from e


@api.get("/sessions/{name}/metadata")
def get_session_metadata(name: str) -> dict[str, Any]:
    """Get metadata for a specific session."""
    cfg = load_config(None)
    sessions_dir = get_sessions_dir(cfg)
    sf = sessions.session_file_for(name, sessions_dir)
    mf = sessions.meta_file_for(sf)

    # Create session files if they don't exist
    if not sessions.session_file_exists(sf):
        sessions.save_session_messages(sf, [])

    meta = sessions.load_session_meta(mf)
    meta = sessions.ensure_meta_defaults(meta)
    sessions.save_session_meta(mf, meta)

    return meta


@api.post("/sessions/{name}/rag/enable")
def enable_session_rag(name: str) -> dict[str, Any]:
    """Enable RAG for a specific session.

    Note: This endpoint uses non-atomic read-modify-write operations.
    For single-user applications, this is acceptable. For multi-user
    deployments, consider file locking or database transactions.
    """
    cfg = load_config(None)
    sessions_dir = get_sessions_dir(cfg)
    sf = sessions.session_file_for(name, sessions_dir)
    mf = sessions.meta_file_for(sf)

    # Create session files if they don't exist
    if not sessions.session_file_exists(sf):
        sessions.save_session_messages(sf, [])

    meta = sessions.load_session_meta(mf)
    meta = sessions.ensure_meta_defaults(meta)
    meta["rag_enabled"] = True
    sessions.save_session_meta(mf, meta)

    return {"session": name, "rag_enabled": True}


@api.post("/sessions/{name}/rag/disable")
def disable_session_rag(name: str) -> dict[str, Any]:
    """Disable RAG for a specific session.

    Note: This endpoint uses non-atomic read-modify-write operations.
    For single-user applications, this is acceptable. For multi-user
    deployments, consider file locking or database transactions.
    """
    cfg = load_config(None)
    sessions_dir = get_sessions_dir(cfg)
    sf = sessions.session_file_for(name, sessions_dir)
    mf = sessions.meta_file_for(sf)

    # Create session files if they don't exist
    if not sessions.session_file_exists(sf):
        sessions.save_session_messages(sf, [])

    meta = sessions.load_session_meta(mf)
    meta = sessions.ensure_meta_defaults(meta)
    meta["rag_enabled"] = False
    sessions.save_session_meta(mf, meta)

    return {"session": name, "rag_enabled": False}


@api.get("/sessions/{name}/documents", response_model=SessionDocumentsResponse)
def list_session_documents(name: str) -> SessionDocumentsResponse:
    """List document IDs force-included for a session's RAG context.

    Returns the list of documents attached to the session that are
    always retrieved when RAG is enabled, regardless of other filters.
    """
    cfg = load_config(None)
    sessions_dir = get_sessions_dir(cfg)
    sf = sessions.session_file_for(name, sessions_dir)
    mf = sessions.meta_file_for(sf)

    meta = sessions.load_session_meta(mf)
    attached = meta.get("attached_doc_ids", [])
    if not isinstance(attached, list):
        attached = []

    return SessionDocumentsResponse(session=name, attached_doc_ids=attached)


@api.post("/sessions/{name}/documents", response_model=SessionDocumentsResponse)
def attach_document_to_session(name: str, req: AttachDocumentRequest) -> SessionDocumentsResponse:
    """Attach a document to a session for force-inclusion in RAG context.

    Attached documents are always retrieved when RAG is enabled for the
    session, regardless of any other metadata filters in the chat request.
    """
    cfg = load_config(None)
    sessions_dir = get_sessions_dir(cfg)
    sf = sessions.session_file_for(name, sessions_dir)
    mf = sessions.meta_file_for(sf)

    if not sessions.session_file_exists(sf):
        sessions.save_session_messages(sf, [])

    meta = sessions.load_session_meta(mf)
    meta = sessions.ensure_meta_defaults(meta)

    raw = meta.get("attached_doc_ids", [])
    attached: list[str] = raw if isinstance(raw, list) else []

    if req.doc_id not in attached:
        attached = attached + [req.doc_id]
        meta["attached_doc_ids"] = attached
        sessions.save_session_meta(mf, meta)

    return SessionDocumentsResponse(session=name, attached_doc_ids=attached)


@api.delete("/sessions/{name}/documents/{doc_id}", response_model=SessionDocumentsResponse)
def detach_document_from_session(name: str, doc_id: str) -> SessionDocumentsResponse:
    """Detach a document from a session, removing it from force-inclusion.

    After detaching, the document will no longer be force-included in
    RAG retrieval for this session.
    """
    cfg = load_config(None)
    sessions_dir = get_sessions_dir(cfg)
    sf = sessions.session_file_for(name, sessions_dir)
    mf = sessions.meta_file_for(sf)

    meta = sessions.load_session_meta(mf)
    meta = sessions.ensure_meta_defaults(meta)

    raw = meta.get("attached_doc_ids", [])
    attached: list[str] = raw if isinstance(raw, list) else []

    if doc_id in attached:
        attached = [d for d in attached if d != doc_id]
        meta["attached_doc_ids"] = attached
        sessions.save_session_meta(mf, meta)

    return SessionDocumentsResponse(session=name, attached_doc_ids=attached)


@api.post("/rag/upload", response_model=RagIngestResponse)
async def rag_upload_file(
    file: UploadFile = File(...),
    doc_id: str | None = None,
    collection: str | None = None,
) -> RagIngestResponse:
    """Upload and ingest a document for RAG with proper markdown parsing.

    `collection` selects the target collection (defaults to "default" when
    omitted); it must already exist (see `POST /rag/collections`).
    """
    resolved_collection = _resolve_and_validate_collection(collection)

    # Validate file type
    allowed_types = {".txt", ".md", ".json", ".pdf", ".pptx", ".docx", ".epub", ".html", ".htm"}
    file_ext = Path(file.filename or "").suffix.lower()
    if file_ext not in allowed_types:
        raise HTTPException(
            status_code=400,
            detail=f"File type {file_ext} not supported. Allowed: {', '.join(allowed_types)}",
        )

    # Read file content
    content = await file.read()

    # Validate file size (prevent memory exhaustion)
    MAX_UPLOAD_SIZE = 10 * 1024 * 1024  # 10MB
    if len(content) > MAX_UPLOAD_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"File too large ({len(content)} bytes). Maximum size: {MAX_UPLOAD_SIZE} bytes (10MB)",
        )

    # Parse based on file type
    if file_ext == ".pdf":
        # Handle PDF with improved extraction (#2663)
        try:
            import pdfplumber
            from pypdf import PdfReader

            # Extract metadata using pypdf
            reader = PdfReader(io.BytesIO(content))
            metadata = {}
            if reader.metadata:
                # Extract common metadata fields
                for key in [
                    "/Title",
                    "/Author",
                    "/Subject",
                    "/Creator",
                    "/Producer",
                    "/CreationDate",
                    "/ModDate",
                ]:
                    if key in reader.metadata:
                        clean_key = key.lstrip("/")
                        metadata[clean_key] = str(reader.metadata[key])

            # Extract text with better formatting using pdfplumber
            text_parts = []

            # Add metadata section if available
            if metadata:
                meta_str = "\n".join(f"{k}: {v}" for k, v in metadata.items())
                text_parts.append(f"[Metadata]\n{meta_str}\n")

            # Try enhanced extraction with pdfplumber
            try:
                # Process each page with pdfplumber for better extraction
                with pdfplumber.open(io.BytesIO(content)) as pdf:
                    for _page_num, page in enumerate(pdf.pages, 1):
                        page_content = []

                        # Extract tables with preserved structure
                        tables = page.extract_tables()
                        if tables:
                            for table in tables:
                                # Format table as text grid
                                table_text = "\n".join(
                                    " | ".join(str(cell) if cell else "" for cell in row)
                                    for row in table
                                )
                                page_content.append(f"[Table]\n{table_text}\n")

                        # Extract text with layout preservation
                        # Use layout mode to preserve formatting and multi-column layouts
                        page_text = page.extract_text(layout=True)
                        if page_text:
                            page_content.append(page_text.strip())

                        # Combine page content
                        if page_content:
                            text_parts.append("\n\n".join(page_content))

            except Exception as plumber_error:
                # Fall back to basic pypdf extraction
                log.warning(f"pdfplumber extraction failed, falling back to pypdf: {plumber_error}")
                text_parts = []
                if metadata:
                    meta_str = "\n".join(f"{k}: {v}" for k, v in metadata.items())
                    text_parts.append(f"[Metadata]\n{meta_str}\n")
                for pdf_page in reader.pages:
                    page_text = pdf_page.extract_text()
                    if page_text:
                        text_parts.append(page_text)

            text = "\n\n".join(text_parts)

            # Check if OCR is needed for image-only pages (#2669)
            cfg = load_config(None)
            ocr_enabled = cfg.getboolean("pdf", "ocr_enabled", fallback=True)
            ocr_min_text_threshold = cfg.getint("pdf", "ocr_min_text_threshold", fallback=50)

            # Detect if PDF is image-only or has minimal text
            needs_ocr = (not text or len(text.strip()) < ocr_min_text_threshold) and ocr_enabled

            if needs_ocr:
                log.info(
                    f"PDF has minimal text ({len(text.strip()) if text else 0} chars), attempting OCR"
                )
                try:
                    import pytesseract
                    from pdf2image import convert_from_bytes

                    # Get OCR configuration
                    ocr_dpi = cfg.getint("pdf", "ocr_dpi", fallback=300)
                    ocr_lang = cfg.get("pdf", "ocr_lang", fallback="eng")
                    ocr_psm = cfg.getint("pdf", "ocr_psm", fallback=3)

                    # Configure tesseract if custom path is specified
                    tesseract_cmd = cfg.get("pdf", "tesseract_cmd", fallback=None)
                    if tesseract_cmd:
                        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd

                    # Convert PDF pages to images
                    log.debug(f"Converting PDF to images at {ocr_dpi} DPI")
                    images = convert_from_bytes(content, dpi=ocr_dpi)

                    # Extract text from each page using OCR
                    ocr_text_parts = []
                    if metadata:
                        meta_str = "\n".join(f"{k}: {v}" for k, v in metadata.items())
                        ocr_text_parts.append(f"[Metadata]\n{meta_str}\n")

                    # Configure OCR with PSM (Page Segmentation Mode)
                    # PSM 3 = Fully automatic page segmentation (default)
                    # PSM 6 = Assume a single uniform block of text
                    # PSM 11 = Sparse text - Find as much text as possible in no particular order
                    custom_config = f"--psm {ocr_psm}"

                    for page_num, image in enumerate(images, 1):
                        try:
                            page_ocr_text = pytesseract.image_to_string(
                                image, lang=ocr_lang, config=custom_config
                            )
                            if page_ocr_text and page_ocr_text.strip():
                                ocr_text_parts.append(
                                    f"[Page {page_num} (OCR)]\n{page_ocr_text.strip()}"
                                )
                                log.debug(
                                    f"OCR extracted {len(page_ocr_text.strip())} chars from page {page_num}"
                                )
                        except Exception as page_error:
                            log.warning(f"OCR failed for page {page_num}: {page_error}")
                            continue

                    # Use OCR text if extraction was successful
                    if ocr_text_parts:
                        ocr_text = "\n\n".join(ocr_text_parts)
                        if len(ocr_text.strip()) > len(text.strip() if text else ""):
                            log.info(
                                f"OCR extracted {len(ocr_text.strip())} chars (vs {len(text.strip()) if text else 0} from standard extraction)"
                            )
                            text = ocr_text
                        else:
                            log.info("OCR did not improve extraction, using standard extraction")
                    else:
                        log.warning("OCR produced no text")

                except ImportError as ocr_import_error:
                    missing_lib = str(ocr_import_error)
                    if "pytesseract" in missing_lib:
                        log.warning(
                            "pytesseract not installed, skipping OCR. Install with: pip install pytesseract"
                        )
                    elif "pdf2image" in missing_lib:
                        log.warning(
                            "pdf2image not installed, skipping OCR. Install with: pip install pdf2image"
                        )
                    else:
                        log.warning(f"OCR dependencies missing: {missing_lib}")
                except Exception as ocr_error:
                    log.warning(f"OCR extraction failed: {ocr_error}")

            # Validate extracted content (after OCR attempt)
            if not text or not text.strip():
                raise HTTPException(
                    status_code=400,
                    detail="PDF extraction produced no text. The file may be empty, image-only, or malformed. "
                    "If this is an image-based PDF, ensure Tesseract OCR is installed and configured.",
                )

        except ImportError as e:
            missing_lib = "pdfplumber" if "pdfplumber" in str(e) else "pypdf"
            raise HTTPException(
                status_code=400,
                detail=f"PDF support not available. Install {missing_lib}: pip install {missing_lib}",
            ) from e
        except HTTPException:
            # Re-raise HTTP exceptions without wrapping
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PDF parsing failed: {e}") from e

    elif file_ext == ".docx":
        # Handle DOCX (Microsoft Word)
        try:
            import zipfile

            from docx import Document
            from docx.opc.exceptions import PackageNotFoundError

            try:
                doc = Document(io.BytesIO(content))
            except PackageNotFoundError:
                raise HTTPException(
                    status_code=400,
                    detail="Invalid DOCX file: file is corrupted or not a valid Word document",
                ) from None
            except zipfile.BadZipFile:
                raise HTTPException(
                    status_code=400,
                    detail="Invalid DOCX file: file structure is corrupted (not a valid ZIP archive)",
                ) from None

            text_parts = []
            image_counter = 0

            for para in doc.paragraphs:
                para_text = para.text.strip()

                # Check for embedded images in this paragraph
                # Images are stored in runs within paragraphs
                for run in para.runs:
                    # Check if this run contains an image
                    if hasattr(run, "_element") and run._element is not None:
                        # Look for drawing elements that contain images
                        for drawing in run._element.findall(
                            ".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing"
                        ):
                            image_counter += 1
                            # Extract image description if available (alt text)
                            desc_elems = drawing.findall(
                                ".//{http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing}docPr"
                            )

                            image_desc = f"[Image {image_counter}"
                            if desc_elems and desc_elems[0].get("descr"):
                                image_desc += f": {desc_elems[0].get('descr')}"
                            elif desc_elems and desc_elems[0].get("name"):
                                image_desc += f": {desc_elems[0].get('name')}"
                            image_desc += "]"

                            text_parts.append(image_desc)

                if para_text:
                    # Preserve heading structure
                    if para.style and para.style.name.startswith("Heading"):
                        text_parts.append(f"\n## {para_text}\n")
                    else:
                        text_parts.append(para_text)

            # Handle tables
            for doc_table in doc.tables:
                table_rows: list[str] = []
                for row in doc_table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text:
                        table_rows.append(row_text)
                if table_rows:
                    text_parts.append("\n[Table]\n" + "\n".join(table_rows) + "\n")

            text = "\n\n".join(text_parts)

            # Check if document is empty
            if not text.strip():
                raise HTTPException(
                    status_code=400,
                    detail="DOCX file is empty or contains no extractable text",
                )

        except ImportError:
            raise HTTPException(
                status_code=400,
                detail="DOCX support not available. Install python-docx: pip install python-docx",
            ) from None
        except HTTPException:
            # Re-raise HTTP exceptions (our specific error messages)
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"DOCX parsing failed: {e}") from None

    elif file_ext == ".md":
        # Handle Markdown with proper parsing (#2667)
        try:
            import frontmatter
            from bs4 import BeautifulSoup
            from markdown import markdown

            # Parse frontmatter and content
            post = frontmatter.loads(content.decode("utf-8"))

            # Extract metadata from frontmatter (values stringified: `metadata`
            # is dict[str, str] in this scope and is only string-formatted below)
            metadata = {str(k): str(v) for k, v in post.metadata.items()} if post.metadata else {}

            # Convert markdown to plain text (strip HTML tags)
            # This preserves structure while making it searchable
            html = markdown(post.content)
            soup = BeautifulSoup(html, "html.parser")
            text = soup.get_text(separator="\n\n")

            # Prepend frontmatter as metadata section
            if metadata:
                meta_str = "\n".join(f"{k}: {v}" for k, v in metadata.items())
                text = f"[Metadata]\n{meta_str}\n\n{text}"

        except ImportError:
            # Fallback to plain text if libraries unavailable
            text = content.decode("utf-8")
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Markdown parsing failed: {e}") from None

    elif file_ext == ".json":
        # JSON files stored as formatted text
        import json

        try:
            data = json.loads(content.decode("utf-8"))
            text = json.dumps(data, indent=2, ensure_ascii=False)
        except json.JSONDecodeError as e:
            raise HTTPException(status_code=400, detail=f"Invalid JSON file: {e}") from None

    elif file_ext == ".pptx":
        # Handle PowerPoint presentations
        try:
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            text_parts = []

            # Process each slide in order
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts = []

                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text.strip())

                # Extract speaker notes if present
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if hasattr(notes_slide, "notes_text_frame") and notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_texts.append(f"[Speaker Notes]\n{notes_text}")

                # Add slide content if any text was found
                if slide_texts:
                    slide_content = f"[Slide {slide_num}]\n" + "\n\n".join(slide_texts)
                    text_parts.append(slide_content)

            # Join all slides with double newlines
            text = "\n\n".join(text_parts) if text_parts else ""

            if not text:
                raise HTTPException(
                    status_code=400, detail="PPTX file contains no extractable text"
                )

        except ImportError:
            raise HTTPException(
                status_code=400,
                detail="PPTX support not available. Install python-pptx: pip install python-pptx",
            ) from None
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"PPTX parsing failed: {e}") from None

    elif file_ext == ".epub":
        # Handle ePUB eBooks
        try:
            import ebooklib
            from bs4 import BeautifulSoup
            from ebooklib import epub

            try:
                book = epub.read_epub(io.BytesIO(content))
            except Exception as e:
                raise HTTPException(status_code=400, detail=f"Invalid ePUB file: {e}") from None

            text_parts = []

            # Helper to extract metadata values (handles tuple format from ebooklib)
            def extract_metadata_values(items):
                """Extract values from metadata items.

                ebooklib returns metadata as tuples (namespace, value).
                This helper extracts just the values.
                """
                values = []
                for item in items:
                    if isinstance(item, tuple):
                        # Tuple format: (namespace, value)
                        values.append(str(item[1] if len(item) > 1 else item[0]))
                    else:
                        # String format (fallback)
                        values.append(str(item))
                return values

            # Extract metadata
            metadata = {}
            if book.get_metadata("DC", "title"):
                metadata["Title"] = ", ".join(
                    extract_metadata_values(book.get_metadata("DC", "title"))
                )
            if book.get_metadata("DC", "creator"):
                metadata["Author"] = ", ".join(
                    extract_metadata_values(book.get_metadata("DC", "creator"))
                )
            if book.get_metadata("DC", "description"):
                metadata["Description"] = ", ".join(
                    extract_metadata_values(book.get_metadata("DC", "description"))
                )
            if book.get_metadata("DC", "publisher"):
                metadata["Publisher"] = ", ".join(
                    extract_metadata_values(book.get_metadata("DC", "publisher"))
                )
            if book.get_metadata("DC", "date"):
                metadata["Date"] = ", ".join(
                    extract_metadata_values(book.get_metadata("DC", "date"))
                )
            if book.get_metadata("DC", "language"):
                metadata["Language"] = ", ".join(
                    extract_metadata_values(book.get_metadata("DC", "language"))
                )

            # Add metadata section if available
            if metadata:
                meta_str = "\n".join(f"{k}: {v}" for k, v in metadata.items())
                text_parts.append(f"[Metadata]\n{meta_str}\n")

            # Process all items in the book
            chapter_num = 0
            has_content = False  # Track if we found actual content (not just metadata)
            for item in book.get_items():
                # Only process document items (XHTML content)
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    chapter_num += 1

                    # Parse HTML content
                    html_content = item.get_body_content()
                    if html_content:
                        soup = BeautifulSoup(html_content, "html.parser")

                        # Remove script and style elements
                        for script in soup(["script", "style"]):
                            script.decompose()

                        # Extract text with some structure preservation
                        chapter_texts = []

                        # Process headings to preserve structure
                        for heading in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
                            heading_text = heading.get_text(strip=True)
                            if heading_text:
                                level = heading.name[1]  # Extract number from h1, h2, etc.
                                chapter_texts.append(f"{'#' * int(level)} {heading_text}")

                        # Extract all text content
                        text_content = soup.get_text(separator="\n")

                        # Clean up excessive whitespace while preserving paragraph breaks
                        lines = [line.strip() for line in text_content.splitlines()]
                        # Remove empty lines and duplicates from heading extraction
                        cleaned_lines = []
                        prev_line = ""
                        for line in lines:
                            if line and line != prev_line:
                                cleaned_lines.append(line)
                                prev_line = line

                        chapter_text = "\n\n".join(cleaned_lines)

                        # Only count as content if it has substantial text (>50 chars)
                        # This filters out navigation-only items that just contain titles
                        if chapter_text.strip() and len(chapter_text.strip()) > 50:
                            has_content = True  # Found actual content
                            # Add chapter marker for multi-chapter books
                            if (
                                chapter_num > 1
                                or len(list(book.get_items_of_type(ebooklib.ITEM_DOCUMENT))) > 1
                            ):
                                text_parts.append(f"[Chapter {chapter_num}]\n{chapter_text}")
                            else:
                                text_parts.append(chapter_text)

            # Validate that we have actual content, not just metadata
            if not has_content:
                raise HTTPException(
                    status_code=400,
                    detail="ePUB extraction produced no text. The file may be empty, image-only, or malformed.",
                )

            # Join all parts with double newlines
            text = "\n\n".join(text_parts) if text_parts else ""

        except ImportError:
            raise HTTPException(
                status_code=400,
                detail="ePUB support not available. Install ebooklib: pip install ebooklib",
            ) from None
        except HTTPException:
            # Re-raise HTTP exceptions without wrapping
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"ePUB parsing failed: {e}") from None

    elif file_ext in {".html", ".htm"}:
        # Handle HTML documents (#2666)
        try:
            from bs4 import BeautifulSoup

            # Decode HTML content
            try:
                html_content = content.decode("utf-8")
            except UnicodeDecodeError:
                # Try common fallback encodings
                for encoding in ["iso-8859-1", "windows-1252", "cp1252"]:
                    try:
                        html_content = content.decode(encoding)
                        break
                    except UnicodeDecodeError:
                        continue
                else:
                    raise HTTPException(
                        status_code=400, detail="Unable to decode HTML file. Unsupported encoding."
                    )

            # Parse HTML
            soup = BeautifulSoup(html_content, "html.parser")

            # Remove boilerplate and non-content elements
            for tag in soup(
                ["script", "style", "nav", "header", "footer", "aside", "iframe", "noscript"]
            ):
                tag.decompose()

            # Remove common ad and tracking elements
            for class_name in [
                "advertisement",
                "ad-container",
                "social-share",
                "comments",
                "sidebar",
            ]:
                for elem in soup.find_all(class_=class_name):
                    elem.decompose()

            text_parts = []

            # Extract metadata from meta tags
            metadata = {}

            # Page title
            title_tag = soup.find("title")
            if title_tag and title_tag.string:
                metadata["Title"] = title_tag.string.strip()

            # Meta tags
            meta_mappings = {
                "description": "Description",
                "author": "Author",
                "keywords": "Keywords",
                "og:title": "OG_Title",
                "og:description": "OG_Description",
            }

            for meta_name, metadata_key in meta_mappings.items():
                # Try name attribute
                meta_tag = soup.find("meta", attrs={"name": meta_name})
                if not meta_tag:
                    # Try property attribute (for Open Graph tags)
                    meta_tag = soup.find("meta", attrs={"property": meta_name})

                if meta_tag:
                    meta_content = meta_tag.get("content")
                    if meta_content and isinstance(meta_content, str):
                        metadata[metadata_key] = meta_content.strip()

            # Add metadata section if available
            if metadata:
                meta_str = "\n".join(f"{k}: {v}" for k, v in metadata.items())
                text_parts.append(f"[Metadata]\n{meta_str}\n")

            # Extract main content with semantic structure preservation
            # Try to find main content area (common patterns)
            main_content = (
                soup.find("main")
                or soup.find("article")
                or soup.find("div", class_="content")
                or soup.find("div", id="content")
                or soup.find("div", class_="main")
                or soup.find("div", id="main")
                or soup.body
                or soup
            )

            # Process content preserving structure
            content_parts = []

            # Extract headings with hierarchy
            for heading in main_content.find_all(["h1", "h2", "h3", "h4", "h5", "h6"]):
                heading_text = heading.get_text(strip=True)
                if heading_text:
                    level = heading.name[1]  # Extract number from h1, h2, etc.
                    content_parts.append(f"{'#' * int(level)} {heading_text}")

            # Extract paragraphs and preserve block structure
            for elem in main_content.find_all(["p", "div", "section", "blockquote", "pre", "code"]):
                elem_text = elem.get_text(strip=True)

                # Skip empty elements
                if not elem_text:
                    continue

                # Skip if this is just a container with nested elements we'll process separately
                if elem.find(["p", "div", "section"]) and elem.name == "div":
                    continue

                # Format blockquotes
                if elem.name == "blockquote":
                    elem_text = "> " + elem_text

                # Format code blocks
                if elem.name in ["pre", "code"]:
                    elem_text = f"```\n{elem_text}\n```"

                content_parts.append(elem_text)

            # Extract list items with structure
            for ul in main_content.find_all("ul"):
                list_items = []
                for li in ul.find_all("li", recursive=False):
                    li_text = li.get_text(strip=True)
                    if li_text:
                        list_items.append(f"• {li_text}")
                if list_items:
                    content_parts.append("\n".join(list_items))

            for ol in main_content.find_all("ol"):
                list_items = []
                for idx, li in enumerate(ol.find_all("li", recursive=False), start=1):
                    li_text = li.get_text(strip=True)
                    if li_text:
                        list_items.append(f"{idx}. {li_text}")
                if list_items:
                    content_parts.append("\n".join(list_items))

            # Extract tables
            for html_table in main_content.find_all("table"):
                html_table_rows: list[str] = []
                for row in html_table.find_all("tr"):
                    cells = row.find_all(["th", "td"])
                    if cells:
                        row_text = " | ".join(cell.get_text(strip=True) for cell in cells)
                        if row_text:
                            html_table_rows.append(row_text)
                if html_table_rows:
                    content_parts.append("[Table]\n" + "\n".join(html_table_rows))

            # If we extracted structured content, use it
            if content_parts:
                # Remove duplicates (headings might be extracted twice)
                seen = set()
                unique_parts = []
                for part in content_parts:
                    if part not in seen:
                        unique_parts.append(part)
                        seen.add(part)
                text_parts.extend(unique_parts)
            else:
                # Fallback: extract all text from main content
                text_content = main_content.get_text(separator="\n")
                # Clean up excessive whitespace
                lines = [line.strip() for line in text_content.splitlines()]
                cleaned_lines = [line for line in lines if line]
                if cleaned_lines:
                    text_parts.append("\n\n".join(cleaned_lines))

            # Combine all parts
            text = "\n\n".join(text_parts)

            # Validate extracted content
            if not text or not text.strip():
                raise HTTPException(
                    status_code=400,
                    detail="HTML extraction produced no text. The file may be empty or contain only boilerplate.",
                )

        except ImportError:
            raise HTTPException(
                status_code=400,
                detail="HTML support not available. Install beautifulsoup4: pip install beautifulsoup4",
            ) from None
        except HTTPException:
            # Re-raise HTTP exceptions
            raise
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"HTML parsing failed: {e}") from None

    else:
        # Plain text
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError as e:
            raise HTTPException(status_code=400, detail=f"File encoding error: {e}") from None

    # Use filename as doc_id if not provided (sanitize to prevent path traversal)
    safe_filename = os.path.basename(file.filename or "").strip() if file.filename else ""
    final_doc_id = doc_id or safe_filename or f"upload_{uuid.uuid4().hex[:8]}"

    # Ingest
    try:
        result = ingest_document(doc_id=final_doc_id, text=text, collection=resolved_collection)
        prom_metrics.RAG_INGESTS_TOTAL.labels(source="upload", result="success").inc()
        return RagIngestResponse(
            doc_id=final_doc_id,
            chunks_ingested=result["chunks_ingested"],
            status=result["status"],
            doc_hash=result["doc_hash"],
            previous_hash=result["previous_hash"],
            collection=resolved_collection,
        )
    except Exception as e:
        prom_metrics.RAG_INGESTS_TOTAL.labels(source="upload", result="failure").inc()
        raise HTTPException(status_code=500, detail=f"Ingestion failed: {e}") from None


@api.post("/rag/index-repo")
def rag_index_repo(_request: Request, req: RagIndexRepoRequest) -> RagIndexRepoResponse:
    """Index a code repository for RAG."""
    from nyxgpt.rag.rag import ingest_repository

    try:
        extensions_set = set(req.extensions) if req.extensions else None

        result = ingest_repository(
            repo_path=req.repo_path,
            doc_id_prefix=req.doc_id_prefix,
            extensions=extensions_set,
            extract_docs_only=req.extract_docs_only,
            ensure_schema=req.ensure_schema,
            collection=req.collection,
        )
        prom_metrics.RAG_INGESTS_TOTAL.labels(source="repo", result="success").inc()

        return RagIndexRepoResponse(
            total_files=result["total_files"],
            total_chunks=result["total_chunks"],
            files=result["files"],
            doc_ids=result["doc_ids"],
        )
    except ValueError as e:
        prom_metrics.RAG_INGESTS_TOTAL.labels(source="repo", result="failure").inc()
        raise HTTPException(status_code=400, detail=str(e)) from e
    except Exception as e:
        prom_metrics.RAG_INGESTS_TOTAL.labels(source="repo", result="failure").inc()
        raise HTTPException(status_code=500, detail=f"Repository indexing failed: {e}") from e


app.include_router(api)
